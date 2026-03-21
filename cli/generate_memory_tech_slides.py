#!/usr/bin/env python3
"""
Generate PowerPoint slides comparing memory technology energy characteristics.

This deck provides the apples-to-apples memory energy comparison needed to
normalize the architecture energy models. It covers the three memory families
used in AI compute systems:

Slide 1: Title -- Why Memory Technology Matters for Energy Comparisons
Slide 2: DRAM Energy Anatomy -- The Four Components Every Access Pays
Slide 3: LPDDR5/5X -- Package-on-Package, Low-Voltage Signaling
Slide 4: GDDR6/6X/7 -- High-Speed PCB Signaling, PAM4/PAM3
Slide 5: HBM2/3 -- TSV Stacking, Silicon Interposer, Wide Bus

References:
  - Horowitz, ISSCC 2014 ("Computing's Energy Problem")
  - Vogelsang, "Understanding the Energy Consumption of DRAM" (MICRO 2010)
  - Ha, "Understanding and Improving the Energy Efficiency of DRAM" (Stanford)
  - JEDEC LPDDR5, GDDR6/7, HBM3 specifications
"""

import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'src'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# -- Color palette (matches all prior decks) ---------------------------------
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xD8)
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
MED_GRAY    = RGBColor(0x99, 0x99, 0x99)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)
RED_ACCENT  = RGBColor(0xE0, 0x40, 0x40)
GREEN       = RGBColor(0x2E, 0xCC, 0x71)
YELLOW      = RGBColor(0xF3, 0xC6, 0x23)

# Memory family colors
LPDDR_COLOR  = RGBColor(0x00, 0x96, 0xD6)  # blue
GDDR_COLOR   = RGBColor(0xDD, 0x66, 0x33)  # orange
HBM_COLOR    = RGBColor(0x8B, 0x5C, 0xF6)  # purple
CORE_COLOR   = RGBColor(0xCC, 0x44, 0x44)  # red (row activation / core)
IO_COLOR     = RGBColor(0x44, 0xAA, 0x44)  # green (I/O driver)
SENSE_COLOR  = RGBColor(0xDD, 0xAA, 0x33)  # gold (sense amplifier)
PHY_COLOR    = RGBColor(0x55, 0x88, 0xBB)  # steel blue (PHY/controller)

ROW_EVEN = RGBColor(0x2A, 0x2A, 0x4E)
ROW_ODD  = RGBColor(0x22, 0x22, 0x3E)
ROW_HDR  = RGBColor(0x33, 0x33, 0x55)


# ============================================================================
# Helpers (same API as all prior decks)
# ============================================================================

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=16, color=WHITE, alignment=PP_ALIGN.LEFT,
                          font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld, clr = item, False, color
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            clr = item[2] if len(item) > 2 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.bold = bld
        p.font.color.rgb = clr
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1.0) * 2)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     text="", font_size=14, font_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_rect(slide, left, top, width, height, fill_color,
             text="", font_size=14, font_color=WHITE, bold=False,
             line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow_right(slide, left, top, width=0.3, height=0.2, color=MED_GRAY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_arrow_down(slide, left, top, width=0.2, height=0.3, color=MED_GRAY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=MED_GRAY, width=1.5):
    connector = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def make_table(slide, left, top, width, height, headers, data,
               header_color=ACCENT_BLUE, font_size=9):
    rows = len(data) + 1
    cols = len(headers)
    tbl = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top),
        Inches(width), Inches(height)).table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_EVEN if i % 2 == 0 else ROW_ODD
    return tbl


# ============================================================================
# SLIDE 1: Title
# ============================================================================

def build_slide_1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.5, 0.6, 12.0, 0.8,
                "Memory Technology Energy Comparison",
                font_size=34, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 1.3, 12.0, 0.5,
                "LPDDR5  vs  GDDR6  vs  HBM3: Energy per Byte at the System Level",
                font_size=22, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    add_line(slide, 2.5, 2.0, 10.5, 2.0, color=ACCENT_BLUE, width=2)

    # Why this matters
    add_rounded_rect(slide, 1.5, 2.3, 10.0, 1.3, RGBColor(0x2A, 0x2A, 0x4E))
    add_multiline_textbox(slide, 1.8, 2.4, 9.4, 1.1, [
        ("Why Memory Technology Dominates the Energy Comparison:", True, ORANGE),
        ("In the validation deck, KPU memory was 57% of total energy (LPDDR5 at 8 pJ/byte)",
         False, LIGHT_GRAY),
        ("while GPU memory was 26% (HBM3 at 5.5 pJ/byte). This 45% difference in off-chip",
         False, LIGHT_GRAY),
        ("energy per byte obscures the architectural comparison. To make a fair argument,",
         False, LIGHT_GRAY),
        ("we must understand exactly what drives these numbers.", True, WHITE),
    ], font_size=14, alignment=PP_ALIGN.CENTER)

    # Three-column memory family overview
    col_y = 4.0
    families = [
        ("LPDDR5 / 5X", LPDDR_COLOR, [
            "Package-on-Package (PoP)",
            "Low-voltage I/O (0.5V)",
            "4 channels x 16-bit",
            "6.4-8.5 GT/s per pin",
            "8.0-7.0 pJ/byte",
            "Target: Edge / Automotive",
        ]),
        ("GDDR6 / 6X / 7", GDDR_COLOR, [
            "Discrete chips on PCB",
            "High-speed PCB traces (1.35V)",
            "8 chips x 32-bit bus",
            "16-32 Gbps per pin",
            "14.0-18.0 pJ/byte",
            "Target: Gaming / Consumer GPU",
        ]),
        ("HBM2 / 3 / 3E", HBM_COLOR, [
            "3D TSV stacking on interposer",
            "Ultra-short wires (~100um)",
            "1024-bit bus per stack",
            "3.2-9.6 Gbps per pin",
            "5.0-7.0 pJ/byte",
            "Target: AI / HPC Datacenter",
        ]),
    ]

    for i, (title, color, items) in enumerate(families):
        x = 0.5 + i * 4.3
        add_rounded_rect(slide, x, col_y, 3.8, 0.45, color,
                         title, font_size=15, bold=True)
        for j, item in enumerate(items):
            is_energy = "pJ/byte" in item
            add_textbox(slide, x + 0.2, col_y + 0.55 + j * 0.28, 3.4, 0.25,
                        item, font_size=11,
                        bold=is_energy,
                        color=YELLOW if is_energy else LIGHT_GRAY)

    # Bottom
    add_textbox(slide, 0.5, 6.8, 12.0, 0.5,
                "The energy per byte varies 3.6x across these technologies (5.0 to 18.0 pJ/byte).\n"
                "This is comparable to the entire architectural energy advantage we are trying to measure.",
                font_size=13, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

    return slide


# ============================================================================
# SLIDE 2: DRAM Energy Anatomy
# ============================================================================

def build_slide_2_dram_anatomy(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.1, 12.5, 0.45,
                "DRAM Energy Anatomy: The Four Components Every Access Pays",
                font_size=26, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.3, 0.55, 12.5, 0.3,
                "Based on Vogelsang (MICRO 2010) and Horowitz (ISSCC 2014) -- all DRAM technologies share this structure",
                font_size=12, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    # -- DRAM Access Sequence Diagram --
    seq_y = 1.0
    step_h = 0.65
    step_w = 2.5
    desc_w = 3.8
    gap = 0.1

    steps = [
        ("1. Row Activation", CORE_COLOR,
         "Open a row in the DRAM array:\n"
         "- Precharge bitlines to VDD/2\n"
         "- Assert wordline (select row)\n"
         "- Sense amplifiers detect charge\n"
         "- Latch entire row (~8-16 KB)",
         "~3-5 pJ/byte\n(amortized\nover burst)"),
        ("2. Column Access", SENSE_COLOR,
         "Select bytes within the open row:\n"
         "- Column decoder selects columns\n"
         "- Mux routes data to I/O gating\n"
         "- Burst length: 16 (LPDDR5) or\n"
         "  32 (HBM3) beats",
         "~1-2 pJ/byte"),
        ("3. I/O & Signaling", IO_COLOR,
         "Drive data off the DRAM die:\n"
         "- I/O driver charges bus capacitance\n"
         "- Signal travels to SoC receiver\n"
         "- Energy = C_bus * V_swing^2 per bit\n"
         "- Dominant factor: wire length + voltage",
         "~1-10 pJ/byte\n(varies with\npackaging)"),
        ("4. PHY & Controller", PHY_COLOR,
         "SoC-side memory interface:\n"
         "- Clock/data recovery (CDR)\n"
         "- Serializer/deserializer (SerDes)\n"
         "- Training, calibration circuits\n"
         "- Command/address encoding",
         "~1-3 pJ/byte"),
    ]

    for i, (name, color, desc, energy) in enumerate(steps):
        sy = seq_y + i * (step_h + gap)

        # Step box
        add_rounded_rect(slide, 0.3, sy, step_w, step_h, color,
                         name, font_size=11, bold=True)

        # Description
        add_textbox(slide, 0.3 + step_w + 0.15, sy, desc_w, step_h,
                    desc, font_size=9, color=LIGHT_GRAY)

        # Energy
        add_textbox(slide, 0.3 + step_w + desc_w + 0.2, sy + 0.05,
                    1.2, step_h - 0.1,
                    energy, font_size=9, bold=True, color=YELLOW,
                    alignment=PP_ALIGN.CENTER)

        # Arrow to next
        if i < len(steps) - 1:
            add_arrow_down(slide, 1.5, sy + step_h, 0.2, gap, color=MED_GRAY)

    # -- Energy Equation --
    eq_y = seq_y + 4 * (step_h + gap) + 0.1
    add_textbox(slide, 0.3, eq_y, 12.5, 0.3,
                "Total Energy per Byte:",
                font_size=14, bold=True, color=ACCENT_TEAL)
    add_textbox(slide, 0.3, eq_y + 0.3, 12.5, 0.3,
                "E_total = E_row_activation/burst_length  +  E_column_access  +  E_io_signaling  +  E_phy_controller",
                font_size=16, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER, font_name="Consolas")

    # -- RIGHT SIDE: What varies across technologies --
    var_x = 8.0
    var_y = 1.0

    add_textbox(slide, var_x, var_y, 5.0, 0.3,
                "What Varies Across Technologies:",
                font_size=14, bold=True, color=ORANGE)

    tech_diffs = [
        ("Row Activation (~same)", "All use same DRAM cell\n"
         "physics (capacitor + transistor).\n"
         "~3-5 pJ/byte regardless of tech.", CORE_COLOR),
        ("Column Access (~same)", "Same array organization.\n"
         "Burst length affects amortization\n"
         "(BL16 for LPDDR5, BL32 for HBM3).", SENSE_COLOR),
        ("I/O Signaling (VARIES 10x!)", "THIS is where the difference lives:\n"
         "LPDDR5: PoP, short wire, 0.5V swing\n"
         "GDDR6: PCB trace, long wire, 1.35V\n"
         "HBM3: TSV, ~100um, 1.1V but wide bus", IO_COLOR),
        ("PHY/Controller (varies 2-3x)", "Complexity tracks signaling:\n"
         "LPDDR5: Simpler PHY, lower rate\n"
         "GDDR7: PAM3 SerDes, complex\n"
         "HBM3: Many channels, simpler/pin", PHY_COLOR),
    ]

    ty = var_y + 0.35
    for name, desc, color in tech_diffs:
        add_rounded_rect(slide, var_x, ty, 4.8, 0.3, color,
                         name, font_size=9, bold=True)
        add_textbox(slide, var_x + 0.1, ty + 0.32, 4.6, 0.55,
                    desc, font_size=8, color=LIGHT_GRAY)
        ty += 1.0

    # Bottom insight
    add_rounded_rect(slide, 1.0, eq_y + 0.7, 11.0, 0.5, RGBColor(0x2A, 0x2A, 0x4E))
    add_textbox(slide, 1.0, eq_y + 0.75, 11.0, 0.4,
                "Key Insight: The DRAM core energy (row + column) is ~4-7 pJ/byte for ALL technologies. "
                "The 3.6x range (5 to 18 pJ/byte) comes almost entirely from I/O signaling and packaging.",
                font_size=13, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)

    return slide


# ============================================================================
# SLIDE 3: LPDDR5 / 5X
# ============================================================================

def build_slide_3_lpddr(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.1, 12.5, 0.45,
                "LPDDR5/5X: Low-Power Package-on-Package Memory",
                font_size=26, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # -- System Diagram --
    diag_y = 0.65

    # SoC die
    add_rect(slide, 1.0, diag_y, 4.5, 2.5, RGBColor(0x2A, 0x3A, 0x5E),
             "", line_color=LPDDR_COLOR)
    add_textbox(slide, 1.0, diag_y + 0.05, 4.5, 0.3,
                "SoC / Processor Die", font_size=12, bold=True,
                color=LPDDR_COLOR, alignment=PP_ALIGN.CENTER)

    # PHY blocks inside SoC
    for ch in range(4):
        px = 1.3 + ch * 1.05
        add_rounded_rect(slide, px, diag_y + 0.4, 0.9, 0.5,
                         PHY_COLOR,
                         f"LPDDR5\nPHY Ch{ch}", font_size=7, bold=True)

    # Memory controller
    add_rounded_rect(slide, 1.5, diag_y + 1.1, 3.5, 0.4, RGBColor(0x44, 0x66, 0x88),
                     "Memory Controller (4-channel, 64-bit total)",
                     font_size=9, bold=True)

    # Compute
    add_rounded_rect(slide, 1.5, diag_y + 1.65, 3.5, 0.6, RGBColor(0x33, 0x55, 0x77),
                     "CPU / GPU / NPU Cores\n(compute logic)",
                     font_size=10, bold=True)

    # PoP DRAM on top
    add_rect(slide, 1.0, diag_y + 2.7, 4.5, 1.2, RGBColor(0x1A, 0x3A, 0x5A),
             "", line_color=LPDDR_COLOR)
    add_textbox(slide, 1.0, diag_y + 2.75, 4.5, 0.3,
                "LPDDR5 DRAM (Package-on-Package)", font_size=12,
                bold=True, color=LPDDR_COLOR, alignment=PP_ALIGN.CENTER)

    # DRAM dies
    for d in range(4):
        dx = 1.3 + d * 1.05
        add_rounded_rect(slide, dx, diag_y + 3.1, 0.9, 0.55,
                         LPDDR_COLOR,
                         f"DRAM\nDie {d}\n(16-bit)", font_size=7, bold=True)

    # Wiring annotation
    add_textbox(slide, 1.0, diag_y + 4.0, 4.5, 0.3,
                "Wire: ~1-3 mm (PoP solder balls)  |  V_IO: 0.5V (LPDDR5) / 0.35V (LPDDR5X)",
                font_size=9, bold=True, color=YELLOW,
                alignment=PP_ALIGN.CENTER)

    # -- RIGHT SIDE: Energy Breakdown --
    eng_x = 6.2
    eng_y = 0.65

    add_textbox(slide, eng_x, eng_y, 6.5, 0.3,
                "Energy Breakdown per Byte (LPDDR5 @ 6400 MT/s)",
                font_size=14, bold=True, color=LPDDR_COLOR)

    # Breakdown table
    breakdown_headers = ["Component", "Energy (pJ/B)", "% of Total", "Key Driver"]
    breakdown_data = [
        ["Row Activation", "3.2", "40%", "DRAM cell physics (same as all)"],
        ["Column Access", "1.2", "15%", "Burst length 16, column mux"],
        ["I/O Signaling", "1.8", "22%", "0.5V swing, short PoP wire (~2mm)"],
        ["PHY / Controller", "1.0", "13%", "Simple PHY, 4 channels"],
        ["Clock / Refresh", "0.8", "10%", "CK differential, tREFI refresh"],
        ["TOTAL", "8.0", "100%", ""],
    ]

    tbl = make_table(slide, eng_x, eng_y + 0.35, 6.5, 2.3,
                     breakdown_headers, breakdown_data,
                     header_color=LPDDR_COLOR, font_size=9)
    # Highlight total row
    for j in range(4):
        cell = tbl.cell(6, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ROW_HDR
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = YELLOW

    # Stacked bar
    bar_y = eng_y + 2.8
    bar_h = 0.45
    bar_w = 6.0
    components = [
        ("Row: 3.2", 3.2, CORE_COLOR),
        ("Col: 1.2", 1.2, SENSE_COLOR),
        ("I/O: 1.8", 1.8, IO_COLOR),
        ("PHY: 1.0", 1.0, PHY_COLOR),
        ("Clk: 0.8", 0.8, MED_GRAY),
    ]
    total = 8.0
    cx = eng_x
    for label, energy, color in components:
        w = bar_w * (energy / total)
        fsz = 7 if w < 0.7 else 8
        add_rect(slide, cx, bar_y, w, bar_h, color,
                 label, font_size=fsz, font_color=WHITE, bold=True)
        cx += w

    # Key characteristics
    char_y = bar_y + 0.65
    add_textbox(slide, eng_x, char_y, 6.5, 0.3,
                "Key LPDDR5 Characteristics:", font_size=13, bold=True, color=WHITE)
    add_multiline_textbox(slide, eng_x, char_y + 0.3, 6.5, 2.5, [
        ("Low I/O voltage: 0.5V (LPDDR5) / 0.35V (LPDDR5X)", False, LIGHT_GRAY),
        ("  E_io = C_wire * V^2 -> 0.5V is 7.3x less than GDDR6's 1.35V", True, GREEN),
        ("", False, WHITE),
        ("Short wires: Package-on-Package stacking = ~1-3 mm trace", False, LIGHT_GRAY),
        ("  vs GDDR6 PCB: ~20-50 mm  |  vs HBM TSV: ~0.05-0.1 mm", False, MED_GRAY),
        ("", False, WHITE),
        ("Narrow bus: 4 channels x 16-bit = 64-bit total bus width", False, LIGHT_GRAY),
        ("  Lower bandwidth (25-34 GB/s) but very power-efficient per byte", False, MED_GRAY),
        ("", False, WHITE),
        ("LPDDR5X improvement: Lower V_IO (0.35V) -> ~7.0 pJ/byte (12% better)", False, YELLOW),
    ], font_size=10)

    return slide


# ============================================================================
# SLIDE 4: GDDR6 / 6X / 7
# ============================================================================

def build_slide_4_gddr(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.1, 12.5, 0.45,
                "GDDR6/6X/7: High-Speed Discrete Graphics Memory",
                font_size=26, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # -- System Diagram --
    diag_y = 0.65

    # GPU die
    add_rect(slide, 0.5, diag_y, 2.8, 2.5, RGBColor(0x2A, 0x3A, 0x5E),
             "", line_color=GDDR_COLOR)
    add_textbox(slide, 0.5, diag_y + 0.05, 2.8, 0.3,
                "GPU Die", font_size=12, bold=True,
                color=GDDR_COLOR, alignment=PP_ALIGN.CENTER)

    # Memory controller inside GPU
    add_rounded_rect(slide, 0.7, diag_y + 0.4, 2.4, 0.4, RGBColor(0x44, 0x66, 0x88),
                     "Memory Controller\n(256-384 bit bus)", font_size=8, bold=True)

    # PHY
    add_rounded_rect(slide, 0.7, diag_y + 0.9, 2.4, 0.35, PHY_COLOR,
                     "GDDR PHY (SerDes, PAM4/PAM3)", font_size=8, bold=True)

    # Compute
    add_rounded_rect(slide, 0.7, diag_y + 1.4, 2.4, 0.9, RGBColor(0x33, 0x55, 0x77),
                     "SMs / Compute Units\n(CUDA cores, TensorCores)",
                     font_size=9, bold=True)

    # PCB traces
    add_rect(slide, 3.5, diag_y + 0.8, 1.5, 0.3, RGBColor(0x55, 0x44, 0x33),
             "PCB Traces\n20-50 mm", font_size=7, font_color=ORANGE, bold=True,
             line_color=ORANGE)

    # GDDR chips (discrete, on PCB)
    for d in range(4):
        dx = 5.3
        dy = diag_y + d * 0.62
        add_rounded_rect(slide, dx, dy, 1.3, 0.5, GDDR_COLOR,
                         f"GDDR6\nChip {d} (32b)", font_size=7, bold=True)
    for d in range(4):
        dx = 6.8
        dy = diag_y + d * 0.62
        add_rounded_rect(slide, dx, dy, 1.3, 0.5, GDDR_COLOR,
                         f"GDDR6\nChip {d+4} (32b)", font_size=7, bold=True)

    add_textbox(slide, 5.3, diag_y + 2.55, 2.8, 0.3,
                "8 chips x 32-bit = 256-bit bus", font_size=9,
                bold=True, color=GDDR_COLOR, alignment=PP_ALIGN.CENTER)

    # Wire annotation
    add_textbox(slide, 0.5, diag_y + 2.6, 7.8, 0.3,
                "Wire: 20-50 mm PCB  |  V_IO: 1.35V (GDDR6), PAM4 (GDDR6X), PAM3 (GDDR7)",
                font_size=9, bold=True, color=YELLOW)

    # -- RIGHT SIDE: Energy Breakdown --
    eng_x = 8.5
    eng_y = 0.65

    add_textbox(slide, eng_x, eng_y, 4.5, 0.3,
                "Energy per Byte (pJ/B)", font_size=14, bold=True, color=GDDR_COLOR)

    # Comparison table across GDDR variants
    gddr_headers = ["Component", "GDDR6", "GDDR6X", "GDDR7"]
    gddr_data = [
        ["Row Activation", "3.5",  "3.5",  "3.2"],
        ["Column Access",  "1.5",  "1.5",  "1.3"],
        ["I/O Signaling",  "6.5",  "9.0",  "5.8"],
        ["PHY / Controller","2.0", "2.5",  "2.2"],
        ["Clock / Refresh", "1.5", "1.5",  "1.5"],
        ["TOTAL",          "15.0", "18.0", "14.0"],
    ]

    tbl = make_table(slide, eng_x, eng_y + 0.35, 4.5, 2.3,
                     gddr_headers, gddr_data,
                     header_color=GDDR_COLOR, font_size=9)
    for j in range(4):
        cell = tbl.cell(6, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ROW_HDR
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = YELLOW

    # GDDR6X stacked bar (worst case)
    bar_y = eng_y + 2.8
    bar_h = 0.4
    bar_w = 4.0
    gddr6x_comps = [
        ("Row: 3.5", 3.5, CORE_COLOR),
        ("Col: 1.5", 1.5, SENSE_COLOR),
        ("I/O: 9.0", 9.0, IO_COLOR),
        ("PHY: 2.5", 2.5, PHY_COLOR),
        ("Clk: 1.5", 1.5, MED_GRAY),
    ]
    cx = eng_x
    for label, energy, color in gddr6x_comps:
        w = bar_w * (energy / 18.0)
        fsz = 7 if w < 0.5 else 8
        add_rect(slide, cx, bar_y, w, bar_h, color,
                 label, font_size=fsz, font_color=WHITE, bold=True)
        cx += w

    add_textbox(slide, eng_x, bar_y + 0.45, 4.5, 0.2,
                "GDDR6X: I/O is 50% of total energy", font_size=9,
                bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

    # Why GDDR is expensive
    why_y = 3.5
    add_textbox(slide, 0.5, why_y, 12.5, 0.3,
                "Why GDDR I/O Energy Is So High:", font_size=16, bold=True, color=ORANGE)

    add_multiline_textbox(slide, 0.5, why_y + 0.35, 7.5, 2.8, [
        ("1. Long PCB traces (20-50 mm):", True, WHITE),
        ("   Capacitance ~ 1-2 pF per trace (vs ~0.05 pF for HBM TSV)", False, LIGHT_GRAY),
        ("   E = C * V^2 per transition = 2 pF * 1.35^2 = 3.6 pJ/bit", False, MED_GRAY),
        ("", False, WHITE),
        ("2. High I/O voltage (1.35V for GDDR6):", True, WHITE),
        ("   Compared to 0.5V for LPDDR5: (1.35/0.5)^2 = 7.3x more energy per bit", False, LIGHT_GRAY),
        ("   This is pure physics: E = C*V^2, voltage dominates", False, MED_GRAY),
        ("", False, WHITE),
        ("3. PAM4 (GDDR6X) adds complexity:", True, WHITE),
        ("   4-level signaling requires tighter voltage margins + more complex receivers", False, LIGHT_GRAY),
        ("   Higher BW per pin, but MORE energy per bit (18 vs 15 pJ/byte)", False, MED_GRAY),
        ("", False, WHITE),
        ("4. PAM3 (GDDR7) partially recovers:", True, WHITE),
        ("   3-level signaling: 50% better noise margin than PAM4", False, LIGHT_GRAY),
        ("   Combined with lower voltage -> 14 pJ/byte (best GDDR)", False, MED_GRAY),
    ], font_size=10)

    return slide


# ============================================================================
# SLIDE 5: HBM2 / 3 / 3E
# ============================================================================

def build_slide_5_hbm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.1, 12.5, 0.45,
                "HBM2/3/3E: 3D-Stacked Memory on Silicon Interposer",
                font_size=26, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # -- System Diagram: Cross-section view --
    diag_y = 0.65

    # Substrate / PCB
    add_rect(slide, 0.5, diag_y + 3.2, 7.5, 0.4, RGBColor(0x55, 0x44, 0x33),
             "Package Substrate / BGA", font_size=10, font_color=LIGHT_GRAY, bold=True)

    # Silicon Interposer
    add_rect(slide, 0.7, diag_y + 2.6, 7.1, 0.5, RGBColor(0x44, 0x33, 0x66),
             "Silicon Interposer (65nm passive Si, ~30mm x 30mm)", font_size=10,
             font_color=HBM_COLOR, bold=True, line_color=HBM_COLOR)

    # GPU/SoC die on interposer
    add_rect(slide, 1.0, diag_y + 0.5, 3.0, 2.0, RGBColor(0x2A, 0x3A, 0x5E),
             "", line_color=ACCENT_BLUE)
    add_textbox(slide, 1.0, diag_y + 0.55, 3.0, 0.3,
                "GPU / Accelerator Die", font_size=11, bold=True,
                color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_rounded_rect(slide, 1.2, diag_y + 0.9, 2.6, 0.35, PHY_COLOR,
                     "HBM PHY (per-stack controller)", font_size=8, bold=True)
    add_rounded_rect(slide, 1.2, diag_y + 1.4, 2.6, 0.9, RGBColor(0x33, 0x55, 0x77),
                     "Compute Logic\n(SMs, TensorCores)", font_size=9, bold=True)

    # HBM stacks (cross-section showing dies)
    for stack in range(3):
        sx = 4.5 + stack * 1.2
        # Stack of dies
        for die in range(4):
            dy = diag_y + 1.8 - die * 0.35
            shade_val = 0x70 + die * 0x15
            die_color = RGBColor(shade_val, 0x50, 0xF0)
            add_rect(slide, sx, dy, 0.9, 0.3, die_color,
                     f"Die {die}", font_size=6, font_color=WHITE)

        # TSV label
        add_textbox(slide, sx, diag_y + 2.1, 0.9, 0.3,
                    "TSV\n1024-bit", font_size=6, bold=True,
                    color=HBM_COLOR, alignment=PP_ALIGN.CENTER)

        # Base die
        add_rect(slide, sx, diag_y + 2.3, 0.9, 0.25,
                 RGBColor(0x44, 0x33, 0x66),
                 "Base", font_size=6, font_color=HBM_COLOR, bold=True)

    add_textbox(slide, 4.5, diag_y + 0.15, 3.5, 0.3,
                "HBM3 Stacks (3-6 stacks, 8-12 dies each)",
                font_size=10, bold=True, color=HBM_COLOR,
                alignment=PP_ALIGN.CENTER)

    # Wire annotations
    add_textbox(slide, 0.5, diag_y + 3.7, 7.5, 0.35,
                "TSV pitch: ~5-10 um  |  Wire length: ~50-100 um (vertical) + ~5-15 mm (interposer)  |  V_IO: 1.1V",
                font_size=9, bold=True, color=YELLOW)

    # -- RIGHT SIDE: Energy Breakdown --
    eng_x = 8.5
    eng_y = 0.65

    add_textbox(slide, eng_x, eng_y, 4.5, 0.3,
                "Energy per Byte (pJ/B)", font_size=14, bold=True, color=HBM_COLOR)

    hbm_headers = ["Component", "HBM2", "HBM2E", "HBM3", "HBM3E"]
    hbm_data = [
        ["Row Activation", "2.8",  "2.6",  "2.2", "2.0"],
        ["Column Access",  "1.0",  "1.0",  "0.8", "0.8"],
        ["I/O (TSV+Interp)","1.5", "1.3",  "1.0", "0.8"],
        ["PHY / Controller","1.0", "1.0",  "0.9", "0.8"],
        ["Clock / Refresh", "0.7",  "0.6",  "0.6", "0.6"],
        ["TOTAL",           "7.0",  "6.5",  "5.5", "5.0"],
    ]

    tbl = make_table(slide, eng_x, eng_y + 0.35, 4.5, 2.3,
                     hbm_headers, hbm_data,
                     header_color=HBM_COLOR, font_size=8)
    for j in range(5):
        cell = tbl.cell(6, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ROW_HDR
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = YELLOW

    # HBM3 stacked bar
    bar_y = eng_y + 2.8
    bar_h = 0.4
    bar_w = 4.0
    hbm3_comps = [
        ("Row: 2.2", 2.2, CORE_COLOR),
        ("Col: 0.8", 0.8, SENSE_COLOR),
        ("I/O: 1.0", 1.0, IO_COLOR),
        ("PHY: 0.9", 0.9, PHY_COLOR),
        ("Clk: 0.6", 0.6, MED_GRAY),
    ]
    cx = eng_x
    for label, energy, color in hbm3_comps:
        w = bar_w * (energy / 5.5)
        add_rect(slide, cx, bar_y, w, bar_h, color,
                 label, font_size=8, font_color=WHITE, bold=True)
        cx += w

    add_textbox(slide, eng_x, bar_y + 0.45, 4.5, 0.2,
                "HBM3: Row activation is 40% -- I/O is only 18%",
                font_size=9, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

    # Why HBM is efficient
    why_y = 3.7
    add_textbox(slide, 0.5, why_y, 12.5, 0.3,
                "Why HBM Achieves the Lowest Energy per Byte:", font_size=16,
                bold=True, color=GREEN)

    add_multiline_textbox(slide, 0.5, why_y + 0.35, 5.8, 2.5, [
        ("1. Ultra-short I/O path (TSV):", True, WHITE),
        ("   TSV length: ~50 um (vs 20-50 mm for GDDR PCB)", False, LIGHT_GRAY),
        ("   Wire capacitance: ~50 fF (vs 2 pF for PCB)", False, MED_GRAY),
        ("   -> 40x less capacitance = 40x less I/O energy", True, GREEN),
        ("", False, WHITE),
        ("2. Ultra-wide bus (1024-bit per stack):", True, WHITE),
        ("   Low speed per pin (3.2-9.6 Gbps)", False, LIGHT_GRAY),
        ("   -> Simple NRZ signaling (no PAM4 overhead)", False, MED_GRAY),
        ("   -> Simpler PHY per pin, lower per-bit energy", True, GREEN),
        ("", False, WHITE),
        ("3. DRAM process optimization:", True, WHITE),
        ("   HBM uses latest DRAM process node", False, LIGHT_GRAY),
        ("   Smaller cells -> less charge -> less row energy", False, MED_GRAY),
    ], font_size=10)

    # Comparison summary (right column)
    add_multiline_textbox(slide, 6.8, why_y + 0.35, 5.8, 2.5, [
        ("The Trade-Off: Cost vs Energy:", True, ORANGE),
        ("", False, WHITE),
        ("   HBM3:   5.5 pJ/B  |  3.35 TB/s  |  12x cost", True, HBM_COLOR),
        ("   LPDDR5: 8.0 pJ/B  |  25 GB/s    |  2x cost", True, LPDDR_COLOR),
        ("   GDDR6:  15.0 pJ/B |  128 GB/s   |  2x cost", True, GDDR_COLOR),
        ("", False, WHITE),
        ("HBM is 2.7x more energy-efficient than GDDR6", False, LIGHT_GRAY),
        ("but costs 6x more per GB.", False, LIGHT_GRAY),
        ("", False, WHITE),
        ("LPDDR5 is 1.9x more energy-efficient than GDDR6", False, LIGHT_GRAY),
        ("at the same cost tier -- but 100x less bandwidth.", False, LIGHT_GRAY),
        ("", False, WHITE),
        ("For energy-optimized AI (KPU), LPDDR5 at 8 pJ/B", True, YELLOW),
        ("is within 1.45x of HBM3 at 5.5 pJ/B, at 6x lower cost.", True, YELLOW),
    ], font_size=10)

    # Source
    add_textbox(slide, 0.5, 7.1, 12.0, 0.3,
                "Sources: Vogelsang MICRO 2010, Horowitz ISSCC 2014, JEDEC LPDDR5/GDDR6/HBM3 specs, vendor datasheets (Micron, SK hynix, Samsung)",
                font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    return slide


# ============================================================================
# MAIN
# ============================================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1_title(prs)
    build_slide_2_dram_anatomy(prs)
    build_slide_3_lpddr(prs)
    build_slide_4_gddr(prs)
    build_slide_5_hbm(prs)

    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "..", "memory_technology_comparison.pptx")
    output_path = os.path.normpath(output_path)
    prs.save(output_path)
    print(f"Saved presentation to: {output_path}")
    print(f"  5 slides generated:")
    print(f"  1. Title: Why Memory Technology Matters for Energy Comparisons")
    print(f"  2. DRAM Energy Anatomy: The Four Components Every Access Pays")
    print(f"  3. LPDDR5/5X: Package-on-Package, Low-Voltage Signaling")
    print(f"  4. GDDR6/6X/7: High-Speed PCB Signaling, PAM4/PAM3")
    print(f"  5. HBM2/3/3E: TSV Stacking, Silicon Interposer, Wide Bus")


if __name__ == "__main__":
    main()
