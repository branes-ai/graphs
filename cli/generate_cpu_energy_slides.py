#!/usr/bin/env python3
"""
Generate PowerPoint slides explaining CPU energy modeling.

This is the first deck in a series that builds toward the KPU 20x energy
advantage argument. These 5 slides establish the baseline: how energy is
consumed in a classical stored-program CPU pipeline.

Slide 1: Title -- The Energy Problem in Computing
Slide 2: 5-Stage CPU Pipeline Architecture Diagram
Slide 3: Execute Stage Energy (Register Read + ALU + Writeback)
Slide 4: Frontend Energy (I-cache + Decode) -- the dominant cost
Slide 5: Summary -- CPU is frontend-heavy, motivating dataflow architectures
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- Color palette ----------------------------------------------------------
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)  # deep navy
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)  # corporate blue
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xD8)  # lighter teal
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
MED_GRAY    = RGBColor(0x99, 0x99, 0x99)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)  # highlight / warning
RED_ACCENT  = RGBColor(0xE0, 0x40, 0x40)
GREEN       = RGBColor(0x2E, 0xCC, 0x71)
YELLOW      = RGBColor(0xF3, 0xC6, 0x23)

# Pipeline stage colors
IF_COLOR    = RGBColor(0x3A, 0x86, 0xC8)  # blue
ID_COLOR    = RGBColor(0x5E, 0x60, 0xCE)  # purple
EX_COLOR    = RGBColor(0x2E, 0xCC, 0x71)  # green
MEM_COLOR   = RGBColor(0xF3, 0xC6, 0x23)  # yellow
WB_COLOR    = RGBColor(0xE0, 0x40, 0x40)  # red

# Cache / storage colors
CACHE_COLOR = RGBColor(0x48, 0x7D, 0xB3)
RF_COLOR    = RGBColor(0x6C, 0x5C, 0xE7)  # register file purple
ALU_COLOR   = RGBColor(0x2E, 0xCC, 0x71)  # ALU green


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a textbox with styled text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=16, color=WHITE, alignment=PP_ALIGN.LEFT,
                          font_name="Calibri", line_spacing=1.2):
    """Add a textbox with multiple paragraphs (list of (text, bold, color) tuples)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld, clr = item, False, color
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            clr = item[2] if len(item) > 2 else color

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.bold = bld
        p.font.color.rgb = clr
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1.0) * 2)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     text="", font_size=14, font_color=WHITE, bold=False):
    """Add a rounded rectangle shape with centered text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # no border
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_rect(slide, left, top, width, height, fill_color,
             text="", font_size=14, font_color=WHITE, bold=False,
             line_color=None):
    """Add a rectangle shape with centered text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow_right(slide, left, top, width=0.3, height=0.2, color=MED_GRAY):
    """Add a right-pointing arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_arrow_down(slide, left, top, width=0.2, height=0.3, color=MED_GRAY):
    """Add a down-pointing arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=MED_GRAY, width=1.5):
    """Add a line connector."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


# ============================================================================
# SLIDE BUILDERS
# ============================================================================

def build_slide_1_title(prs):
    """Slide 1: Title slide -- The Energy Problem in Computing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide, DARK_BG)

    # Title
    add_textbox(slide, 0.8, 1.0, 11.5, 1.0,
                "The Energy Problem in Computing",
                font_size=36, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, 0.8, 1.8, 11.5, 0.8,
                "Understanding Where Energy Goes in a CPU Pipeline",
                font_size=22, bold=False, color=ACCENT_TEAL,
                alignment=PP_ALIGN.CENTER)

    # Horizontal divider line
    add_line(slide, 2.5, 2.7, 10.5, 2.7, color=ACCENT_BLUE, width=2)

    # Key insight box
    add_rounded_rect(slide, 2.5, 3.2, 8.0, 1.2, RGBColor(0x2A, 0x2A, 0x4E),
                     "", font_size=16)
    add_multiline_textbox(slide, 2.8, 3.3, 7.4, 1.0, [
        ("Key Insight:", True, ORANGE),
        ("The ALU -- where useful computation happens -- consumes only", False, LIGHT_GRAY),
        ("~20% of total operation energy. The other 80% is overhead.", True, WHITE),
    ], font_size=16, alignment=PP_ALIGN.CENTER)

    # Three pillars
    pillar_y = 4.8
    pillar_w = 2.8
    pillars = [
        ("Instruction Frontend", "Fetch + Decode\n~40% of energy\n~66% of transistors", IF_COLOR),
        ("Data Movement", "Register File + Caches\n~30% of energy", RF_COLOR),
        ("Useful Compute", "ALU / FPU\n~20% of energy", ALU_COLOR),
    ]
    for i, (title, body, color) in enumerate(pillars):
        x = 1.0 + i * 3.5
        add_rounded_rect(slide, x, pillar_y, pillar_w, 0.55, color,
                         title, font_size=14, font_color=WHITE, bold=True)
        add_textbox(slide, x, pillar_y + 0.6, pillar_w, 0.9,
                    body, font_size=12, color=LIGHT_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # Footer
    add_textbox(slide, 0.5, 7.0, 12.0, 0.4,
                "Source: Horowitz, ISSCC 2014 -- scaled to modern process nodes (4-7nm)",
                font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    return slide


def build_slide_2_pipeline(prs):
    """Slide 2: 5-Stage CPU Pipeline Architecture Diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Title
    add_textbox(slide, 0.5, 0.2, 12.0, 0.6,
                "The 5-Stage CPU Pipeline: A Stored-Program Machine",
                font_size=28, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 0.5, 0.75, 12.0, 0.4,
                "Random Access Memory Stored Program Machine (Von Neumann Architecture)",
                font_size=14, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    # -- Pipeline stages (horizontal flow) --
    stage_y = 1.6
    stage_h = 0.75
    stage_w = 1.6
    gap = 0.15
    arrow_w = 0.3
    start_x = 1.3

    stages = [
        ("IF", "Instruction\nFetch",  IF_COLOR),
        ("ID", "Instruction\nDecode", ID_COLOR),
        ("EX", "Execute",            EX_COLOR),
        ("MEM", "Memory\nAccess",    MEM_COLOR),
        ("WB", "Write\nBack",        WB_COLOR),
    ]

    for i, (abbr, name, color) in enumerate(stages):
        x = start_x + i * (stage_w + gap + arrow_w)
        add_rounded_rect(slide, x, stage_y, stage_w, stage_h, color,
                         f"{abbr}\n{name}", font_size=12,
                         font_color=WHITE, bold=True)
        # Arrow between stages
        if i < len(stages) - 1:
            ax = x + stage_w + 0.02
            add_arrow_right(slide, ax, stage_y + stage_h/2 - 0.1,
                            width=arrow_w - 0.04, height=0.2, color=LIGHT_GRAY)

    # -- L1 Instruction Cache (above IF stage) --
    icache_x = start_x - 0.1
    icache_y = 2.8
    add_rect(slide, icache_x, icache_y, 1.8, 0.55, CACHE_COLOR,
             "L1 I-Cache\n(32-64 KB)", font_size=11, font_color=WHITE, bold=True)
    # Arrow from I-cache to IF
    add_arrow_down(slide, start_x + 0.7, icache_y - 0.35,
                   width=0.2, height=0.35, color=CACHE_COLOR)

    # -- Register File (below ID/EX, centered) --
    rf_x = start_x + 1 * (stage_w + gap + arrow_w) + 0.6
    rf_y = 2.8
    add_rect(slide, rf_x, rf_y, 2.4, 0.55, RF_COLOR,
             "Register File (RF)\n128-256 entries x 64-bit",
             font_size=11, font_color=WHITE, bold=True)

    # Arrow from RF up to EX (operand read)
    ex_x = start_x + 2 * (stage_w + gap + arrow_w)
    add_line(slide, rf_x + 1.2, rf_y, ex_x + 0.8, stage_y + stage_h,
             color=RF_COLOR, width=1.5)

    # -- ALU block (below EX stage) --
    alu_x = ex_x + 0.15
    alu_y = 3.65
    add_rect(slide, alu_x, alu_y, 1.3, 0.5, ALU_COLOR,
             "ALU / FPU / SIMD", font_size=11, font_color=WHITE, bold=True)
    # Line from EX to ALU
    add_line(slide, ex_x + 0.8, stage_y + stage_h, alu_x + 0.65, alu_y,
             color=ALU_COLOR, width=1.5)

    # -- L1 Data Cache (below MEM stage) --
    mem_x = start_x + 3 * (stage_w + gap + arrow_w)
    dcache_y = 2.8
    add_rect(slide, mem_x - 0.1, dcache_y, 1.8, 0.55, CACHE_COLOR,
             "L1 D-Cache\n(32-64 KB)", font_size=11, font_color=WHITE, bold=True)
    # Arrow from MEM down to D-cache
    add_line(slide, mem_x + 0.8, stage_y + stage_h, mem_x + 0.8, dcache_y,
             color=CACHE_COLOR, width=1.5)

    # -- WB arrow back to RF --
    wb_x = start_x + 4 * (stage_w + gap + arrow_w)
    add_line(slide, wb_x + 0.8, stage_y + stage_h, rf_x + 1.8, rf_y + 0.55,
             color=WB_COLOR, width=1.5)

    # -- Cache hierarchy below --
    hier_y = 3.65
    caches = [
        ("L2 Cache\n(256 KB - 1 MB)", 1.6),
        ("L3 / LLC\n(8 - 32 MB)", 1.6),
        ("DRAM\n(DDR5 / LPDDR5)", 1.6),
    ]
    cache_x = 5.8
    for i, (label, w) in enumerate(caches):
        cy = hier_y + i * 0.7
        darkness = max(0x28, 0x48 - i * 0x10)
        c = RGBColor(darkness, 0x7D - i * 0x15, 0xB3 - i * 0x10)
        add_rect(slide, cache_x, cy, w, 0.55, c,
                 label, font_size=10, font_color=WHITE, bold=False)
        if i > 0:
            add_line(slide, cache_x + 0.8, cy, cache_x + 0.8, cy - 0.15,
                     color=MED_GRAY, width=1)

    # Connect D-cache to L2
    add_line(slide, mem_x + 0.8, dcache_y + 0.55, cache_x + 0.8, hier_y,
             color=MED_GRAY, width=1)

    # -- Annotations --
    # Frontend bracket
    add_textbox(slide, 0.3, 1.35, 1.0, 0.3, "FRONTEND", font_size=10,
                bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
    add_rect(slide, 0.4, 1.55, 0.08, 0.85, ORANGE)  # vertical bar

    # Backend bracket
    add_textbox(slide, 10.5, 1.35, 1.2, 0.3, "BACKEND", font_size=10,
                bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)
    add_rect(slide, 11.5, 1.55, 0.08, 0.85, GREEN)  # vertical bar

    # Key observations
    obs_y = 5.2
    add_textbox(slide, 0.5, obs_y, 12.0, 0.4,
                "Key Architectural Observations",
                font_size=18, bold=True, color=WHITE)
    add_multiline_textbox(slide, 0.5, obs_y + 0.45, 12.0, 2.0, [
        ("  * Every operation requires an instruction fetch + decode cycle (the \"instruction tax\")", False, LIGHT_GRAY),
        ("  * The register file is accessed on BOTH sides: read before EX, write after EX", False, LIGHT_GRAY),
        ("  * WAR/RAW hazards in writeback add scoreboarding and forwarding energy", False, LIGHT_GRAY),
        ("  * ~66% of all transistors serve the frontend (fetch, decode, branch prediction, rename)", True, ORANGE),
    ], font_size=13)

    return slide


def build_slide_3_execute_energy(prs):
    """Slide 3: Execute Stage Energy Breakdown."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Title
    add_textbox(slide, 0.5, 0.2, 12.0, 0.6,
                "Execute Stage Energy: The Useful Work",
                font_size=28, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # -- Diagram: RF Read -> ALU -> RF Write --
    diag_y = 1.2
    box_h = 0.7

    # RF Read
    add_rounded_rect(slide, 1.5, diag_y, 2.0, box_h, RF_COLOR,
                     "Register File\nREAD (2 operands)",
                     font_size=12, font_color=WHITE, bold=True)
    add_arrow_right(slide, 3.55, diag_y + 0.25, 0.4, 0.2, LIGHT_GRAY)

    # ALU
    add_rounded_rect(slide, 4.1, diag_y, 2.0, box_h, ALU_COLOR,
                     "ALU / FPU\nFMA Operation",
                     font_size=12, font_color=WHITE, bold=True)
    add_arrow_right(slide, 6.15, diag_y + 0.25, 0.4, 0.2, LIGHT_GRAY)

    # RF Write
    add_rounded_rect(slide, 6.7, diag_y, 2.0, box_h, WB_COLOR,
                     "Register File\nWRITE (result)",
                     font_size=12, font_color=WHITE, bold=True)

    # Energy equation
    eq_y = 2.15
    add_textbox(slide, 0.5, eq_y, 12.0, 0.35,
                "Energy Equation for the Execute Stage:",
                font_size=16, bold=True, color=ACCENT_TEAL)

    add_textbox(slide, 0.5, eq_y + 0.4, 12.0, 0.4,
                "E_execute = 2 x E_rf_read  +  E_alu  +  E_rf_write",
                font_size=20, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER, font_name="Consolas")

    # -- Table: Energy values by process node --
    table_y = 3.2
    add_textbox(slide, 0.5, table_y - 0.35, 12.0, 0.35,
                "Typical Energy Values by Process Technology (pJ per 32-bit FP32 FMA)",
                font_size=14, bold=True, color=WHITE)

    # Create actual table
    rows, cols = 6, 6
    tbl = slide.shapes.add_table(rows, cols,
                                  Inches(0.8), Inches(table_y),
                                  Inches(11.4), Inches(2.2)).table

    headers = ["Process Node", "RF Read (x2)", "ALU (FMA)", "RF Write", "E_execute Total", "ALU % of Total"]
    data = [
        ["28nm (IoT)",       "2 x 1.8 = 3.6",  "4.0",  "2.0",  "9.6 pJ",   "42%"],
        ["16nm (Auto)",      "2 x 1.2 = 2.4",  "2.7",  "1.3",  "6.4 pJ",   "42%"],
        ["7nm (Datacenter)", "2 x 0.8 = 1.6",  "1.8",  "0.9",  "4.3 pJ",   "42%"],
        ["5nm (Mobile SoC)", "2 x 0.6 = 1.2",  "1.5",  "0.7",  "3.4 pJ",   "44%"],
        ["4nm (GPU/AI)",     "2 x 0.5 = 1.0",  "1.3",  "0.6",  "2.9 pJ",   "45%"],
    ]

    # Style header
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE

    # Style data rows
    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x4E)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)

    # Key takeaway
    tk_y = 5.6
    add_rounded_rect(slide, 1.0, tk_y, 11.0, 0.9, RGBColor(0x2A, 0x2A, 0x4E))
    add_multiline_textbox(slide, 1.3, tk_y + 0.05, 10.4, 0.8, [
        ("Takeaway:", True, ORANGE),
        ("Even in the execute stage alone, the ALU is less than half the energy.", False, LIGHT_GRAY),
        ("Register file access (read 2 operands + write 1 result) costs more than the computation itself.", False, WHITE),
    ], font_size=13)

    # Source
    add_textbox(slide, 0.5, 7.0, 12.0, 0.4,
                "Values derived from Horowitz ISSCC 2014, scaled by TSMC process node capacitance ratios",
                font_size=10, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    return slide


def build_slide_4_frontend_energy(prs):
    """Slide 4: Frontend Energy -- the dominant cost."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Title
    add_textbox(slide, 0.5, 0.2, 12.0, 0.6,
                "Frontend Energy: The Hidden Dominant Cost",
                font_size=28, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # -- Frontend diagram --
    diag_y = 1.15
    box_h = 0.65

    add_rounded_rect(slide, 0.8, diag_y, 1.8, box_h, CACHE_COLOR,
                     "L1 I-Cache\nRead", font_size=12, font_color=WHITE, bold=True)
    add_arrow_right(slide, 2.65, diag_y + 0.22, 0.35, 0.2, LIGHT_GRAY)
    add_rounded_rect(slide, 3.1, diag_y, 1.8, box_h, ID_COLOR,
                     "Instruction\nDecode", font_size=12, font_color=WHITE, bold=True)
    add_arrow_right(slide, 4.95, diag_y + 0.22, 0.35, 0.2, LIGHT_GRAY)
    add_rounded_rect(slide, 5.4, diag_y, 1.8, box_h, RGBColor(0x88, 0x55, 0xBB),
                     "Branch\nPredict", font_size=12, font_color=WHITE, bold=True)
    add_arrow_right(slide, 7.25, diag_y + 0.22, 0.35, 0.2, LIGHT_GRAY)
    add_rounded_rect(slide, 7.7, diag_y, 1.8, box_h, RGBColor(0x44, 0x77, 0xAA),
                     "Register\nRename", font_size=12, font_color=WHITE, bold=True)

    # Bracket label
    add_textbox(slide, 0.8, diag_y + 0.72, 8.7, 0.3,
                "--- FRONTEND (instruction overhead, paid for EVERY operation) ---",
                font_size=11, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

    # Energy equation
    eq_y = 2.25
    add_textbox(slide, 0.5, eq_y, 12.0, 0.35,
                "Frontend Energy Equation:",
                font_size=16, bold=True, color=ACCENT_TEAL)
    add_textbox(slide, 0.5, eq_y + 0.35, 12.0, 0.4,
                "E_frontend = E_icache_read  +  E_decode  +  E_branch_predict  +  E_rename",
                font_size=18, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER, font_name="Consolas")

    # -- Comparison table: Frontend vs Execute --
    table_y = 3.2
    add_textbox(slide, 0.5, table_y - 0.35, 12.0, 0.35,
                "Frontend vs Execute Energy Comparison (per scalar FP32 operation, 7nm)",
                font_size=14, bold=True, color=WHITE)

    rows, cols = 8, 4
    tbl = slide.shapes.add_table(rows, cols,
                                  Inches(1.5), Inches(table_y),
                                  Inches(10.0), Inches(2.6)).table

    headers = ["Component", "Energy (pJ)", "Category", "% of Total"]
    tdata = [
        ["L1 I-Cache Read (32B line)", "5.0",  "Frontend",  ""],
        ["Instruction Decode",         "2.0",  "Frontend",  ""],
        ["Branch Prediction",          "1.5",  "Frontend",  ""],
        ["Register Rename / Scoreboard","1.0", "Frontend",  ""],
        ["Frontend Subtotal",          "9.5",  "FRONTEND",  "63%"],
        ["Execute Stage (from Slide 3)","4.3", "BACKEND",   "29%"],
        ["Pipeline/Clock/Wire",        "1.2",  "Overhead",  "8%"],
    ]

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE

    for i, row_data in enumerate(tdata):
        is_subtotal = ("Subtotal" in row_data[0] or "Execute" in row_data[0]
                       or "Pipeline" in row_data[0])
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.bold = is_subtotal
                p.alignment = PP_ALIGN.CENTER
                if "FRONTEND" in row_data[2] and is_subtotal:
                    p.font.color.rgb = ORANGE
                elif "BACKEND" in row_data[2]:
                    p.font.color.rgb = GREEN
                else:
                    p.font.color.rgb = WHITE
            if is_subtotal:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x4E)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)

    # Transistor allocation visual
    vis_y = 6.0
    add_textbox(slide, 0.5, vis_y, 12.0, 0.35,
                "Transistor Budget Allocation (typical OoO CPU):",
                font_size=14, bold=True, color=WHITE)

    # Stacked bar
    bar_y = vis_y + 0.4
    bar_h = 0.45
    total_w = 10.0
    start_x = 1.5
    # Frontend: 66%, Backend: 20%, Cache: 14%
    fe_w = total_w * 0.66
    be_w = total_w * 0.20
    ca_w = total_w * 0.14

    add_rect(slide, start_x, bar_y, fe_w, bar_h, ORANGE,
             "Frontend: 66%", font_size=12, font_color=WHITE, bold=True)
    add_rect(slide, start_x + fe_w, bar_y, be_w, bar_h, GREEN,
             "Backend: 20%", font_size=11, font_color=WHITE, bold=True)
    add_rect(slide, start_x + fe_w + be_w, bar_y, ca_w, bar_h, CACHE_COLOR,
             "Caches: 14%", font_size=10, font_color=WHITE, bold=True)

    # Bottom insight
    add_multiline_textbox(slide, 0.5, 7.0, 12.0, 0.5, [
        ("The CPU spends 2/3 of its energy and transistors on figuring out WHAT to compute, not computing it.", True, ORANGE),
    ], font_size=14, alignment=PP_ALIGN.CENTER)

    return slide


def build_slide_5_summary(prs):
    """Slide 5: Summary -- CPU energy is frontend-dominated."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Title
    add_textbox(slide, 0.5, 0.2, 12.0, 0.6,
                "Total CPU Energy per Operation: The Full Picture",
                font_size=28, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # Full energy equation
    add_textbox(slide, 0.3, 1.0, 12.5, 0.35,
                "Complete Energy Equation:",
                font_size=16, bold=True, color=ACCENT_TEAL)
    add_textbox(slide, 0.3, 1.35, 12.5, 0.7,
                "E_total = E_frontend + E_execute + E_memory + E_static",
                font_size=22, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER, font_name="Consolas")
    add_textbox(slide, 0.3, 1.8, 12.5, 0.4,
                "= (I-cache + Decode + Branch + Rename) + (2*RF_read + ALU + RF_write) + (Cache hierarchy) + (Leakage)",
                font_size=13, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER, font_name="Consolas")

    # Stacked energy bar chart (horizontal)
    bar_y = 2.5
    add_textbox(slide, 0.5, bar_y, 12.0, 0.3,
                "Energy Breakdown per Scalar FP32 Operation at 7nm (pJ):",
                font_size=14, bold=True, color=WHITE)

    bar_y2 = bar_y + 0.4
    bar_h = 0.6
    total_e = 15.0  # total pJ for scalar
    total_w = 10.0
    start_x = 1.5
    components = [
        ("I-cache: 5.0", 5.0, IF_COLOR),
        ("Decode: 2.0", 2.0, ID_COLOR),
        ("Branch: 1.5", 1.5, RGBColor(0x88, 0x55, 0xBB)),
        ("Rename: 1.0", 1.0, RGBColor(0x44, 0x77, 0xAA)),
        ("RF Read: 1.6", 1.6, RF_COLOR),
        ("ALU: 1.8", 1.8, ALU_COLOR),
        ("RF Write: 0.9", 0.9, WB_COLOR),
        ("Wire: 1.2", 1.2, MED_GRAY),
    ]
    cx = start_x
    for label, energy, color in components:
        w = total_w * (energy / total_e)
        font_sz = 9 if w < 0.8 else 10
        add_rect(slide, cx, bar_y2, w, bar_h, color,
                 label, font_size=font_sz, font_color=WHITE, bold=False)
        cx += w

    # Bracket labels
    fe_w = total_w * (9.5 / total_e)
    ex_w = total_w * (4.3 / total_e)
    add_textbox(slide, start_x, bar_y2 + bar_h + 0.02, fe_w, 0.25,
                "Frontend: 9.5 pJ (63%)", font_size=11, bold=True, color=ORANGE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, start_x + fe_w, bar_y2 + bar_h + 0.02, ex_w, 0.25,
                "Execute: 4.3 pJ (29%)", font_size=11, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)

    # -- The argument: what if we eliminate the frontend? --
    arg_y = 4.0
    add_line(slide, 1.0, arg_y, 12.0, arg_y, color=ACCENT_BLUE, width=2)

    add_textbox(slide, 0.5, arg_y + 0.15, 12.0, 0.4,
                "The Architectural Question: What If We Eliminate the Frontend?",
                font_size=20, bold=True, color=ACCENT_TEAL,
                alignment=PP_ALIGN.CENTER)

    # Three architecture comparison boxes
    arch_y = 4.8
    arch_w = 3.5
    arch_h = 1.8

    # CPU
    add_rounded_rect(slide, 0.7, arch_y, arch_w, 0.5,
                     RGBColor(0x44, 0x44, 0x66),
                     "CPU (Stored Program)", font_size=13,
                     font_color=WHITE, bold=True)
    add_multiline_textbox(slide, 0.7, arch_y + 0.55, arch_w, 1.2, [
        ("~15 pJ / scalar FP32 op", True, RED_ACCENT),
        ("Frontend: 9.5 pJ (63%)", False, LIGHT_GRAY),
        ("Execute: 4.3 pJ (29%)", False, LIGHT_GRAY),
        ("Overhead: 1.2 pJ (8%)", False, LIGHT_GRAY),
        ("Every op pays full instruction tax", False, MED_GRAY),
    ], font_size=11)

    # GPU
    add_rounded_rect(slide, 4.5, arch_y, arch_w, 0.5,
                     RGBColor(0x44, 0x44, 0x66),
                     "GPU (SIMT Parallel)", font_size=13,
                     font_color=WHITE, bold=True)
    add_multiline_textbox(slide, 4.5, arch_y + 0.55, arch_w, 1.2, [
        ("~1.65 pJ / FMA (TensorCore)", True, YELLOW),
        ("SIMT amortizes fetch over 32 threads", False, LIGHT_GRAY),
        ("But adds coherence + scheduling", False, LIGHT_GRAY),
        ("Operand collector still ~1.5x ALU", False, LIGHT_GRAY),
        ("Still a stored-program machine", False, MED_GRAY),
    ], font_size=11)

    # KPU (preview)
    add_rounded_rect(slide, 8.3, arch_y, arch_w, 0.5,
                     ACCENT_BLUE,
                     "KPU (Domain Flow)", font_size=13,
                     font_color=WHITE, bold=True)
    add_multiline_textbox(slide, 8.3, arch_y + 0.55, arch_w, 1.2, [
        ("Target: <0.1 pJ / MAC", True, GREEN),
        ("NO instruction fetch/decode", False, LIGHT_GRAY),
        ("Data-driven token execution", False, LIGHT_GRAY),
        ("Operands arrive at PE, not fetched", False, LIGHT_GRAY),
        ("Eliminates the 63% frontend tax", False, MED_GRAY),
    ], font_size=11)

    # Bottom line
    add_rounded_rect(slide, 1.5, 7.0, 10.0, 0.5, RGBColor(0x2A, 0x2A, 0x4E))
    add_textbox(slide, 1.5, 7.05, 10.0, 0.4,
                "Eliminating the instruction frontend is the key to achieving 10-20x energy efficiency",
                font_size=16, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)

    return slide


# ============================================================================
# MAIN
# ============================================================================

def main():
    prs = Presentation()
    # Set widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1_title(prs)
    build_slide_2_pipeline(prs)
    build_slide_3_execute_energy(prs)
    build_slide_4_frontend_energy(prs)
    build_slide_5_summary(prs)

    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "..", "cpu_energy_pipeline.pptx")
    output_path = os.path.normpath(output_path)
    prs.save(output_path)
    print(f"Saved presentation to: {output_path}")
    print(f"  5 slides generated:")
    print(f"  1. Title: The Energy Problem in Computing")
    print(f"  2. 5-Stage CPU Pipeline Architecture Diagram")
    print(f"  3. Execute Stage Energy (RF Read + ALU + RF Write)")
    print(f"  4. Frontend Energy: The Hidden Dominant Cost")
    print(f"  5. Summary: Total Energy & Architecture Comparison")


if __name__ == "__main__":
    main()
