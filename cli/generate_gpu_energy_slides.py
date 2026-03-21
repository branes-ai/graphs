#!/usr/bin/env python3
"""
Generate PowerPoint slides explaining GPU energy modeling.

This is the second deck in the series building toward the KPU 20x energy
advantage argument. These 5 slides show how GPUs trade instruction frontend
overhead for massive data movement energy to support thread concurrency.

Slide 1: Title -- GPU: Trading Frontend for Data Movement
Slide 2: SM Architecture -- 4 Partitions, Banked Register File, Operand Collectors
Slide 3: The 64K Register File -- Concurrent Banked SRAM Architecture
Slide 4: GPU Execution Cycle Energy Equation
Slide 5: CPU vs GPU Energy Comparison -- Different overhead, same problem
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# -- Color palette (matches CPU deck) ---------------------------------------
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xD8)
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
MED_GRAY    = RGBColor(0x99, 0x99, 0x99)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)
RED_ACCENT  = RGBColor(0xE0, 0x40, 0x40)
GREEN       = RGBColor(0x2E, 0xCC, 0x71)
YELLOW      = RGBColor(0xF3, 0xC6, 0x23)

# GPU-specific colors
NVIDIA_GREEN  = RGBColor(0x76, 0xB9, 0x00)
SM_COLOR      = RGBColor(0x3A, 0x5F, 0x8A)
PARTITION_CLR = RGBColor(0x44, 0x66, 0x88)
WARP_SCHED    = RGBColor(0x8B, 0x5C, 0xF6)  # purple for scheduler
SCOREBOARD    = RGBColor(0xE0, 0x40, 0x40)   # red for hazard detection
COLLECTOR_CLR = RGBColor(0xF3, 0x96, 0x23)   # orange for operand collector
BANK_COLOR    = RGBColor(0x48, 0x7D, 0xB3)   # blue for register banks
CROSSBAR_CLR  = RGBColor(0xCC, 0x55, 0x77)   # pink for crossbar
RF_COLOR      = RGBColor(0x6C, 0x5C, 0xE7)   # purple for register file
ALU_COLOR     = RGBColor(0x2E, 0xCC, 0x71)   # green for ALU
WB_COLOR      = RGBColor(0xE0, 0x40, 0x40)


# ============================================================================
# Helper functions (same API as CPU deck)
# ============================================================================

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=16, color=WHITE, alignment=PP_ALIGN.LEFT,
                          font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld, clr = item, False, color
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            clr = item[2] if len(item) > 2 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.bold = bld
        p.font.color.rgb = clr
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1.0) * 2)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     text="", font_size=14, font_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_rect(slide, left, top, width, height, fill_color,
             text="", font_size=14, font_color=WHITE, bold=False,
             line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow_right(slide, left, top, width=0.3, height=0.2, color=MED_GRAY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_arrow_down(slide, left, top, width=0.2, height=0.3, color=MED_GRAY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=MED_GRAY, width=1.5):
    connector = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


# ============================================================================
# SLIDE 1: Title
# ============================================================================

def build_slide_1_title(prs):
    """GPU Energy: Trading Frontend Overhead for Data Movement Overhead."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.8, 11.5, 1.0,
                "GPU Energy: From Instruction Tax to Data Movement Tax",
                font_size=34, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 0.8, 1.7, 11.5, 0.6,
                "How Thread Concurrency Hides Latency but Shifts the Energy Problem",
                font_size=22, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    add_line(slide, 2.5, 2.5, 10.5, 2.5, color=NVIDIA_GREEN, width=2)

    # Key insight box
    add_rounded_rect(slide, 2.0, 2.9, 9.0, 1.4, RGBColor(0x2A, 0x2A, 0x4E))
    add_multiline_textbox(slide, 2.3, 3.0, 8.4, 1.2, [
        ("The GPU Bargain:", True, NVIDIA_GREEN),
        ("GPUs reduce instruction frontend overhead by amortizing one instruction", False, LIGHT_GRAY),
        ("across 32 threads (a warp). But to hide memory latency, they run hundreds", False, LIGHT_GRAY),
        ("of concurrent threads -- each needing its own registers.", False, LIGHT_GRAY),
        ("The result: 64K registers per SM, fed through banked SRAM with operand", True, ORANGE),
        ("collectors. This data movement infrastructure is always active.", True, ORANGE),
    ], font_size=15, alignment=PP_ALIGN.CENTER)

    # Two-column comparison
    col_y = 4.8
    col_w = 4.5

    # CPU column
    add_rounded_rect(slide, 1.5, col_y, col_w, 0.5, RGBColor(0x44, 0x44, 0x66),
                     "CPU: Frontend-Heavy", font_size=15, bold=True)
    add_multiline_textbox(slide, 1.5, col_y + 0.55, col_w, 1.2, [
        ("~15 pJ per scalar FP32 op", True, RED_ACCENT),
        ("63% energy in instruction frontend", False, LIGHT_GRAY),
        ("I-cache, decode, branch predict, rename", False, MED_GRAY),
        ("Small register file (16-64 entries)", False, MED_GRAY),
    ], font_size=12)

    # GPU column
    add_rounded_rect(slide, 7.0, col_y, col_w, 0.5, NVIDIA_GREEN,
                     "GPU: Data-Movement-Heavy", font_size=15,
                     font_color=BLACK, bold=True)
    add_multiline_textbox(slide, 7.0, col_y + 0.55, col_w, 1.2, [
        ("~1.65 pJ per TensorCore MAC", True, NVIDIA_GREEN),
        ("Frontend amortized over 32 threads", False, LIGHT_GRAY),
        ("But: 64K registers, banked SRAM,", False, ORANGE),
        ("operand collectors, crossbar routing", False, ORANGE),
    ], font_size=12)

    # Arrow between columns
    add_textbox(slide, 5.7, col_y + 0.6, 1.5, 0.5,
                "-->", font_size=28, bold=True, color=ACCENT_TEAL,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 0.5, 7.0, 12.0, 0.4,
                "Both architectures are stored-program machines -- operands must be explicitly fetched",
                font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    return slide


# ============================================================================
# SLIDE 2: SM Architecture with 4 Partitions
# ============================================================================

def build_slide_2_sm_architecture(prs):
    """SM Architecture: 4 Partitions with independent operand delivery."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.15, 12.5, 0.5,
                "Streaming Multiprocessor (SM) Architecture: H100-Class",
                font_size=26, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 0.3, 0.6, 12.5, 0.4,
                "Each SM is divided into 4 independent partitions -- each a complete SIMD execution pipeline",
                font_size=13, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    # -- Draw 4 partitions side by side --
    part_start_x = 0.5
    part_w = 2.9
    part_gap = 0.25
    part_y = 1.1

    for p_idx in range(4):
        px = part_start_x + p_idx * (part_w + part_gap)

        # Partition border
        add_rect(slide, px, part_y, part_w, 3.75,
                 RGBColor(0x22, 0x22, 0x3E), "",
                 line_color=PARTITION_CLR)
        add_textbox(slide, px, part_y + 0.02, part_w, 0.25,
                    f"Partition {p_idx}", font_size=11, bold=True,
                    color=PARTITION_CLR, alignment=PP_ALIGN.CENTER)

        # Warp Scheduler
        add_rounded_rect(slide, px + 0.1, part_y + 0.3, part_w - 0.2, 0.35,
                         WARP_SCHED, "Warp Scheduler",
                         font_size=10, bold=True)

        # Scoreboard
        add_rounded_rect(slide, px + 0.1, part_y + 0.72, part_w - 0.2, 0.32,
                         SCOREBOARD, "Scoreboard (RAW/WAW/WAR)",
                         font_size=9, bold=True)

        # Operand Collector
        add_rounded_rect(slide, px + 0.1, part_y + 1.12, part_w - 0.2, 0.38,
                         COLLECTOR_CLR, "Operand Collector\n(gather + buffer)",
                         font_size=9, font_color=BLACK, bold=True)

        # Register Bank
        add_rounded_rect(slide, px + 0.1, part_y + 1.58, part_w - 0.2, 0.35,
                         BANK_COLOR, "Register Bank (16K regs)",
                         font_size=9, bold=True)

        # Crossbar
        add_rounded_rect(slide, px + 0.1, part_y + 2.0, part_w - 0.2, 0.3,
                         CROSSBAR_CLR, "Crossbar (32-wide routing)",
                         font_size=9, bold=True)

        # CUDA Cores
        add_rounded_rect(slide, px + 0.1, part_y + 2.38, 1.25, 0.35,
                         ALU_COLOR, "32 CUDA\nCores",
                         font_size=9, font_color=BLACK, bold=True)

        # TensorCore
        add_rounded_rect(slide, px + 1.5, part_y + 2.38, 1.25, 0.35,
                         NVIDIA_GREEN, "1 Tensor\nCore (64 MACs)",
                         font_size=8, font_color=BLACK, bold=True)

        # Writeback arrow
        add_rounded_rect(slide, px + 0.1, part_y + 2.82, part_w - 0.2, 0.3,
                         WB_COLOR, "Result Writeback",
                         font_size=9, bold=True)

        # Back arrow label
        add_textbox(slide, px + part_w - 0.7, part_y + 3.15, 0.6, 0.25,
                    "^ RF", font_size=8, bold=True, color=WB_COLOR,
                    alignment=PP_ALIGN.CENTER)

    # -- Shared resources below partitions --
    shared_y = 5.05

    # L1 / Shared Memory
    add_rect(slide, 0.5, shared_y, 12.0, 0.45, RGBColor(0x33, 0x55, 0x77),
             "Shared L1 Cache / Shared Memory -- 192 KB (configurable split)",
             font_size=12, bold=True)

    # L2
    add_rect(slide, 0.5, shared_y + 0.55, 12.0, 0.4,
             RGBColor(0x28, 0x45, 0x65),
             "L2 Cache -- 50 MB (shared across all 132 SMs)",
             font_size=12, bold=True)

    # HBM
    add_rect(slide, 0.5, shared_y + 1.05, 12.0, 0.4,
             RGBColor(0x1E, 0x35, 0x55),
             "HBM3 -- 80 GB (3.35 TB/s bandwidth)",
             font_size=12, bold=True)

    # Key numbers
    add_textbox(slide, 0.5, 6.6, 12.5, 0.3,
                "Key SM Numbers (H100)",
                font_size=14, bold=True, color=WHITE)

    add_multiline_textbox(slide, 0.5, 6.9, 12.5, 0.8, [
        ("  64K registers per SM (256 KB)  |  2048 max resident threads  |  64 max warps  |  "
         "128 CUDA cores + 4 TensorCores  |  32 register banks", False, LIGHT_GRAY),
        ("  Every cycle: 4 warp schedulers issue 4 instructions -> feeding 128 CUDA cores or 4 TensorCores simultaneously",
         True, ORANGE),
    ], font_size=11)

    return slide


# ============================================================================
# SLIDE 3: The 64K Register File -- Banked SRAM Architecture
# ============================================================================

def build_slide_3_register_file(prs):
    """The 64K Register File: Concurrent Banked SRAM with Operand Collectors."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.15, 12.5, 0.5,
                "The 64K Register File: Banked SRAM Feeding the SIMD Engine",
                font_size=26, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 0.3, 0.6, 12.5, 0.35,
                "256 KB of register storage per SM, organized as 32 banks to support concurrent warp access",
                font_size=13, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    # -- LEFT SIDE: Register Bank Diagram --
    # Show the banked SRAM structure
    bank_start_x = 0.5
    bank_start_y = 1.2
    bank_w = 0.65
    bank_h = 0.32
    bank_gap = 0.04
    banks_per_row = 8

    add_textbox(slide, bank_start_x, bank_start_y - 0.25, 6.0, 0.25,
                "32 Register Banks (8 KB SRAM each)",
                font_size=13, bold=True, color=BANK_COLOR)

    for b in range(32):
        row = b // banks_per_row
        col = b % banks_per_row
        bx = bank_start_x + col * (bank_w + bank_gap)
        by = bank_start_y + row * (bank_h + bank_gap)
        add_rect(slide, bx, by, bank_w, bank_h, BANK_COLOR,
                 f"Bank {b}", font_size=7, font_color=WHITE, bold=False)

    # Bank conflict explanation
    bank_end_y = bank_start_y + 4 * (bank_h + bank_gap)
    add_textbox(slide, bank_start_x, bank_end_y + 0.05, 5.8, 0.25,
                "Thread 0 -> Bank 0, Thread 1 -> Bank 1, ... Thread 31 -> Bank 31",
                font_size=10, color=LIGHT_GRAY)
    add_textbox(slide, bank_start_x, bank_end_y + 0.3, 5.8, 0.4,
                "Bank Conflict: When 2+ threads in a warp access the same bank,\n"
                "accesses are serialized -> variable latency -> needs buffering",
                font_size=10, color=ORANGE, bold=True)

    # -- RIGHT SIDE: Operand Delivery Pipeline --
    pipe_x = 6.5
    pipe_y = 1.2
    pipe_w = 5.8
    pipe_h = 0.42
    pipe_gap = 0.08

    add_textbox(slide, pipe_x, pipe_y - 0.25, pipe_w, 0.25,
                "Operand Delivery Pipeline (per warp instruction)",
                font_size=13, bold=True, color=COLLECTOR_CLR)

    steps = [
        ("1. Scoreboard Lookup", "Check RAW/WAW/WAR hazards for this warp",
         SCOREBOARD, "0.3 pJ"),
        ("2. Register Address Gen", "Decode src1, src2, dst register IDs (x32 threads)",
         RGBColor(0x55, 0x77, 0x99), "0.6 pJ"),
        ("3. Bank Arbitration", "Detect & resolve bank conflicts across 32 threads",
         RGBColor(0x44, 0x77, 0xAA), "0.3 pJ"),
        ("4. SRAM Read (32 banks)", "Read 32 x 32-bit operands from banked register file",
         BANK_COLOR, "0.5 pJ"),
        ("5. Operand Collector", "Buffer operands until ALL 32 are ready (variable latency)",
         COLLECTOR_CLR, "0.8 pJ"),
        ("6. Crossbar Routing", "Route 32 operands from collector to SIMD lanes",
         CROSSBAR_CLR, "0.4 pJ"),
        ("7. SIMD ALU Execute", "32 FP32 FMA operations in parallel",
         ALU_COLOR, "0.7 pJ/op"),
        ("8. Result Writeback", "Route 32 results back through crossbar to register banks",
         WB_COLOR, "0.3 pJ + 0.5 pJ"),
    ]

    for i, (name, desc, color, energy) in enumerate(steps):
        sy = pipe_y + i * (pipe_h + pipe_gap)

        # Step box
        add_rounded_rect(slide, pipe_x, sy, 3.5, pipe_h, color,
                         name, font_size=10, bold=True)
        # Description
        add_textbox(slide, pipe_x + 3.6, sy + 0.02, 2.2, pipe_h - 0.04,
                    desc, font_size=8, color=LIGHT_GRAY)

    # Energy column on far right
    add_textbox(slide, 12.1, pipe_y - 0.25, 0.9, 0.25,
                "Energy", font_size=10, bold=True, color=ACCENT_TEAL,
                alignment=PP_ALIGN.CENTER)
    for i, (_, _, _, energy) in enumerate(steps):
        sy = pipe_y + i * (pipe_h + pipe_gap)
        clr = GREEN if "ALU" in steps[i][0] else ORANGE
        add_textbox(slide, 12.0, sy + 0.05, 1.0, pipe_h - 0.04,
                    energy, font_size=9, bold=True, color=clr,
                    alignment=PP_ALIGN.CENTER)

    # -- Bottom insight --
    insight_y = 5.45
    add_rounded_rect(slide, 0.5, insight_y, 12.0, 0.7, RGBColor(0x2A, 0x2A, 0x4E))
    add_multiline_textbox(slide, 0.8, insight_y + 0.05, 11.4, 0.6, [
        ("The operand collector exists because bank conflicts create variable latency.",
         True, COLLECTOR_CLR),
        ("It buffers partial results until ALL 32 thread operands are ready, "
         "then issues them as a unit into the SIMD engine.", False, LIGHT_GRAY),
        ("This collector + the entire banked SRAM infrastructure is active on EVERY cycle, "
         "for EVERY warp instruction.", True, ORANGE),
    ], font_size=12)

    # Context numbers
    add_multiline_textbox(slide, 0.5, 6.35, 12.5, 0.8, [
        ("Why 64K registers?  2048 threads x 32 registers/thread = 65,536 registers. "
         "All threads must be resident to hide memory latency (100+ cycle HBM access).",
         False, LIGHT_GRAY),
        ("The GPU cannot reduce the register file -- it IS the latency hiding mechanism. "
         "But powering 256 KB of banked SRAM every cycle is the energy cost of concurrency.",
         True, YELLOW),
    ], font_size=11)

    return slide


# ============================================================================
# SLIDE 4: GPU Energy Equation
# ============================================================================

def build_slide_4_energy_equation(prs):
    """GPU Execution Cycle Energy: The Analytical Equation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.15, 12.5, 0.5,
                "GPU Energy Equation: Per-MAC Cost at the SM Level",
                font_size=26, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # -- Full energy equation --
    eq_y = 0.8
    add_textbox(slide, 0.3, eq_y, 12.5, 0.3,
                "Energy per CUDA Core MAC (amortized over 32-thread warp):",
                font_size=14, bold=True, color=ACCENT_TEAL)

    add_textbox(slide, 0.3, eq_y + 0.35, 12.5, 0.35,
                "E_mac = (E_sched + E_scoreboard + E_addr_gen + E_bank_arb + E_sram_read + E_collector + E_crossbar) / 32  +  E_alu  +  (E_result_route + E_sram_write) / 32",
                font_size=13, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER, font_name="Consolas")

    # Simplify
    add_textbox(slide, 0.3, eq_y + 0.8, 12.5, 0.3,
                "Simplified:",
                font_size=14, bold=True, color=ACCENT_TEAL)
    add_textbox(slide, 0.3, eq_y + 1.1, 12.5, 0.35,
                "E_mac  =  E_overhead/32  +  E_alu  +  E_writeback/32",
                font_size=20, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER, font_name="Consolas")

    # Bracket labels
    add_textbox(slide, 1.5, eq_y + 1.5, 3.5, 0.25,
                "\"Concurrency Tax\"", font_size=12, bold=True,
                color=ORANGE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 5.5, eq_y + 1.5, 2.5, 0.25,
                "Useful Work", font_size=12, bold=True,
                color=GREEN, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 8.5, eq_y + 1.5, 3.0, 0.25,
                "Result Delivery", font_size=12, bold=True,
                color=WB_COLOR, alignment=PP_ALIGN.CENTER)

    # -- Numerical evaluation table --
    table_y = 2.8
    add_textbox(slide, 0.5, table_y - 0.35, 12.0, 0.3,
                "Numerical Evaluation: CUDA Core Path at 4nm (H100-class)",
                font_size=14, bold=True, color=WHITE)

    rows, cols = 10, 5
    tbl = slide.shapes.add_table(rows, cols,
                                  Inches(0.8), Inches(table_y),
                                  Inches(11.5), Inches(3.0)).table

    headers = ["Component", "Per-Warp Energy (pJ)", "Per-MAC (pJ)", "Category", "% of Total"]

    # Data: per-warp values from gpu.py at 4nm (process_scale = 1.0)
    # Total per-warp overhead = sum of steps 1-6 + 8
    # Per-MAC = per-warp / 32 (for overhead) or direct (for ALU)
    data = [
        ["Warp Scheduler",       "0.50",  "0.016",  "Control",     ""],
        ["Scoreboard Lookup",    "0.30",  "0.009",  "Control",     ""],
        ["Register Addr Gen (x3)","0.60", "0.019",  "Operand Fetch",""],
        ["Bank Arbitration",     "0.30",  "0.009",  "Operand Fetch",""],
        ["SRAM Read (2 src)",    "0.50",  "0.016",  "Operand Fetch",""],
        ["Operand Collector",    "0.80",  "0.025",  "Operand Fetch",""],
        ["Crossbar Routing",     "0.40",  "0.013",  "Operand Fetch",""],
        ["CUDA Core ALU (FMA)",  "---",   "1.105",  "Compute",     ""],
        ["Result Route + Write", "0.80",  "0.025",  "Writeback",   ""],
    ]

    # Calculate totals
    overhead_per_mac = sum(float(r[2]) for r in data if r[3] != "Compute")
    alu_per_mac = 1.105
    total_per_mac = overhead_per_mac + alu_per_mac

    # Fill in percentages
    for row in data:
        val = float(row[2])
        row[4] = f"{val/total_per_mac*100:.1f}%"

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE

    cat_colors = {
        "Control": WARP_SCHED,
        "Operand Fetch": COLLECTOR_CLR,
        "Compute": ALU_COLOR,
        "Writeback": WB_COLOR,
    }

    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.alignment = PP_ALIGN.CENTER
                if j == 3:
                    p.font.color.rgb = cat_colors.get(val, WHITE)
                    p.font.bold = True
                elif row_data[3] == "Compute":
                    p.font.color.rgb = GREEN
                    p.font.bold = True
                else:
                    p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = (RGBColor(0x2A, 0x2A, 0x4E) if i % 2 == 0
                                        else RGBColor(0x22, 0x22, 0x3E))

    # -- Summary bar --
    bar_y = 5.95
    add_textbox(slide, 0.5, bar_y, 12.0, 0.3,
                f"Total Energy per CUDA Core MAC at 4nm:  {total_per_mac:.2f} pJ",
                font_size=16, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # Stacked bar
    bar_y2 = bar_y + 0.35
    bar_h = 0.45
    total_w = 10.0
    start_x = 1.5

    control_e = 0.016 + 0.009  # sched + scoreboard
    fetch_e = 0.019 + 0.009 + 0.016 + 0.025 + 0.013  # addr + arb + sram + collector + xbar
    compute_e = alu_per_mac
    wb_e = 0.025

    total_e = control_e + fetch_e + compute_e + wb_e
    c_w = total_w * (control_e / total_e)
    f_w = total_w * (fetch_e / total_e)
    a_w = total_w * (compute_e / total_e)
    w_w = total_w * (wb_e / total_e)

    cx = start_x
    add_rect(slide, cx, bar_y2, c_w, bar_h, WARP_SCHED,
             f"Ctrl: {control_e:.3f}", font_size=8, bold=True)
    cx += c_w
    add_rect(slide, cx, bar_y2, f_w, bar_h, COLLECTOR_CLR,
             f"Operand Fetch: {fetch_e:.3f}", font_size=9,
             font_color=BLACK, bold=True)
    cx += f_w
    add_rect(slide, cx, bar_y2, a_w, bar_h, ALU_COLOR,
             f"ALU: {compute_e:.3f} pJ", font_size=10,
             font_color=BLACK, bold=True)
    cx += a_w
    add_rect(slide, cx, bar_y2, w_w, bar_h, WB_COLOR,
             f"WB: {wb_e:.3f}", font_size=8, bold=True)

    # Percentage annotations
    overhead_pct = (control_e + fetch_e + wb_e) / total_e * 100
    compute_pct = compute_e / total_e * 100
    add_textbox(slide, start_x, bar_y2 + bar_h + 0.02,
                total_w * ((control_e + fetch_e) / total_e), 0.25,
                f"Overhead: {overhead_pct:.0f}%", font_size=11, bold=True,
                color=ORANGE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, start_x + total_w * ((control_e + fetch_e) / total_e),
                bar_y2 + bar_h + 0.02, total_w * (compute_e / total_e), 0.25,
                f"Useful Compute: {compute_pct:.0f}%", font_size=11, bold=True,
                color=GREEN, alignment=PP_ALIGN.CENTER)

    # Insight
    add_textbox(slide, 0.5, 7.0, 12.0, 0.4,
                "CUDA cores: ALU dominates because overhead is amortized over 32 threads. "
                "But at the SM level, the banked SRAM + collectors draw continuous power.",
                font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    return slide


# ============================================================================
# SLIDE 5: TensorCore Energy + CPU vs GPU Comparison
# ============================================================================

def build_slide_5_comparison(prs):
    """CPU vs GPU Energy: Different overhead, same fundamental problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.15, 12.5, 0.5,
                "CPU vs GPU: Different Overhead, Same Stored-Program Problem",
                font_size=26, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # -- TensorCore equation (more efficient path) --
    tc_y = 0.7
    add_textbox(slide, 0.3, tc_y, 12.5, 0.3,
                "TensorCore Path (4x4x4 MMA = 64 MACs per instruction):",
                font_size=14, bold=True, color=NVIDIA_GREEN)
    add_textbox(slide, 0.3, tc_y + 0.3, 12.5, 0.3,
                "E_tc_mac = E_overhead/64  +  E_tensor_alu  +  E_writeback/64  =  ~1.65 pJ/MAC",
                font_size=16, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER, font_name="Consolas")

    add_textbox(slide, 0.3, tc_y + 0.65, 12.5, 0.25,
                "TensorCores amortize overhead over 64 MACs (vs 32 for CUDA), reducing per-MAC overhead by 2x",
                font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # -- Side-by-side comparison table --
    table_y = 1.7
    add_textbox(slide, 0.5, table_y - 0.3, 12.0, 0.3,
                "Energy Breakdown Comparison at 4nm Process",
                font_size=16, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    rows, cols = 8, 4
    tbl = slide.shapes.add_table(rows, cols,
                                  Inches(1.5), Inches(table_y),
                                  Inches(10.0), Inches(2.6)).table

    headers = ["Energy Component", "CPU (scalar FP32)", "GPU CUDA Core", "GPU TensorCore"]
    tdata = [
        ["Instruction Frontend",   "9.5 pJ (63%)",    "0.025 pJ (2%)",  "0.013 pJ (1%)"],
        ["Operand Fetch / Delivery","1.6 pJ (11%)",   "0.082 pJ (7%)",  "0.041 pJ (4%)"],
        ["ALU / Compute",           "1.8 pJ (12%)",   "1.105 pJ (89%)", "0.34 pJ (29%)"],
        ["Writeback",               "0.9 pJ (6%)",    "0.025 pJ (2%)",  "0.010 pJ (1%)"],
        ["Pipeline / Wire / Other", "1.2 pJ (8%)",    "--- ",           "--- "],
        ["TOTAL per operation",     "~15.0 pJ",       "~1.24 pJ",       "~0.40 pJ"],
        ["Compute Efficiency",      "12%",            "89%",            "85%"],
    ]

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE

    for i, row_data in enumerate(tdata):
        is_total = ("TOTAL" in row_data[0] or "Ratio" in row_data[0])
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.bold = is_total
                p.alignment = PP_ALIGN.CENTER
                if is_total and "TOTAL" in row_data[0]:
                    p.font.color.rgb = YELLOW
                elif is_total:
                    p.font.color.rgb = ACCENT_TEAL
                elif "ALU" in row_data[0]:
                    p.font.color.rgb = GREEN
                elif "Frontend" in row_data[0]:
                    p.font.color.rgb = ORANGE
                else:
                    p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = (RGBColor(0x33, 0x33, 0x55) if is_total
                                        else RGBColor(0x2A, 0x2A, 0x4E) if i % 2 == 0
                                        else RGBColor(0x22, 0x22, 0x3E))

    # -- The "but" -- SM-level power --
    but_y = 4.5
    add_line(slide, 1.0, but_y, 12.0, but_y, color=ORANGE, width=2)
    add_textbox(slide, 0.5, but_y + 0.1, 12.0, 0.35,
                "But: The Full SM Picture -- Where the Energy Really Goes",
                font_size=18, bold=True, color=ORANGE,
                alignment=PP_ALIGN.CENTER)

    add_multiline_textbox(slide, 0.5, but_y + 0.5, 12.0, 1.8, [
        ("Per-MAC numbers look great for the GPU. But the SM has:", False, LIGHT_GRAY),
        ("", False, LIGHT_GRAY),
        ("  * 64K registers (256 KB of banked SRAM) -- powered every cycle regardless of utilization",
         False, WHITE),
        ("  * 4 operand collectors -- continuously buffering, arbitrating, routing",
         False, WHITE),
        ("  * 4 warp schedulers + scoreboards -- tracking 64 warps of hazard state",
         False, WHITE),
        ("  * 32 register banks with crossbar -- always ready for next warp switch",
         False, WHITE),
        ("", False, LIGHT_GRAY),
        ("This infrastructure exists to hide memory latency through thread concurrency.",
         True, YELLOW),
        ("At the SM level, the register file + operand delivery burns ~60% of compute power.",
         True, ORANGE),
    ], font_size=12)

    # -- Bottom: the question --
    add_rounded_rect(slide, 1.5, 7.0, 10.0, 0.45, RGBColor(0x2A, 0x2A, 0x4E))
    add_textbox(slide, 1.5, 7.03, 10.0, 0.35,
                "What if we could eliminate the register file entirely? (Next: Domain Flow Architecture)",
                font_size=15, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)

    return slide


# ============================================================================
# MAIN
# ============================================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1_title(prs)
    build_slide_2_sm_architecture(prs)
    build_slide_3_register_file(prs)
    build_slide_4_energy_equation(prs)
    build_slide_5_comparison(prs)

    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "..", "gpu_energy_architecture.pptx")
    output_path = os.path.normpath(output_path)
    prs.save(output_path)
    print(f"Saved presentation to: {output_path}")
    print(f"  5 slides generated:")
    print(f"  1. Title: GPU Energy -- From Instruction Tax to Data Movement Tax")
    print(f"  2. SM Architecture: 4 Partitions with Operand Delivery Infrastructure")
    print(f"  3. The 64K Register File: Banked SRAM with Operand Collectors")
    print(f"  4. GPU Energy Equation: Per-MAC Cost Breakdown")
    print(f"  5. CPU vs GPU Comparison: Different Overhead, Same Problem")


if __name__ == "__main__":
    main()
