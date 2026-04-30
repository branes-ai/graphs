#!/usr/bin/env python3
"""
Core Compute Engine Energy v2 -- The Register File is the Architecture.

This deck focuses on the central architectural difference: the GPU's shared
banked SRAM register file vs the TPU/KPU distributed PE-local registers.
All architectures at 7nm, core compute only.

Slide 1: Title -- The Register File IS the Architecture
Slide 2: CPU Core -- Scalar + SIMD with Private Register File
Slide 3: GPU SM -- The 256KB Banked SRAM and Its Infrastructure
Slide 4: TPU/KPU -- Distributed PE-Local Registers (No Shared RF)
Slide 5: Dynamic Energy -- Quantifying the Per-Op RF Cost
Slide 6: Static Power -- The Always-On Cost of 256KB SRAM
Slide 7: Summary -- Core Compute Energy at 7nm
"""

import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'src'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xD8)
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
MED_GRAY    = RGBColor(0x99, 0x99, 0x99)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)
RED_ACCENT  = RGBColor(0xE0, 0x40, 0x40)
GREEN       = RGBColor(0x2E, 0xCC, 0x71)
YELLOW      = RGBColor(0xF3, 0xC6, 0x23)

CPU_COLOR     = RGBColor(0x3A, 0x86, 0xC8)
GPU_COLOR     = RGBColor(0x76, 0xB9, 0x00)
TPU_COLOR     = RGBColor(0xDD, 0x88, 0x33)
KPU_COLOR     = RGBColor(0x00, 0x96, 0xD6)
ALU_COLOR     = RGBColor(0x2E, 0xCC, 0x71)
RF_COLOR      = RGBColor(0x6C, 0x5C, 0xE7)
SIMT_COLOR    = RGBColor(0x8B, 0x5C, 0xF6)
COLLECTOR_CLR = RGBColor(0xF3, 0x96, 0x23)
CROSSBAR_CLR  = RGBColor(0xCC, 0x55, 0x77)
BANK_COLOR    = RGBColor(0x48, 0x7D, 0xB3)
SURE_COLOR    = RGBColor(0x00, 0xCC, 0x88)
PE_COLOR      = RGBColor(0x3A, 0xB0, 0xD0)
BUF_COLOR     = RGBColor(0x55, 0x88, 0xBB)
STATIC_COLOR  = RGBColor(0xBB, 0x33, 0x33)

ROW_EVEN = RGBColor(0x2A, 0x2A, 0x4E)
ROW_ODD  = RGBColor(0x22, 0x22, 0x3E)

def set_slide_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def add_textbox(slide, l, t, w, h, text, fs=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, fn="Calibri"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(fs)
    p.font.bold = bold; p.font.color.rgb = color; p.font.name = fn; p.alignment = align
    return tb

def add_ml(slide, l, t, w, h, lines, fs=16, color=WHITE, align=PP_ALIGN.LEFT,
           fn="Calibri", ls=1.2):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(lines):
        txt, bld, clr = (item, False, color) if isinstance(item, str) else (
            item[0], item[1] if len(item)>1 else False, item[2] if len(item)>2 else color)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(fs); p.font.bold = bld
        p.font.color.rgb = clr; p.font.name = fn; p.alignment = align
        p.space_after = Pt(fs * (ls - 1.0) * 2)
    return tb

def add_rr(slide, l, t, w, h, fc, text="", fs=14, ftc=WHITE, bold=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fc; s.line.fill.background()
    if text:
        tf = s.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        p.text = text; p.font.size = Pt(fs); p.font.color.rgb = ftc
        p.font.bold = bold; p.alignment = PP_ALIGN.CENTER
    return s

def add_rect(slide, l, t, w, h, fc, text="", fs=14, ftc=WHITE, bold=False, lc=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if lc: s.line.color.rgb = lc; s.line.width = Pt(1.5)
    else: s.line.fill.background()
    if text:
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_left = Pt(4); tf.margin_right = Pt(4)
        tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(fs)
        p.font.color.rgb = ftc; p.font.bold = bold; p.alignment = PP_ALIGN.CENTER
    return s

def add_arrow_down(slide, l, t, w=0.2, h=0.3, c=MED_GRAY):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background(); return s

def add_arrow_right(slide, l, t, w=0.3, h=0.2, c=MED_GRAY):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background(); return s

def add_line(slide, x1, y1, x2, y2, c=MED_GRAY, w=1.5):
    cn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = c; cn.line.width = Pt(w); return cn


# ============================================================================
# SLIDE 1: Title
# ============================================================================
def build_slide_1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(s, DARK_BG)

    add_textbox(s, 0.5, 0.5, 12.0, 0.8,
                "The Register File IS the Architecture",
                fs=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 1.3, 12.0, 0.5,
                "Why 256 KB of Banked SRAM Defines GPU Energy -- and What TPU/KPU Do Instead",
                fs=20, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    add_line(s, 2.5, 2.0, 10.5, 2.0, c=RF_COLOR, w=2)

    add_rr(s, 1.5, 2.3, 10.0, 1.8, RGBColor(0x2A, 0x2A, 0x4E))
    add_ml(s, 1.8, 2.4, 9.4, 1.6, [
        ("The Central Question:", True, ORANGE),
        ("Every compute architecture must deliver operands to ALUs. The question is HOW:", False, LIGHT_GRAY),
        ("", False, WHITE),
        ("GPU: Centralized 256 KB banked SRAM shared by 2048 threads, accessed via", False, WHITE),
        ("     address decode -> bank arbitration -> operand collector -> crossbar -> ALU", False, LIGHT_GRAY),
        ("", False, WHITE),
        ("TPU/KPU: Distributed PE-local registers (latches). Each PE has 2-4 registers.", False, WHITE),
        ("     Data arrives via spatial wiring (compile-time routed). No shared SRAM.", False, LIGHT_GRAY),
    ], fs=13, align=PP_ALIGN.CENTER)

    # Three boxes showing the continuum
    bx_y = 4.5
    archs = [
        ("CPU", CPU_COLOR,
         "Small RF (16-64 regs)\nMulti-ported SRAM\nPer-core, private\n\nOverhead: moderate\n(addressed by SIMD width)"),
        ("GPU SM", GPU_COLOR,
         "MASSIVE RF (64K regs = 256 KB)\n32-banked SRAM\nShared by 2048 threads\n\nOverhead: LARGE\n(hidden by amortization)"),
        ("TPU / KPU", KPU_COLOR,
         "DISTRIBUTED (PE-local latches)\nNo shared SRAM structure\nNo address decode or arbitration\n\nOverhead: NEAR ZERO\n(operands arrive spatially)"),
    ]
    for i, (name, color, desc) in enumerate(archs):
        x = 0.5 + i * 4.2
        add_rr(s, x, bx_y, 3.8, 0.4, color, name, fs=15, bold=True)
        add_textbox(s, x + 0.1, bx_y + 0.5, 3.6, 1.8, desc, fs=11, color=LIGHT_GRAY)

    add_textbox(s, 0.5, 7.0, 12.0, 0.3,
                "All comparisons at 7nm TSMC N7. Core compute only (no external memory).",
                fs=11, color=MED_GRAY, align=PP_ALIGN.CENTER)
    return s


# ============================================================================
# SLIDE 2: CPU Core
# ============================================================================
def build_slide_2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(s, DARK_BG)

    add_textbox(s, 0.3, 0.1, 12.5, 0.45,
                "CPU: Small Register File, Amortized by SIMD Width",
                fs=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Left: diagram
    dx, dy = 0.5, 0.7
    add_rect(s, dx, dy, 5.5, 4.5, RGBColor(0x22, 0x22, 0x3E), "", lc=CPU_COLOR)
    add_textbox(s, dx, dy + 0.03, 5.5, 0.25, "CPU Core (1 of N)",
                fs=12, bold=True, color=CPU_COLOR, align=PP_ALIGN.CENTER)

    # Register file
    add_rect(s, dx + 0.3, dy + 0.4, 4.9, 1.0, RF_COLOR,
             "Register File: 16 architectural (x86-64) + ~100 physical (rename)\n"
             "Multi-ported SRAM: 2 read + 1 write per cycle\n"
             "Size: ~1-2 KB  |  Private to this core",
             fs=9, bold=True, lc=RF_COLOR)

    # Scalar ALU
    add_rr(s, dx + 0.3, dy + 1.6, 2.2, 0.7, ALU_COLOR,
           "Scalar ALU\nFP32 FMA\n1 op/cycle", fs=10, bold=True)

    # SIMD ALU
    add_rr(s, dx + 2.7, dy + 1.6, 2.5, 0.7, ALU_COLOR,
           "SIMD ALU (AVX-512)\n16 x FP32 FMA\n16 ops/cycle", fs=10, bold=True)

    # Instruction pipeline
    add_rr(s, dx + 0.3, dy + 2.5, 4.9, 0.5, RGBColor(0x3A, 0x55, 0x88),
           "Instruction Pipeline: I-cache -> Decode -> Dispatch -> Rename -> Issue",
           fs=9, bold=True)

    # L1 caches
    add_rr(s, dx + 0.3, dy + 3.2, 2.2, 0.45, BANK_COLOR,
           "L1 I-Cache (32 KB)", fs=9, bold=True)
    add_rr(s, dx + 2.7, dy + 3.2, 2.5, 0.45, BANK_COLOR,
           "L1 D-Cache (48 KB)", fs=9, bold=True)

    # Exclusive L2
    add_rr(s, dx + 0.3, dy + 3.8, 4.9, 0.4, RGBColor(0x33, 0x55, 0x77),
           "L2 Cache (256 KB - 1 MB, exclusive to core)", fs=9, bold=True)

    # Right: energy equations
    ex, ey = 6.5, 0.7
    add_textbox(s, ex, ey, 6.0, 0.3,
                "Energy Equations at 7nm:", fs=14, bold=True, color=ACCENT_TEAL)

    add_textbox(s, ex, ey + 0.35, 6.0, 0.25,
                "Scalar (1 op/instr):", fs=12, bold=True, color=CPU_COLOR)
    add_ml(s, ex, ey + 0.6, 6.0, 1.5, [
        ("E = E_ifetch + E_decode + 2*E_rf_read + E_alu + E_rf_write", True, WHITE),
        ("  = 0.540 + 0.288 + 2*1.013 + 1.800 + 1.215", False, MED_GRAY),
        ("  = 5.868 pJ (overhead) + 1.800 pJ (ALU) = 7.668 pJ/op", True, ORANGE),
        ("  x2 instructions per FMA = 9.936 pJ/op  |  Eff: 18%", False, LIGHT_GRAY),
    ], fs=10, fn="Consolas")

    add_textbox(s, ex, ey + 1.8, 6.0, 0.25,
                "AVX-512 (16 ops/instr):", fs=12, bold=True, color=CPU_COLOR)
    add_ml(s, ex, ey + 2.05, 6.0, 1.5, [
        ("E = (E_ifetch + E_decode + 2*E_vrf_read + E_vrf_write)/16 + E_alu", True, WHITE),
        ("  = (0.540 + 0.288 + 2*1.013 + 1.215) / 16  +  1.800", False, MED_GRAY),
        ("  = 0.254 pJ (overhead) + 1.800 pJ (ALU) = 2.054 pJ/op", True, GREEN),
        ("  x2 instructions per FMA / 16 = 2.309 pJ/op  |  Eff: 78%", False, LIGHT_GRAY),
    ], fs=10, fn="Consolas")

    # Key point
    ky = 4.2
    add_line(s, ex, ky, ex + 6.0, ky, c=CPU_COLOR, w=1.5)
    add_ml(s, ex, ky + 0.1, 6.0, 1.5, [
        ("CPU RF is SMALL (~1-2 KB) and PRIVATE:", True, CPU_COLOR),
        ("  * Multi-ported but only 2R+1W per cycle", False, LIGHT_GRAY),
        ("  * No bank conflicts (only 1 thread using it)", False, LIGHT_GRAY),
        ("  * No operand collector (data immediately available)", False, LIGHT_GRAY),
        ("  * SIMD amortizes the overhead over 16 lanes", False, LIGHT_GRAY),
        ("", False, WHITE),
        ("The CPU RF is a solved problem at ~1-2 KB.", True, WHITE),
        ("The GPU scales this to 256 KB -- that changes everything.", True, ORANGE),
    ], fs=11)

    add_textbox(s, 0.3, 6.7, 12.5, 0.25,
                "Granularity: 1 core  |  RF: 1-2 KB private  |  1 scalar + 1 SIMD ALU  |  Private L1i + L1d + L2",
                fs=11, color=MED_GRAY, align=PP_ALIGN.CENTER)
    return s


# ============================================================================
# SLIDE 3: GPU SM -- The 256KB Banked SRAM
# ============================================================================
def build_slide_3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(s, DARK_BG)

    add_textbox(s, 0.3, 0.1, 12.5, 0.45,
                "GPU SM: 256 KB Banked SRAM Register File + Delivery Infrastructure",
                fs=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.3, 0.5, 12.5, 0.3,
                "The register file IS the SM. Everything else exists to feed it or be fed by it.",
                fs=13, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    # -- LEFT: The RF structure --
    rx, ry = 0.3, 0.9

    # Big RF box
    add_rect(s, rx, ry, 5.8, 3.3, RGBColor(0x1E, 0x1E, 0x38), "", lc=RF_COLOR)
    add_textbox(s, rx, ry + 0.03, 5.8, 0.3,
                "Register File: 64K x 32-bit = 256 KB Banked SRAM",
                fs=12, bold=True, color=RF_COLOR, align=PP_ALIGN.CENTER)

    # 32 banks (show 8x4 grid)
    for row in range(4):
        for col in range(8):
            bx = rx + 0.15 + col * 0.7
            by = ry + 0.4 + row * 0.4
            add_rect(s, bx, by, 0.62, 0.32, BANK_COLOR,
                     f"Bank {row*8+col}\n2K regs", fs=6, ftc=WHITE)

    add_textbox(s, rx, ry + 2.05, 5.8, 0.25,
                "32 banks  |  2048 regs/bank  |  8 KB SRAM per bank  |  All banks active every cycle",
                fs=9, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

    # Infrastructure below RF
    infra_y = ry + 2.35

    add_rr(s, rx + 0.1, infra_y, 2.6, 0.35, COLLECTOR_CLR,
           "Operand Collector (x4)", fs=9, ftc=RGBColor(0x1A, 0x1A, 0x2E), bold=True)
    add_rr(s, rx + 2.8, infra_y, 2.8, 0.35, CROSSBAR_CLR,
           "32-wide Crossbar + Bank Arb", fs=9, bold=True)

    add_arrow_down(s, rx + 2.8, infra_y + 0.35, 0.2, 0.15, MED_GRAY)

    add_rr(s, rx + 0.1, infra_y + 0.55, 5.5, 0.3, ALU_COLOR,
           "128 CUDA Cores + 4 TensorCores (shared across 4 partitions)",
           fs=9, ftc=RGBColor(0x1A, 0x1A, 0x2E), bold=True)

    # -- RIGHT: Why it must be 256 KB --
    ex, ey = 6.5, 0.9
    add_textbox(s, ex, ey, 6.0, 0.3,
                "Why 256 KB? The Latency-Hiding Requirement:", fs=14, bold=True, color=ORANGE)

    add_ml(s, ex, ey + 0.35, 6.0, 2.5, [
        ("Thread concurrency hides memory latency:", True, WHITE),
        ("  HBM access: ~400 cycles at 1.5 GHz", False, LIGHT_GRAY),
        ("  To hide this: need ~32 warps ready to execute", False, LIGHT_GRAY),
        ("  Each warp: 32 threads x 32 registers x 4 bytes = 4 KB", False, LIGHT_GRAY),
        ("  Total: 32 warps x 4 KB = 128 KB minimum", False, LIGHT_GRAY),
        ("  Actual: 64 warps x 4 KB = 256 KB (headroom)", True, YELLOW),
        ("", False, WHITE),
        ("This is NOT a design choice -- it is a REQUIREMENT.", True, ORANGE),
        ("Fewer registers = fewer resident threads = exposed latency = stalls.", False, LIGHT_GRAY),
        ("The GPU MUST have 256 KB of fast, banked SRAM to function.", False, LIGHT_GRAY),
    ], fs=11)

    # Infrastructure costs
    add_line(s, ex, ey + 2.8, ex + 6.0, ey + 2.8, c=RF_COLOR, w=1)
    add_textbox(s, ex, ey + 2.9, 6.0, 0.25,
                "Infrastructure Required to Serve 256 KB RF:", fs=13, bold=True, color=RF_COLOR)

    infra = [
        ("32 SRAM banks", "Parallel access for 32 threads per warp", BANK_COLOR),
        ("Bank arbitration logic", "Resolve conflicts when 2+ warps access same bank", BANK_COLOR),
        ("4 Operand collectors", "Buffer partially-read operands until all 32 are ready", COLLECTOR_CLR),
        ("32-wide crossbar", "Route operands from any bank to any ALU lane", CROSSBAR_CLR),
        ("4 Warp schedulers", "Select which of 16 warps to issue each cycle", SIMT_COLOR),
        ("4 Scoreboards", "Track RAW/WAW/WAR hazards across all active warps", RED_ACCENT),
    ]

    for i, (name, desc, color) in enumerate(infra):
        y = ey + 3.2 + i * 0.28
        add_textbox(s, ex + 0.1, y, 2.2, 0.25, name, fs=9, bold=True, color=color)
        add_textbox(s, ex + 2.4, y, 3.5, 0.25, desc, fs=8, color=LIGHT_GRAY)

    # Bottom
    add_rr(s, 0.5, 6.6, 12.0, 0.7, RGBColor(0x2A, 0x2A, 0x4E))
    add_ml(s, 0.8, 6.65, 11.4, 0.6, [
        ("The GPU cannot shrink or remove the register file. It is the latency-hiding mechanism.", True, ORANGE),
        ("Every component above (banks, arbitration, collectors, crossbar, schedulers, scoreboards) exists", False, LIGHT_GRAY),
        ("solely to serve this 256 KB SRAM. This is the architecture the TPU and KPU eliminate.", True, WHITE),
    ], fs=12, align=PP_ALIGN.CENTER)
    return s


# ============================================================================
# SLIDE 4: TPU/KPU -- Distributed PE-Local Registers
# ============================================================================
def build_slide_4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(s, DARK_BG)

    add_textbox(s, 0.3, 0.1, 12.5, 0.45,
                "TPU / KPU: Distributed PE-Local Registers Replace Shared SRAM",
                fs=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.3, 0.5, 12.5, 0.3,
                "The spatial schedule eliminates resource conflicts -- each PE owns its data",
                fs=13, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    # -- LEFT: TPU --
    tx, ty = 0.3, 0.9
    add_rect(s, tx, ty, 6.0, 3.5, RGBColor(0x22, 0x22, 0x3E), "", lc=TPU_COLOR)
    add_textbox(s, tx, ty + 0.03, 6.0, 0.25,
                "TPU: 128x128 Systolic Array (Weight-Stationary)",
                fs=12, bold=True, color=TPU_COLOR, align=PP_ALIGN.CENTER)

    # PE detail
    add_rect(s, tx + 0.2, ty + 0.4, 5.6, 1.4, RGBColor(0x1E, 0x2E, 0x3E), "", lc=TPU_COLOR)

    # Show a row of PEs with internal registers
    for c in range(6):
        px = tx + 0.4 + c * 0.85
        py = ty + 0.55
        add_rect(s, px, py, 0.7, 1.1, RGBColor(0x3A, 0x5A, 0x3A), "", lc=TPU_COLOR)
        add_textbox(s, px, py + 0.02, 0.7, 0.15, "PE", fs=7, bold=True,
                    color=TPU_COLOR, align=PP_ALIGN.CENTER)
        # Weight register
        add_rect(s, px + 0.05, py + 0.2, 0.6, 0.2,
                 RGBColor(0xAA, 0x66, 0x22), "W reg", fs=6, ftc=WHITE, bold=True)
        # MAC
        add_rect(s, px + 0.05, py + 0.45, 0.6, 0.2,
                 ALU_COLOR, "MAC", fs=7, ftc=WHITE, bold=True)
        # Accumulator
        add_rect(s, px + 0.05, py + 0.7, 0.6, 0.2,
                 RGBColor(0x88, 0x44, 0x44), "Accum", fs=6, ftc=WHITE, bold=True)
        # Horizontal wire
        if c < 5:
            add_arrow_right(s, px + 0.7, py + 0.53, 0.15, 0.1, TPU_COLOR)

    add_textbox(s, tx + 0.2, ty + 1.9, 5.6, 0.25,
                "Inputs stream left-to-right  |  Partial sums flow top-to-bottom",
                fs=9, color=MED_GRAY, align=PP_ALIGN.CENTER)

    # TPU register summary
    add_ml(s, tx + 0.2, ty + 2.2, 5.6, 1.2, [
        ("Total storage: 16,384 PEs x 2 regs x 32b = 128 KB", True, TPU_COLOR),
        ("But: each register is a LOCAL latch inside its PE", False, LIGHT_GRAY),
        ("  * No address decode (hardwired)", False, GREEN),
        ("  * No bank arbitration (no sharing)", False, GREEN),
        ("  * No operand collector (data already local)", False, GREEN),
        ("  * No crossbar (PE-to-PE wire, compile-time routed)", False, GREEN),
        ("  * Weight-stationary: W loaded once, reused 128+ times", False, LIGHT_GRAY),
    ], fs=10)

    # -- RIGHT: KPU --
    kx, ky = 6.7, 0.9
    add_rect(s, kx, ky, 6.0, 3.5, RGBColor(0x22, 0x22, 0x3E), "", lc=KPU_COLOR)
    add_textbox(s, kx, ky + 0.03, 6.0, 0.25,
                "KPU: 16x16 PE Array (Programmable Domain Flow)",
                fs=12, bold=True, color=KPU_COLOR, align=PP_ALIGN.CENTER)

    # PE detail
    add_rect(s, kx + 0.2, ky + 0.4, 5.6, 1.4, RGBColor(0x1E, 0x2E, 0x3E), "", lc=KPU_COLOR)

    for c in range(6):
        px = kx + 0.4 + c * 0.85
        py = ky + 0.55
        add_rect(s, px, py, 0.7, 1.1, RGBColor(0x2A, 0x4A, 0x5A), "", lc=KPU_COLOR)
        add_textbox(s, px, py + 0.02, 0.7, 0.15, "PE", fs=7, bold=True,
                    color=KPU_COLOR, align=PP_ALIGN.CENTER)
        # A register
        add_rect(s, px + 0.05, py + 0.2, 0.28, 0.2,
                 RGBColor(0x44, 0x88, 0x66), "A", fs=6, ftc=WHITE, bold=True)
        # B register
        add_rect(s, px + 0.37, py + 0.2, 0.28, 0.2,
                 RGBColor(0x44, 0x66, 0x88), "B", fs=6, ftc=WHITE, bold=True)
        # MAC
        add_rect(s, px + 0.05, py + 0.45, 0.6, 0.2,
                 ALU_COLOR, "MAC", fs=7, ftc=WHITE, bold=True)
        # C register
        add_rect(s, px + 0.15, py + 0.7, 0.4, 0.2,
                 RGBColor(0x66, 0x44, 0x88), "C", fs=6, ftc=WHITE, bold=True)
        # SURE wires
        if c < 5:
            add_line(s, px + 0.7, py + 0.3, px + 0.85, py + 0.3, c=SURE_COLOR, w=1)
            add_line(s, px + 0.7, py + 0.8, px + 0.85, py + 0.8, c=SURE_COLOR, w=1)

    add_textbox(s, kx + 0.2, ky + 1.9, 5.6, 0.25,
                "SURE network: compile-time routed wires between PEs (~0.02 pJ/transfer)",
                fs=9, color=SURE_COLOR, bold=True, align=PP_ALIGN.CENTER)

    # KPU register summary
    add_ml(s, kx + 0.2, ky + 2.2, 5.6, 1.2, [
        ("Total storage: 256 PEs x 4 regs x 32b = 4 KB per tile", True, KPU_COLOR),
        ("64 tiles = 16,384 PEs x 4 regs = 256 KB total (same as GPU!)", True, YELLOW),
        ("But: each register is a LOCAL latch inside its PE", False, LIGHT_GRAY),
        ("  * Same benefits as TPU (no decode, no arbitration, ...)", False, GREEN),
        ("  * PLUS: programmable schedule (A/B/C-stationary)", False, GREEN),
        ("  * PLUS: any SURE recurrence equation, not just GEMM", False, GREEN),
        ("  * Smaller tile (256 MACs) = better utilization at batch=1", False, GREEN),
    ], fs=10)

    # -- Bottom: the comparison --
    cmp_y = 4.65
    add_line(s, 0.5, cmp_y, 12.5, cmp_y, c=ACCENT_BLUE, w=2)
    add_textbox(s, 0.3, cmp_y + 0.1, 12.5, 0.3,
                "The Fundamental Difference: Shared SRAM vs Distributed Latches",
                fs=18, bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    headers = ["", "GPU SM", "TPU Array", "KPU Tile"]
    data = [
        ["Total register storage", "256 KB", "128 KB", "256 KB (64 tiles)"],
        ["Storage structure", "Centralized banked SRAM", "PE-local latches", "PE-local latches"],
        ["Access mechanism", "Address decode + bank arb\n+ collector + crossbar", "Spatial wiring\n(weight-stationary)", "SURE network\n(compile-time routed)"],
        ["Concurrent accessors", "2048 threads (contention!)", "16,384 PEs (no contention)", "16,384 PEs (no contention)"],
        ["Infrastructure needed", "Scheduler, scoreboard,\ncollector, crossbar", "Input/output buffers only", "Streamer + L1 buffer only"],
        ["Can GPU copy this?", "---", "NO (would require\nremoving SIMT model)", "NO (would require\nremoving SIMT model)"],
    ]

    rows = len(data) + 1
    cols = len(headers)
    tbl = s.shapes.add_table(rows, cols, Inches(0.5), Inches(cmp_y + 0.45),
                              Inches(12.0), Inches(2.2)).table
    hdr_colors = [RGBColor(0x33, 0x33, 0x55), GPU_COLOR, TPU_COLOR, KPU_COLOR]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j); cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = hdr_colors[j]
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = tbl.cell(i+1, j); cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8); p.font.color.rgb = WHITE
                p.font.bold = (j == 0 or "NO" in val)
                p.alignment = PP_ALIGN.CENTER
                if "NO" in val: p.font.color.rgb = RED_ACCENT
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_EVEN if i % 2 == 0 else ROW_ODD

    return s


# ============================================================================
# SLIDE 5: Dynamic Energy -- Per-Op RF Cost
# ============================================================================
def build_slide_5(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(s, DARK_BG)

    add_textbox(s, 0.3, 0.1, 12.5, 0.45,
                "Dynamic Energy: The Per-Operation Cost of Operand Delivery",
                fs=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.3, 0.5, 12.5, 0.3,
                "Same computation, different operand delivery -- all at 7nm, per FP32 MAC",
                fs=13, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    # GPU operand delivery chain
    gy = 0.9
    add_textbox(s, 0.3, gy, 12.5, 0.25,
                "GPU: Operand Delivery Chain (per warp instruction, CUDA path)",
                fs=14, bold=True, color=GPU_COLOR)

    gpu_chain = [
        ("Warp\nScheduler", SIMT_COLOR, "0.875"),
        ("Score-\nboard", RED_ACCENT, "0.525"),
        ("Addr\nGen x3", BANK_COLOR, "1.050"),
        ("Bank\nArbitr.", BANK_COLOR, "0.525"),
        ("SRAM\nRead x2", RF_COLOR, "0.562"),
        ("Operand\nCollector", COLLECTOR_CLR, "1.400"),
        ("Cross-\nbar", CROSSBAR_CLR, "0.700"),
        ("ALU\n(32 MACs)", ALU_COLOR, "1.530\n/op"),
        ("Result\nRoute+WB", RED_ACCENT, "0.863"),
    ]

    bw = 1.15
    bh = 0.65
    for i, (name, color, energy) in enumerate(gpu_chain):
        x = 0.3 + i * (bw + 0.12)
        add_rr(s, x, gy + 0.3, bw, bh, color, name, fs=7, bold=True,
               ftc=WHITE if color != COLLECTOR_CLR else RGBColor(0x1A, 0x1A, 0x2E))
        e_color = GREEN if "ALU" in name else ORANGE
        add_textbox(s, x, gy + bh + 0.32, bw, 0.2,
                    f"{energy} pJ", fs=8, bold=True, color=e_color, align=PP_ALIGN.CENTER)
        if i < len(gpu_chain) - 1:
            add_arrow_right(s, x + bw + 0.01, gy + 0.55, 0.1, 0.12, MED_GRAY)

    # Sum
    overhead_sum = 0.875 + 0.525 + 1.050 + 0.525 + 0.562 + 1.400 + 0.700 + 0.863
    add_textbox(s, 0.3, gy + 1.2, 12.5, 0.25,
                f"Overhead per warp: {overhead_sum:.3f} pJ  |  Per MAC (/32): {overhead_sum/32:.3f} pJ  "
                f"|  ALU: 1.530 pJ  |  Total: {overhead_sum/32 + 1.530:.3f} pJ/op  |  Eff: {1.530/(overhead_sum/32+1.530)*100:.0f}%",
                fs=11, bold=True, color=GPU_COLOR, fn="Consolas", align=PP_ALIGN.CENTER)

    # TPU chain
    ty2 = gy + 1.6
    add_textbox(s, 0.3, ty2, 12.5, 0.25,
                "TPU: Operand Delivery (spatial, weight-stationary)", fs=14, bold=True, color=TPU_COLOR)

    tpu_chain = [
        ("Input\nBuffer", BUF_COLOR, "0.300"),
        ("PE wire\n(spatial)", SURE_COLOR, "~0.05"),
        ("Systolic\nMAC", ALU_COLOR, "1.440\n/op"),
        ("PE wire\n(drain)", SURE_COLOR, "~0.01"),
    ]

    for i, (name, color, energy) in enumerate(tpu_chain):
        x = 0.3 + i * (bw + 0.6)
        add_rr(s, x, ty2 + 0.3, bw, bh * 0.85, color, name, fs=9, bold=True,
               ftc=WHITE if color != SURE_COLOR else RGBColor(0x1A, 0x1A, 0x2E))
        e_color = GREEN if "MAC" in name else LIGHT_GRAY
        add_textbox(s, x, ty2 + bh * 0.85 + 0.32, bw, 0.2,
                    f"{energy} pJ", fs=9, bold=True, color=e_color, align=PP_ALIGN.CENTER)
        if i < len(tpu_chain) - 1:
            add_arrow_right(s, x + bw + 0.1, ty2 + 0.55, 0.4, 0.12, TPU_COLOR)

    add_textbox(s, 5.5, ty2 + 0.3, 7.0, 0.6,
                "Overhead: ~0.36 pJ/op  |  Total: ~1.80 pJ/op  |  Eff: 80%\n"
                "No scheduler, no scoreboard, no collector, no crossbar.",
                fs=11, bold=True, color=TPU_COLOR)

    # KPU chain
    ky2 = ty2 + 1.3
    add_textbox(s, 0.3, ky2, 12.5, 0.25,
                "KPU: Operand Delivery (SURE network, compile-time routed)", fs=14, bold=True, color=KPU_COLOR)

    kpu_chain = [
        ("L1 Stream\nBuffer", BUF_COLOR, "~0.01"),
        ("PE wire\n(SURE)", SURE_COLOR, "0.022"),
        ("Domain\nFlow MAC", ALU_COLOR, "0.675-\n2.160"),
        ("Local\nreg write", PE_COLOR, "0.009"),
    ]

    for i, (name, color, energy) in enumerate(kpu_chain):
        x = 0.3 + i * (bw + 0.6)
        add_rr(s, x, ky2 + 0.3, bw, bh * 0.85, color, name, fs=9, bold=True,
               ftc=WHITE if color not in (SURE_COLOR, PE_COLOR) else RGBColor(0x1A, 0x1A, 0x2E))
        e_color = GREEN if "MAC" in name else LIGHT_GRAY
        add_textbox(s, x, ky2 + bh * 0.85 + 0.32, bw, 0.2,
                    f"{energy} pJ", fs=9, bold=True, color=e_color, align=PP_ALIGN.CENTER)
        if i < len(kpu_chain) - 1:
            add_arrow_right(s, x + bw + 0.1, ky2 + 0.55, 0.4, 0.12, KPU_COLOR)

    add_textbox(s, 5.5, ky2 + 0.3, 7.0, 0.6,
                "Overhead: ~0.053 pJ/op  |  Total: 0.958 pJ/op (mixed)  |  Eff: 94%\n"
                "Same zero-infrastructure as TPU, plus programmable schedule.",
                fs=11, bold=True, color=KPU_COLOR)

    # Comparison summary
    cmp_y = ky2 + 1.35
    add_line(s, 0.5, cmp_y, 12.5, cmp_y, c=ACCENT_BLUE, w=2)
    add_rr(s, 0.5, cmp_y + 0.1, 12.0, 0.7, RGBColor(0x2A, 0x2A, 0x4E))
    add_ml(s, 0.8, cmp_y + 0.13, 11.4, 0.65, [
        ("Operand delivery overhead per MAC at 7nm:", True, ACCENT_TEAL),
        (f"  GPU: {overhead_sum/32:.3f} pJ (9 pipeline stages, all serving 256 KB SRAM)"
         f"  |  TPU: ~0.36 pJ (buffer + wire)  |  KPU: ~0.053 pJ (wire + local reg)", False, WHITE),
        (f"  GPU-to-KPU overhead ratio: {overhead_sum/32/0.053:.0f}x.  "
         "This is the structural advantage the GPU cannot replicate without removing SIMT.",
         True, YELLOW),
    ], fs=12, align=PP_ALIGN.CENTER)

    return s


# ============================================================================
# SLIDE 6: Static Power
# ============================================================================
def build_slide_6(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(s, DARK_BG)

    add_textbox(s, 0.3, 0.1, 12.5, 0.45,
                "Static Power: The Always-On Cost of 256 KB SRAM",
                fs=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.3, 0.5, 12.5, 0.3,
                "A separate argument from dynamic energy: leakage power is paid regardless of utilization",
                fs=13, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    # Static power explanation
    add_ml(s, 0.5, 0.95, 6.0, 3.5, [
        ("What is Static Power?", True, STATIC_COLOR),
        ("", False, WHITE),
        ("Every transistor in the 256 KB SRAM leaks current", False, LIGHT_GRAY),
        ("even when no computation is happening.", False, LIGHT_GRAY),
        ("This is a consequence of FinFET physics at 7nm:", False, LIGHT_GRAY),
        ("", False, WHITE),
        ("  P_static = N_transistors * I_leak * V_dd", True, WHITE),
        ("", False, WHITE),
        ("At 7nm FinFET, typical SRAM leakage:", False, LIGHT_GRAY),
        ("  ~0.2 W per MB of SRAM (typical 7nm datacenter)", False, LIGHT_GRAY),
        ("  256 KB SRAM -> ~50 mW per SM", True, YELLOW),
        ("", False, WHITE),
        ("This leakage is paid:", True, ORANGE),
        ("  * When computing (overlaps with dynamic energy)", False, LIGHT_GRAY),
        ("  * When waiting for memory (100s of cycles)", False, LIGHT_GRAY),
        ("  * When partially utilized (batch=1 workloads)", False, LIGHT_GRAY),
        ("  * When idle (unless power-gated -- but GPU SMs", False, LIGHT_GRAY),
        ("    rarely power-gate RF due to thread context)", False, LIGHT_GRAY),
    ], fs=11)

    # Scale to full GPU
    ex = 7.0
    add_textbox(s, ex, 0.95, 5.5, 0.3,
                "Scaling to Full GPU (H100-class, 7nm):", fs=14, bold=True, color=GPU_COLOR)

    add_ml(s, ex, 1.3, 5.5, 2.5, [
        ("Per SM:", True, WHITE),
        ("  RF SRAM: 256 KB", False, LIGHT_GRAY),
        ("  Leakage: ~50 mW", False, LIGHT_GRAY),
        ("  At 1.5 GHz: 33.3 pJ/cycle", False, LIGHT_GRAY),
        ("  At 128 MACs/cycle: 0.26 pJ/MAC (CUDA)", False, LIGHT_GRAY),
        ("  At 256 MACs/cycle: 0.13 pJ/MAC (TC)", False, YELLOW),
        ("", False, WHITE),
        ("Full GPU (132 SMs):", True, WHITE),
        ("  Total RF: 132 x 256 KB = 33 MB SRAM", False, LIGHT_GRAY),
        ("  Total RF leakage: 132 x 50 mW = 6.6 W", True, YELLOW),
        ("  Share of 700W TDP: 0.9%", False, LIGHT_GRAY),
        ("", False, WHITE),
        ("  + Shared L1/SMEM (192 KB x 132 = 25 MB): ~5 W", False, MED_GRAY),
        ("  + L2 Cache (50 MB): ~10 W", False, MED_GRAY),
        ("  Total SRAM leakage: ~22 W (3% of TDP)", False, MED_GRAY),
    ], fs=10)

    # Comparison
    cmp_y = 4.2
    add_line(s, 0.5, cmp_y, 12.5, cmp_y, c=ACCENT_BLUE, w=1.5)
    add_textbox(s, 0.3, cmp_y + 0.1, 12.5, 0.3,
                "TPU/KPU Comparison: Distributed Latches Have ~10x Lower Leakage",
                fs=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    add_ml(s, 0.5, cmp_y + 0.45, 12.0, 2.0, [
        ("PE-local latches (flip-flops) leak ~10x less than SRAM cells per bit, because:", True, WHITE),
        ("  * SRAM cell: 6 transistors, continuous feedback loop, both P/N paths leak", False, LIGHT_GRAY),
        ("  * Flip-flop latch: simpler structure, clock-gated when PE is idle", False, LIGHT_GRAY),
        ("", False, WHITE),
        ("Furthermore, spatial architectures enable fine-grained power gating:", True, WHITE),
        ("  * GPU: cannot power-gate individual SM register banks (thread context would be lost)", False, LIGHT_GRAY),
        ("  * TPU: can power-gate unused rows/columns of systolic array (but array is large)", False, LIGHT_GRAY),
        ("  * KPU: can power-gate individual tiles (16x16 = fine granularity, 256 MACs per tile)", False, GREEN),
    ], fs=11)

    # But caveat
    add_rr(s, 1.0, 6.6, 11.0, 0.7, RGBColor(0x2A, 0x2A, 0x4E))
    add_ml(s, 1.3, 6.65, 10.4, 0.6, [
        ("Important: Static power is a secondary effect (~0.13-0.26 pJ/MAC) compared to", False, LIGHT_GRAY),
        ("dynamic RF infrastructure cost (~0.20 pJ/MAC overhead). The dynamic cost is the primary argument.", True, ORANGE),
        ("Static power matters most at low utilization (batch=1, idle SMs, memory-bound workloads).", False, MED_GRAY),
    ], fs=11, align=PP_ALIGN.CENTER)

    return s


# ============================================================================
# SLIDE 7: Summary
# ============================================================================
def build_slide_7(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(s, DARK_BG)

    add_textbox(s, 0.3, 0.1, 12.5, 0.45,
                "Core Compute Energy Summary: All at 7nm, No External Memory",
                fs=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Summary table
    headers = ["", "CPU Scalar", "CPU AVX-512", "GPU CUDA", "GPU TC", "TPU 128x128", "KPU 16x16"]
    data = [
        ["RF / operand storage",    "~1 KB\nprivate SRAM", "~1 KB\nprivate SRAM",
         "256 KB\nbanked SRAM", "256 KB\nbanked SRAM", "128 KB\nPE latches", "4 KB/tile\nPE latches"],
        ["Operand delivery",        "2R+1W\nports", "2R+1W\nvector", "32-bank\n+collector\n+crossbar",
         "32-bank\n+collector\n+crossbar", "Spatial\nwire", "SURE\nwire"],
        ["MAC energy (pJ)",         "1.800", "1.800", "1.530", "1.530", "1.440", "0.905*"],
        ["Overhead (pJ/op)",        "8.136", "0.509", "0.418", "0.284", "0.639", "0.053"],
        ["  of which: RF infra",    "4.050", "0.253", "0.127", "0.075", "~0.01", "~0.009"],
        ["Static RF (pJ/MAC)",      "~0",    "~0",    "0.260", "0.130", "~0.03", "~0.01"],
        ["TOTAL (pJ/op)",           "9.936", "2.309", "1.948", "1.814", "2.085", "0.958"],
        ["Compute Efficiency",      "18%",   "78%",   "79%",   "84%",   "69%",   "94%"],
    ]

    rows = len(data) + 1; cols = len(headers)
    tbl = s.shapes.add_table(rows, cols, Inches(0.3), Inches(0.6),
                              Inches(12.6), Inches(3.5)).table

    hc = [RGBColor(0x33, 0x33, 0x55), CPU_COLOR, CPU_COLOR,
          GPU_COLOR, GPU_COLOR, TPU_COLOR, KPU_COLOR]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j); cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = hc[j]

    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = tbl.cell(i+1, j); cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8); p.font.bold = (j == 0 or i in [2,3,6,7])
                p.alignment = PP_ALIGN.CENTER
                if i == 6: p.font.color.rgb = YELLOW
                elif i == 7: p.font.color.rgb = GREEN
                elif i == 4: p.font.color.rgb = ORANGE  # RF infra row
                elif i == 5: p.font.color.rgb = STATIC_COLOR  # static row
                elif j == 0: p.font.color.rgb = ACCENT_TEAL
                else: p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_EVEN if i % 2 == 0 else ROW_ODD

    # Energy bars
    bar_y = 4.3
    add_textbox(s, 0.5, bar_y, 12.0, 0.25,
                "Energy per Op:  Green = MAC  |  Orange = RF Infrastructure  |  Gray = Other Overhead  |  Red = Static",
                fs=11, bold=True, color=WHITE)

    bar_data = [
        ("CPU Scalar",  1.800, 4.050, 4.086, 0.0,   CPU_COLOR),
        ("CPU AVX-512", 1.800, 0.253, 0.256, 0.0,   CPU_COLOR),
        ("GPU CUDA",    1.530, 0.127, 0.291, 0.260, GPU_COLOR),
        ("GPU TC",      1.530, 0.075, 0.209, 0.130, GPU_COLOR),
        ("TPU",         1.440, 0.010, 0.629, 0.030, TPU_COLOR),
        ("KPU",         0.905, 0.009, 0.044, 0.010, KPU_COLOR),
    ]

    max_e = 10.0
    bw = 8.0; bs = 3.0; bh = 0.3

    for i, (name, mac, rf, other, static, color) in enumerate(bar_data):
        by = bar_y + 0.3 + i * (bh + 0.06)
        add_textbox(s, 0.5, by, 2.4, bh, name, fs=10, bold=True, color=color, align=PP_ALIGN.RIGHT)

        cx = bs
        cw = bw * (mac / max_e)
        add_rect(s, cx, by, cw, bh, ALU_COLOR, f"{mac:.2f}" if cw > 0.5 else "", fs=7, bold=True)
        cx += cw

        rw = bw * (rf / max_e)
        if rw > 0.08:
            add_rect(s, cx, by, rw, bh, ORANGE,
                     f"{rf:.2f}" if rw > 0.3 else "", fs=7, bold=True)
        cx += rw

        ow = bw * (other / max_e)
        if ow > 0.08:
            add_rect(s, cx, by, ow, bh, RGBColor(0x55, 0x55, 0x77),
                     f"{other:.2f}" if ow > 0.25 else "", fs=7, ftc=MED_GRAY, bold=True)
        cx += ow

        sw = bw * (static / max_e)
        if sw > 0.05:
            add_rect(s, cx, by, sw, bh, STATIC_COLOR,
                     f"{static:.2f}" if sw > 0.15 else "", fs=7, bold=True)
        cx += sw

        total = mac + rf + other + static
        add_textbox(s, cx + 0.05, by, 1.2, bh,
                    f"{total:.2f}", fs=9, bold=True, color=LIGHT_GRAY)

    # Takeaway
    tk_y = bar_y + 0.3 + 6 * (bh + 0.06) + 0.15
    add_rr(s, 0.5, tk_y, 12.0, 0.85, RGBColor(0x2A, 0x2A, 0x4E))
    add_ml(s, 0.8, tk_y + 0.05, 11.4, 0.75, [
        ("Core Compute Takeaway:", True, ACCENT_TEAL),
        ("The 256 KB banked SRAM register file and its delivery infrastructure (operand collector, crossbar, "
         "bank arbitration) is the GPU's structural energy cost. It adds ~0.13-0.26 pJ/MAC in dynamic overhead "
         "plus ~0.13-0.26 pJ/MAC in static leakage.",
         False, LIGHT_GRAY),
        ("TPU/KPU eliminate this entirely with distributed PE-local latches. This is a 1.9x core compute "
         "advantage (GPU TC: 1.81 vs KPU: 0.96 pJ/op) that the GPU cannot close without abandoning SIMT.",
         True, YELLOW),
        ("* KPU mixed precision: 70% INT8 + 20% BF16 + 10% FP32. Single-precision BF16 KPU: ~1.13 pJ/op.",
         False, MED_GRAY),
    ], fs=11, align=PP_ALIGN.CENTER)

    return s


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)
    build_slide_7(prs)

    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "..", "core_compute_energy_v2.pptx")
    output_path = os.path.normpath(output_path)
    prs.save(output_path)
    print(f"Saved: {output_path}")
    print(f"  7 slides:")
    print(f"  1. Title: The Register File IS the Architecture")
    print(f"  2. CPU Core: Small RF, Amortized by SIMD")
    print(f"  3. GPU SM: 256 KB Banked SRAM + Infrastructure")
    print(f"  4. TPU/KPU: Distributed PE-Local Registers")
    print(f"  5. Dynamic Energy: Per-Op Operand Delivery Cost")
    print(f"  6. Static Power: Always-On SRAM Leakage")
    print(f"  7. Summary: Core Compute at 7nm")


if __name__ == "__main__":
    main()
