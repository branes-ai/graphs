#!/usr/bin/env python3
"""
Generate PowerPoint slides for energy model validation and troubleshooting.

This is the fourth deck -- a diagnostic/validation deck that spells out every
energy equation, every parameter source, and every assumption so that the
numbers in the investor slides can be independently verified.

Slide 1: Title + Methodology -- Where the numbers come from
Slide 2: Parameter Sources -- Process node, circuit multipliers, memory tech
Slide 3: Per-Architecture Energy Equations (all terms spelled out)
Slide 4: Side-by-Side Numerical Evaluation (1M ops, 1MB data)
Slide 5: Compute Efficiency Analysis + Model Gaps
"""

import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'src'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from graphs.hardware.technology_profile import (
    _create_category_cpu_profile, _create_category_gpu_profile,
    _create_category_tpu_profile, _create_category_kpu_profile,
    MemoryType, CIRCUIT_TYPE_MULTIPLIER, PROCESS_NODE_BASE_ENERGY_PJ,
)
from graphs.hardware.cycle_energy.cpu import build_cpu_cycle_energy
from graphs.hardware.cycle_energy.gpu import (
    build_gpu_cuda_cycle_energy, build_gpu_tensorcore_cycle_energy, SM_CONFIG,
)
from graphs.hardware.cycle_energy.tpu import build_tpu_cycle_energy
from graphs.hardware.cycle_energy.kpu import build_kpu_cycle_energy

# -- Color palette (matches all prior decks) ---------------------------------
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xD8)
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
MED_GRAY    = RGBColor(0x99, 0x99, 0x99)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)
RED_ACCENT  = RGBColor(0xE0, 0x40, 0x40)
GREEN       = RGBColor(0x2E, 0xCC, 0x71)
YELLOW      = RGBColor(0xF3, 0xC6, 0x23)

# Architecture colors
IF_COLOR      = RGBColor(0x3A, 0x86, 0xC8)
NVIDIA_GREEN  = RGBColor(0x76, 0xB9, 0x00)
TPU_COLOR     = RGBColor(0xDD, 0x88, 0x33)
KPU_BLUE      = RGBColor(0x00, 0x96, 0xD6)
ALU_COLOR     = RGBColor(0x2E, 0xCC, 0x71)

ROW_EVEN = RGBColor(0x2A, 0x2A, 0x4E)
ROW_ODD  = RGBColor(0x22, 0x22, 0x3E)
ROW_HDR  = RGBColor(0x33, 0x33, 0x55)


# ============================================================================
# Helpers (same API)
# ============================================================================

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=16, color=WHITE, alignment=PP_ALIGN.LEFT,
                          font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld, clr = item, False, color
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            clr = item[2] if len(item) > 2 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.bold = bld
        p.font.color.rgb = clr
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1.0) * 2)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     text="", font_size=14, font_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_rect(slide, left, top, width, height, fill_color,
             text="", font_size=14, font_color=WHITE, bold=False,
             line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_line(slide, x1, y1, x2, y2, color=MED_GRAY, width=1.5):
    connector = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def make_table(slide, left, top, width, height, headers, data,
               header_color=ACCENT_BLUE, font_size=9):
    """Helper to create a styled table."""
    rows = len(data) + 1
    cols = len(headers)
    tbl = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top),
        Inches(width), Inches(height)).table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color

    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_EVEN if i % 2 == 0 else ROW_ODD
    return tbl


# ============================================================================
# Run all energy models
# ============================================================================

N_OPS = 1_000_000
N_BYTES = 1_048_576  # 1 MB

cpu_prof = _create_category_cpu_profile(7, MemoryType.DDR5, 'datacenter')
gpu_prof = _create_category_gpu_profile(4, MemoryType.HBM3, 'datacenter')
tpu_prof = _create_category_tpu_profile(7, MemoryType.HBM2E, 'datacenter')
kpu_prof = _create_category_kpu_profile(7, MemoryType.LPDDR5, 'edge')

cpu_scalar_bd = build_cpu_cycle_energy(
    num_ops=N_OPS, bytes_transferred=N_BYTES, tech_profile=cpu_prof, simd_width=1)
cpu_avx_bd = build_cpu_cycle_energy(
    num_ops=N_OPS, bytes_transferred=N_BYTES, tech_profile=cpu_prof, simd_width=16)
gpu_cuda_bd = build_gpu_cuda_cycle_energy(
    num_ops=N_OPS, bytes_transferred=N_BYTES, tech_profile=gpu_prof)
gpu_tc_bd = build_gpu_tensorcore_cycle_energy(
    num_ops=N_OPS, bytes_transferred=N_BYTES, tech_profile=gpu_prof)
tpu_bd = build_tpu_cycle_energy(
    num_ops=N_OPS, bytes_transferred=N_BYTES, tech_profile=tpu_prof)
kpu_bd = build_kpu_cycle_energy(
    num_ops=N_OPS, bytes_transferred=N_BYTES, tech_profile=kpu_prof)


def categorize_events(bd):
    """Group events into categories, return dict of {cat: total_pJ}."""
    cats = {}
    for e in bd.events:
        pn = e.phase.name
        if 'INSTRUCTION' in pn:
            cat = 'Frontend'
        elif 'OPERAND' in pn:
            cat = 'Operand Fetch'
        elif 'EXECUTE' in pn or ('COMPUTE' in pn and 'SPATIAL' not in pn):
            cat = 'Compute (ALU)'
        elif 'SPATIAL_COMPUTE' in pn:
            cat = 'Compute (ALU)'
        elif 'WRITEBACK' in pn:
            cat = 'Writeback'
        elif 'MEM' in pn or 'EDDO' in pn:
            cat = 'Memory'
        elif 'SIMT' in pn:
            cat = 'SIMT Overhead'
        elif 'SPATIAL' in pn:
            cat = 'Spatial/Config'
        elif 'SYSTOLIC' in pn:
            cat = 'Systolic Ctrl'
        else:
            cat = 'Other'
        cats[cat] = cats.get(cat, 0) + e.total_energy_pj
    return cats


# ============================================================================
# SLIDE 1: Title + Methodology
# ============================================================================

def build_slide_1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.5, 0.5, 12.0, 0.8,
                "Energy Model Validation: Checking Our Numbers",
                font_size=32, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 1.2, 12.0, 0.5,
                "Every parameter, every equation, every assumption -- spelled out for verification",
                font_size=18, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    add_line(slide, 2.0, 1.9, 11.0, 1.9, color=ACCENT_BLUE, width=2)

    # Methodology
    add_textbox(slide, 0.5, 2.2, 12.0, 0.4,
                "Methodology", font_size=20, bold=True, color=WHITE)

    add_multiline_textbox(slide, 0.5, 2.7, 12.0, 2.5, [
        ("1. Base ALU Energy: Horowitz ISSCC 2014, scaled by process node capacitance (C*V^2)",
         False, LIGHT_GRAY),
        ("   E_base(node) = {3nm: 1.20, 4nm: 1.30, 5nm: 1.50, 7nm: 1.80, 16nm: 2.70} pJ",
         True, WHITE),
        ("", False, WHITE),
        ("2. Circuit Type Multiplier: Accounts for design complexity vs efficiency trade-off",
         False, LIGHT_GRAY),
        ("   E_alu = E_base * multiplier   (x86: 2.50x, cuda: 0.95x, tensor: 0.85x, systolic: 0.80x, domain: 0.75x)",
         True, WHITE),
        ("", False, WHITE),
        ("3. Architecture-Specific Overhead: Per-cycle infrastructure energy",
         False, LIGHT_GRAY),
        ("   CPU: instruction pipeline + register file (per instruction, amortized over SIMD width)",
         False, MED_GRAY),
        ("   GPU: warp scheduler + scoreboard + operand collector + crossbar (per warp, amortized over 32/64 MACs)",
         False, MED_GRAY),
        ("   TPU: weight loading + input feeding + output draining (per tile, amortized over MxM MACs)",
         False, MED_GRAY),
        ("   KPU: PE-to-PE wire + local register write + EDDO memory (per op, near-zero overhead)",
         False, MED_GRAY),
        ("", False, WHITE),
        ("4. Memory Hierarchy: Cache/scratchpad + off-chip access, workload-dependent",
         False, LIGHT_GRAY),
    ], font_size=12)

    # Workload specification
    add_textbox(slide, 0.5, 5.5, 12.0, 0.4,
                "Reference Workload for All Evaluations", font_size=16, bold=True, color=ORANGE)

    add_multiline_textbox(slide, 0.5, 5.95, 12.0, 1.2, [
        (f"  Operations: {N_OPS:,} MACs  |  Data transferred: {N_BYTES:,} bytes (1 MB)  |  "
         "Operating mode: DRAM-resident", False, WHITE),
        ("  This workload is large enough to amortize fixed overheads (kernel launch, domain program load)",
         False, MED_GRAY),
        ("  KPU precision mix: 70% INT8, 20% BF16, 10% FP32 (as configured in kpu.py)",
         False, MED_GRAY),
    ], font_size=12)

    # Source
    add_textbox(slide, 0.5, 7.0, 12.0, 0.3,
                "Source: graphs/hardware/technology_profile.py, graphs/hardware/cycle_energy/{cpu,gpu,tpu,kpu}.py",
                font_size=10, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    return slide


# ============================================================================
# SLIDE 2: Parameter Sources
# ============================================================================

def build_slide_2_parameters(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.15, 12.5, 0.5,
                "Parameter Sources: Process, Circuit, and Memory Technology",
                font_size=26, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # Table 1: Architecture configuration
    t1_y = 0.65
    add_textbox(slide, 0.5, t1_y, 12.0, 0.3,
                "Architecture Configuration (from technology_profile.py)",
                font_size=13, bold=True, color=ACCENT_TEAL)

    config_headers = [
        "Parameter", "CPU (7nm)", "GPU (4nm)", "TPU (7nm)", "KPU (7nm)", "Formula"
    ]
    config_data = [
        ["Process Node", "7nm (TSMC N7)", "4nm (TSMC N4P)", "7nm", "7nm", "---"],
        ["E_base (FP32 ALU)", f"{cpu_prof.base_alu_energy_pj:.2f} pJ",
         f"{gpu_prof.base_alu_energy_pj:.2f} pJ",
         f"{tpu_prof.base_alu_energy_pj:.2f} pJ",
         f"{kpu_prof.base_alu_energy_pj:.2f} pJ",
         "PROCESS_NODE_BASE_ENERGY_PJ[nm]"],
        ["Circuit Multiplier",
         f"x86_perf: {CIRCUIT_TYPE_MULTIPLIER['x86_performance']:.2f}x",
         f"tensor_core: {CIRCUIT_TYPE_MULTIPLIER['tensor_core']:.2f}x",
         f"systolic: {CIRCUIT_TYPE_MULTIPLIER['systolic_mac']:.2f}x",
         f"domain_flow: {CIRCUIT_TYPE_MULTIPLIER['domain_flow']:.2f}x",
         "CIRCUIT_TYPE_MULTIPLIER[type]"],
        ["E_mac (effective)",
         f"{cpu_prof.base_alu_energy_pj:.2f} pJ",
         f"{gpu_prof.tensor_core_mac_energy_pj:.3f} pJ",
         f"{tpu_prof.systolic_mac_energy_pj:.3f} pJ",
         f"{kpu_prof.domain_flow_mac_energy_pj:.3f} pJ",
         "E_base * multiplier"],
        ["Memory Technology", "DDR5", "HBM3", "HBM2e", "LPDDR5", "---"],
        ["Off-chip (pJ/byte)",
         f"{cpu_prof.offchip_energy_per_byte_pj:.1f}",
         f"{gpu_prof.offchip_energy_per_byte_pj:.1f}",
         f"{tpu_prof.offchip_energy_per_byte_pj:.1f}",
         f"{kpu_prof.offchip_energy_per_byte_pj:.1f}",
         "MemoryType spec"],
        ["SRAM/L1 (pJ/byte)",
         f"{cpu_prof.sram_energy_per_byte_pj:.3f}",
         f"{gpu_prof.sram_energy_per_byte_pj:.4f}",
         f"{tpu_prof.sram_energy_per_byte_pj:.3f}",
         f"{kpu_prof.sram_energy_per_byte_pj:.3f}",
         "get_sram_energy_per_byte_pj()"],
        ["I-fetch (pJ)",
         f"{cpu_prof.instruction_fetch_energy_pj:.3f}",
         f"{gpu_prof.instruction_fetch_energy_pj:.3f}",
         f"{tpu_prof.instruction_fetch_energy_pj:.3f}",
         f"{kpu_prof.instruction_fetch_energy_pj:.3f}",
         "E_base * pipeline_factor"],
        ["RF read (pJ)",
         f"{cpu_prof.register_read_energy_pj:.3f}",
         f"{gpu_prof.register_read_energy_pj:.3f}",
         f"{tpu_prof.register_read_energy_pj:.3f}",
         f"{kpu_prof.register_read_energy_pj:.3f}",
         "E_base * rf_factor"],
        ["RF write (pJ)",
         f"{cpu_prof.register_write_energy_pj:.3f}",
         f"{gpu_prof.register_write_energy_pj:.3f}",
         f"{tpu_prof.register_write_energy_pj:.3f}",
         f"{kpu_prof.register_write_energy_pj:.3f}",
         "E_base * rf_factor"],
    ]

    make_table(slide, 0.3, t1_y + 0.3, 12.6, 3.5,
               config_headers, config_data, font_size=8)

    # Key observations
    obs_y = 4.3
    add_line(slide, 0.5, obs_y, 12.5, obs_y, color=ORANGE, width=1.5)
    add_textbox(slide, 0.5, obs_y + 0.1, 12.0, 0.3,
                "Critical Observations About Parameter Consistency",
                font_size=14, bold=True, color=ORANGE)

    add_multiline_textbox(slide, 0.5, obs_y + 0.45, 12.0, 2.8, [
        ("1. Process Node Asymmetry: GPU at 4nm (E_base=1.30) vs TPU/KPU at 7nm (E_base=1.80)",
         True, YELLOW),
        ("   The 4nm GPU has a 28% base energy advantage before any circuit multiplier is applied.",
         False, LIGHT_GRAY),
        ("   For fair comparison, all architectures should be evaluated at the same process node.",
         False, MED_GRAY),
        ("", False, WHITE),
        ("2. Memory Technology Asymmetry: GPU uses HBM3 (5.5 pJ/byte) vs KPU uses LPDDR5 (8.0 pJ/byte)",
         True, YELLOW),
        ("   The GPU has a 31% memory energy advantage per byte. KPU compensates via EDDO (no tag lookup),",
         False, LIGHT_GRAY),
        ("   but the raw off-chip cost is higher. HBM variant of KPU would close this gap.",
         False, MED_GRAY),
        ("", False, WHITE),
        ("3. MAC Energy: TPU systolic = 1.440 pJ, KPU domain_flow = 1.350 pJ (both at 7nm)",
         True, YELLOW),
        ("   These ARE comparable, as expected. The 0.75x vs 0.80x multiplier reflects domain flow's",
         False, LIGHT_GRAY),
        ("   slightly simpler PE design (no weight-stationary control logic).",
         False, MED_GRAY),
        ("", False, WHITE),
        ("4. GPU tensor_core_mac = 1.105 pJ. This is E_base(4nm)*0.85 = 1.30*0.85.",
         True, YELLOW),
        ("   The Horowitz-derived switching energy for FP16 mult + FP32 add is 0.34 pJ.",
         False, LIGHT_GRAY),
        ("   The model uses the circuit-multiplier method, not Horowitz component analysis.",
         False, MED_GRAY),
    ], font_size=11)

    return slide


# ============================================================================
# SLIDE 3: Per-Architecture Energy Equations
# ============================================================================

def build_slide_3_equations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.1, 12.5, 0.45,
                "Energy Equations: Every Term Spelled Out",
                font_size=24, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    y = 0.55
    eq_font = "Consolas"
    eq_sz = 10
    label_sz = 12
    val_sz = 9

    # CPU
    add_rounded_rect(slide, 0.3, y, 1.6, 0.3, IF_COLOR,
                     "CPU (7nm scalar)", font_size=10, bold=True)
    y += 0.35
    add_textbox(slide, 0.3, y, 12.5, 0.25,
                "E_cpu = 2*(E_ifetch + E_idecode) + 2*E_rf_read + E_alu + E_rf_write + E_memory",
                font_size=eq_sz, bold=True, color=WHITE, font_name=eq_font)
    y += 0.25
    cpu_e = cpu_prof
    add_textbox(slide, 0.3, y, 12.5, 0.25,
                f"     = 2*({cpu_e.instruction_fetch_energy_pj:.3f} + {cpu_e.instruction_decode_energy_pj:.3f}) "
                f"+ 2*{cpu_e.register_read_energy_pj:.3f} + {cpu_e.base_alu_energy_pj:.3f} "
                f"+ {cpu_e.register_write_energy_pj:.3f} + E_mem",
                font_size=val_sz, color=LIGHT_GRAY, font_name=eq_font)
    y += 0.22
    # Evaluate
    cpu_overhead = 2*(cpu_e.instruction_fetch_energy_pj + cpu_e.instruction_decode_energy_pj)
    cpu_rf = 2*cpu_e.register_read_energy_pj + cpu_e.register_write_energy_pj
    cpu_total_no_mem = cpu_overhead + cpu_rf + cpu_e.base_alu_energy_pj
    cpu_cats = categorize_events(cpu_scalar_bd)
    cpu_mem = cpu_cats.get('Memory', 0) / N_OPS
    cpu_total = cpu_scalar_bd.total_energy_pj / N_OPS
    add_textbox(slide, 0.3, y, 12.5, 0.25,
                f"     = {cpu_overhead:.3f} (frontend) + {cpu_rf:.3f} (RF) + "
                f"{cpu_e.base_alu_energy_pj:.3f} (ALU) + {cpu_mem:.3f} (mem) "
                f"= {cpu_total:.2f} pJ/op  |  Compute eff: {cpu_e.base_alu_energy_pj/cpu_total*100:.1f}%",
                font_size=val_sz, bold=True, color=GREEN, font_name=eq_font)

    y += 0.38

    # GPU TensorCore
    add_rounded_rect(slide, 0.3, y, 2.0, 0.3, NVIDIA_GREEN,
                     "GPU TensorCore (4nm)", font_size=10, bold=True, font_color=BLACK)
    y += 0.35
    add_textbox(slide, 0.3, y, 12.5, 0.25,
                "E_gpu_tc = (E_ifetch + E_idecode + E_sched + E_scoreboard + E_addrgen + E_collector + E_arb + E_route)/64 + E_mac + (E_result_route + E_rf_write)/64 + E_simt + E_memory",
                font_size=eq_sz - 1, bold=True, color=WHITE, font_name=eq_font)
    y += 0.25
    gpu_cats = categorize_events(gpu_tc_bd)
    gpu_fe = gpu_cats.get('Frontend', 0) / N_OPS
    gpu_of = gpu_cats.get('Operand Fetch', 0) / N_OPS
    gpu_alu = gpu_cats.get('Compute (ALU)', 0) / N_OPS
    gpu_wb = gpu_cats.get('Writeback', 0) / N_OPS
    gpu_simt = gpu_cats.get('SIMT Overhead', 0) / N_OPS
    gpu_mem = gpu_cats.get('Memory', 0) / N_OPS
    gpu_total = gpu_tc_bd.total_energy_pj / N_OPS
    add_textbox(slide, 0.3, y, 12.5, 0.25,
                f"     = {gpu_fe:.4f} (front) + {gpu_of:.4f} (fetch) + {gpu_alu:.4f} (ALU) "
                f"+ {gpu_wb:.4f} (WB) + {gpu_simt:.4f} (SIMT) + {gpu_mem:.4f} (mem)",
                font_size=val_sz, color=LIGHT_GRAY, font_name=eq_font)
    y += 0.22
    add_textbox(slide, 0.3, y, 12.5, 0.25,
                f"     = {gpu_total:.4f} pJ/op  |  ALU = {gpu_prof.tensor_core_mac_energy_pj:.3f} pJ "
                f"= {gpu_prof.base_alu_energy_pj:.2f} * {CIRCUIT_TYPE_MULTIPLIER['tensor_core']:.2f}  "
                f"|  Compute eff: {gpu_alu/gpu_total*100:.1f}%",
                font_size=val_sz, bold=True, color=GREEN, font_name=eq_font)

    y += 0.38

    # GPU CUDA (for comparison)
    add_rounded_rect(slide, 0.3, y, 2.0, 0.3, NVIDIA_GREEN,
                     "GPU CUDA Core (4nm)", font_size=10, bold=True, font_color=BLACK)
    y += 0.35
    add_textbox(slide, 0.3, y, 12.5, 0.25,
                "E_gpu_cuda = (E_ifetch + E_idecode + E_sched + E_scoreboard + E_addrgen + E_collector + E_arb + E_route)/32 + E_mac + (E_result_route + E_rf_write)/32 + E_simt + E_memory",
                font_size=eq_sz - 1, bold=True, color=WHITE, font_name=eq_font)
    y += 0.25
    gc_cats = categorize_events(gpu_cuda_bd)
    gc_fe = gc_cats.get('Frontend', 0) / N_OPS
    gc_of = gc_cats.get('Operand Fetch', 0) / N_OPS
    gc_alu = gc_cats.get('Compute (ALU)', 0) / N_OPS
    gc_wb = gc_cats.get('Writeback', 0) / N_OPS
    gc_simt = gc_cats.get('SIMT Overhead', 0) / N_OPS
    gc_mem = gc_cats.get('Memory', 0) / N_OPS
    gc_total = gpu_cuda_bd.total_energy_pj / N_OPS
    add_textbox(slide, 0.3, y, 12.5, 0.25,
                f"     = {gc_fe:.4f} (front) + {gc_of:.4f} (fetch) + {gc_alu:.4f} (ALU) "
                f"+ {gc_wb:.4f} (WB) + {gc_simt:.4f} (SIMT) + {gc_mem:.4f} (mem) "
                f"= {gc_total:.4f} pJ/op  |  Compute eff: {gc_alu/gc_total*100:.1f}%",
                font_size=val_sz, bold=True, color=GREEN, font_name=eq_font)

    y += 0.38

    # TPU
    add_rounded_rect(slide, 0.3, y, 1.8, 0.3, TPU_COLOR,
                     "TPU Systolic (7nm)", font_size=10, bold=True, font_color=BLACK)
    y += 0.35
    add_textbox(slide, 0.3, y, 12.5, 0.25,
                "E_tpu = E_config/N_ops + E_weight_load/reuse + E_input_feed + E_systolic_mac + E_output_drain + E_memory",
                font_size=eq_sz, bold=True, color=WHITE, font_name=eq_font)
    y += 0.25
    tpu_cats = categorize_events(tpu_bd)
    tpu_ctrl = (tpu_cats.get('Systolic Ctrl', 0) + tpu_cats.get('Other', 0)) / N_OPS
    tpu_alu = tpu_cats.get('Compute (ALU)', 0) / N_OPS
    tpu_mem = tpu_cats.get('Memory', 0) / N_OPS
    tpu_total = tpu_bd.total_energy_pj / N_OPS
    add_textbox(slide, 0.3, y, 12.5, 0.25,
                f"     = {tpu_ctrl:.4f} (ctrl+data_move) + {tpu_alu:.4f} (ALU) + {tpu_mem:.4f} (mem)"
                f"  |  MAC = {tpu_prof.systolic_mac_energy_pj:.3f} pJ "
                f"= {tpu_prof.base_alu_energy_pj:.2f} * {CIRCUIT_TYPE_MULTIPLIER['systolic_mac']:.2f}",
                font_size=val_sz, color=LIGHT_GRAY, font_name=eq_font)
    y += 0.22
    add_textbox(slide, 0.3, y, 12.5, 0.25,
                f"     = {tpu_total:.4f} pJ/op  |  Compute eff: {tpu_alu/tpu_total*100:.1f}%",
                font_size=val_sz, bold=True, color=GREEN, font_name=eq_font)

    y += 0.38

    # KPU
    add_rounded_rect(slide, 0.3, y, 2.0, 0.3, KPU_BLUE,
                     "KPU Domain Flow (7nm)", font_size=10, bold=True)
    y += 0.35
    add_textbox(slide, 0.3, y, 12.5, 0.25,
                "E_kpu = E_config/N_ops + 2*E_pe_transfer + E_domain_mac + E_local_write + E_eddo_memory",
                font_size=eq_sz, bold=True, color=WHITE, font_name=eq_font)
    y += 0.25
    kpu_cats = categorize_events(kpu_bd)
    kpu_spatial = kpu_cats.get('Spatial/Config', 0) / N_OPS
    kpu_alu = kpu_cats.get('Compute (ALU)', 0) / N_OPS
    kpu_mem = kpu_cats.get('Memory', 0) / N_OPS
    kpu_total = kpu_bd.total_energy_pj / N_OPS
    add_textbox(slide, 0.3, y, 12.5, 0.25,
                f"     = {kpu_spatial:.4f} (spatial+config) + {kpu_alu:.4f} (ALU) + {kpu_mem:.4f} (mem)"
                f"  |  MAC(BF16) = {kpu_prof.domain_flow_mac_energy_pj:.3f}*0.8 = "
                f"{kpu_prof.domain_flow_mac_energy_pj*0.8:.3f} pJ "
                f"= {kpu_prof.base_alu_energy_pj:.2f} * {CIRCUIT_TYPE_MULTIPLIER['domain_flow']:.2f} * 0.8",
                font_size=val_sz, color=LIGHT_GRAY, font_name=eq_font)
    y += 0.22
    add_textbox(slide, 0.3, y, 12.5, 0.25,
                f"     = {kpu_total:.4f} pJ/op  (mixed: 70% INT8 + 20% BF16 + 10% FP32)  "
                f"|  Compute eff: {kpu_alu/kpu_total*100:.1f}%",
                font_size=val_sz, bold=True, color=GREEN, font_name=eq_font)

    return slide


# ============================================================================
# SLIDE 4: Side-by-Side Numerical Evaluation
# ============================================================================

def build_slide_4_evaluation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.1, 12.5, 0.45,
                f"Numerical Evaluation: {N_OPS:,} ops, {N_BYTES:,} bytes",
                font_size=24, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # Build the comprehensive comparison table
    all_bds = [
        ("CPU scalar\n(7nm)", cpu_scalar_bd, cpu_prof),
        ("CPU AVX-512\n(7nm)", cpu_avx_bd, cpu_prof),
        ("GPU CUDA\n(4nm)", gpu_cuda_bd, gpu_prof),
        ("GPU TC\n(4nm)", gpu_tc_bd, gpu_prof),
        ("TPU\n(7nm)", tpu_bd, tpu_prof),
        ("KPU\n(7nm)", kpu_bd, kpu_prof),
    ]

    categories = ['Frontend', 'Operand Fetch', 'Compute (ALU)', 'Writeback',
                  'SIMT Overhead', 'Systolic Ctrl', 'Spatial/Config', 'Memory']

    headers = ["Category (pJ/op)"] + [name for name, _, _ in all_bds]

    data = []
    for cat in categories:
        row = [cat]
        any_nonzero = False
        for _, bd, _ in all_bds:
            cats = categorize_events(bd)
            val = cats.get(cat, 0) / N_OPS
            if val > 0.0001:
                row.append(f"{val:.4f}")
                any_nonzero = True
            else:
                row.append("---")
        if any_nonzero:
            data.append(row)

    # Add totals row
    total_row = ["TOTAL"]
    for _, bd, _ in all_bds:
        total_row.append(f"{bd.total_energy_pj / N_OPS:.4f}")
    data.append(total_row)

    # Add compute efficiency row
    eff_row = ["Compute Efficiency"]
    for _, bd, _ in all_bds:
        cats = categorize_events(bd)
        compute = cats.get('Compute (ALU)', 0)
        total = bd.total_energy_pj
        eff_row.append(f"{compute/total*100:.1f}%")
    data.append(eff_row)

    # Add vs GPU TC row
    tc_total = gpu_tc_bd.total_energy_pj / N_OPS
    ratio_row = ["Energy ratio vs GPU TC"]
    for _, bd, _ in all_bds:
        t = bd.total_energy_pj / N_OPS
        ratio_row.append(f"{t/tc_total:.2f}x")
    data.append(ratio_row)

    tbl = make_table(slide, 0.3, 0.55, 12.6, 4.2,
                     headers, data, font_size=9)

    # Color the header cells by architecture
    arch_colors = [ROW_HDR, IF_COLOR, IF_COLOR, NVIDIA_GREEN,
                   NVIDIA_GREEN, TPU_COLOR, KPU_BLUE]
    for j in range(len(headers)):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = arch_colors[j] if j < len(arch_colors) else ACCENT_BLUE

    # Highlight total and efficiency rows
    n_data = len(data)
    for j in range(len(headers)):
        for special_row in [n_data - 2, n_data - 1, n_data]:  # total, eff, ratio
            cell = tbl.cell(special_row, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_HDR
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                if special_row == n_data:  # ratio row
                    p.font.color.rgb = ACCENT_TEAL

    # Analysis below table
    analysis_y = 5.0
    add_line(slide, 0.5, analysis_y, 12.5, analysis_y, color=ORANGE, width=1.5)
    add_textbox(slide, 0.5, analysis_y + 0.1, 12.0, 0.3,
                "Analysis: What Drives Each Architecture's Total",
                font_size=16, bold=True, color=ORANGE)

    # Compute the actual dominant cost for each
    analyses = []
    for name, bd, prof in all_bds:
        cats = categorize_events(bd)
        total = bd.total_energy_pj
        dominant = max(cats.items(), key=lambda x: x[1])
        analyses.append(f"{name.split(chr(10))[0]}: {dominant[0]} = {dominant[1]/total*100:.0f}%")

    add_multiline_textbox(slide, 0.5, analysis_y + 0.45, 12.0, 2.0, [
        ("Dominant energy category per architecture:", True, WHITE),
        (f"  {analyses[0]}  |  {analyses[1]}", False, LIGHT_GRAY),
        (f"  {analyses[2]}  |  {analyses[3]}", False, LIGHT_GRAY),
        (f"  {analyses[4]}  |  {analyses[5]}", False, LIGHT_GRAY),
        ("", False, WHITE),
        ("Key: GPU SIMT overhead includes kernel launch (100K pJ, amortized over 1M ops = 0.10 pJ/op).",
         False, MED_GRAY),
        ("KPU memory is 57% of total because LPDDR5 costs 8.0 pJ/byte vs GPU's HBM3 at 5.5 pJ/byte.",
         True, YELLOW),
        ("The KPU's 70% INT8 mix lowers its average MAC to 0.90 pJ vs TPU's pure 1.44 pJ systolic MAC.",
         False, MED_GRAY),
    ], font_size=11)

    return slide


# ============================================================================
# SLIDE 5: Compute Efficiency + Model Gaps
# ============================================================================

def build_slide_5_gaps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.1, 12.5, 0.45,
                "Compute Efficiency and Known Model Gaps",
                font_size=24, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # Compute efficiency visual
    eff_y = 0.6
    add_textbox(slide, 0.5, eff_y, 12.0, 0.3,
                "Compute Efficiency = ALU Energy / Total Energy  (what fraction is useful work)",
                font_size=13, bold=True, color=ACCENT_TEAL)

    bar_y = eff_y + 0.35
    bar_h = 0.4
    total_w = 10.0
    start_x = 2.8

    archs = [
        ("CPU scalar", cpu_scalar_bd, IF_COLOR),
        ("CPU AVX-512", cpu_avx_bd, IF_COLOR),
        ("GPU CUDA", gpu_cuda_bd, NVIDIA_GREEN),
        ("GPU TC", gpu_tc_bd, NVIDIA_GREEN),
        ("TPU", tpu_bd, TPU_COLOR),
        ("KPU (mixed)", kpu_bd, KPU_BLUE),
    ]

    for i, (name, bd, color) in enumerate(archs):
        cats = categorize_events(bd)
        eff = cats.get('Compute (ALU)', 0) / bd.total_energy_pj
        by = bar_y + i * (bar_h + 0.08)

        # Label
        add_textbox(slide, 0.3, by + 0.02, 2.4, bar_h,
                    name, font_size=10, bold=True, color=color,
                    alignment=PP_ALIGN.RIGHT)

        # Compute bar (green)
        eff_w = total_w * eff
        add_rect(slide, start_x, by, eff_w, bar_h, ALU_COLOR,
                 f"{eff*100:.1f}%", font_size=9, font_color=BLACK, bold=True)
        # Overhead bar (dark)
        add_rect(slide, start_x + eff_w, by, total_w - eff_w, bar_h,
                 RGBColor(0x44, 0x44, 0x66),
                 f"overhead: {(1-eff)*100:.1f}%", font_size=8, font_color=MED_GRAY)

    # Model gaps section
    gaps_y = bar_y + 6 * (bar_h + 0.08) + 0.15
    add_line(slide, 0.5, gaps_y, 12.5, gaps_y, color=RED_ACCENT, width=2)
    add_textbox(slide, 0.5, gaps_y + 0.1, 12.0, 0.3,
                "Known Model Gaps and Issues to Resolve",
                font_size=18, bold=True, color=RED_ACCENT)

    add_multiline_textbox(slide, 0.5, gaps_y + 0.45, 12.0, 3.0, [
        ("1. GPU SM Baseline Power Not Modeled", True, YELLOW),
        ("   The 64K register file (256 KB SRAM), 4 operand collectors, 4 scoreboards, and 32-bank crossbar",
         False, LIGHT_GRAY),
        ("   draw continuous power even between instructions. This static infrastructure energy is NOT captured",
         False, LIGHT_GRAY),
        ("   in the per-instruction cycle model. It would increase GPU total and decrease GPU compute efficiency.",
         False, MED_GRAY),
        ("", False, WHITE),
        ("2. Process Node Mismatch", True, YELLOW),
        ("   GPU at 4nm (E_base=1.30) vs TPU/KPU at 7nm (E_base=1.80). The GPU gets a 28% base advantage.",
         False, LIGHT_GRAY),
        ("   Need: same-node comparison (all at 4nm or all at 7nm) to isolate architectural effects.",
         False, MED_GRAY),
        ("", False, WHITE),
        ("3. Memory Technology Mismatch", True, YELLOW),
        ("   GPU HBM3 (5.5 pJ/byte) vs KPU LPDDR5 (8.0 pJ/byte). The GPU has 31% lower off-chip energy.",
         False, LIGHT_GRAY),
        ("   At 1MB data, memory is 26% of GPU energy but 57% of KPU energy. Need: same memory comparison.",
         False, MED_GRAY),
        ("", False, WHITE),
        ("4. KPU Precision Mix", True, YELLOW),
        ("   KPU default is 70% INT8 + 20% BF16 + 10% FP32 (avg MAC = 0.90 pJ).",
         False, LIGHT_GRAY),
        ("   TPU is pure systolic at 1.44 pJ. Single-precision comparison would change the efficiency ratio.",
         False, MED_GRAY),
    ], font_size=10)

    return slide


# ============================================================================
# MAIN
# ============================================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1_title(prs)
    build_slide_2_parameters(prs)
    build_slide_3_equations(prs)
    build_slide_4_evaluation(prs)
    build_slide_5_gaps(prs)

    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "..", "energy_validation.pptx")
    output_path = os.path.normpath(output_path)
    prs.save(output_path)
    print(f"Saved presentation to: {output_path}")
    print(f"  5 slides generated:")
    print(f"  1. Title: Methodology and Workload Specification")
    print(f"  2. Parameter Sources: Process, Circuit, Memory Technology")
    print(f"  3. Energy Equations: Every Term Spelled Out (all 5 architectures)")
    print(f"  4. Numerical Evaluation: Side-by-Side Comparison Table")
    print(f"  5. Compute Efficiency and Known Model Gaps")
    print()

    # Also print summary to console
    print("=" * 80)
    print("SUMMARY (from actual model runs)")
    print("=" * 80)
    for name, bd, prof in [
        ("CPU scalar (7nm)", cpu_scalar_bd, cpu_prof),
        ("CPU AVX-512 (7nm)", cpu_avx_bd, cpu_prof),
        ("GPU CUDA (4nm)", gpu_cuda_bd, gpu_prof),
        ("GPU TensorCore (4nm)", gpu_tc_bd, gpu_prof),
        ("TPU Systolic (7nm)", tpu_bd, tpu_prof),
        ("KPU Domain Flow (7nm)", kpu_bd, kpu_prof),
    ]:
        cats = categorize_events(bd)
        total = bd.total_energy_pj / N_OPS
        compute = cats.get('Compute (ALU)', 0) / N_OPS
        eff = compute / total * 100
        print(f"  {name:<28s}  {total:>8.4f} pJ/op  "
              f"ALU: {compute:.4f}  Eff: {eff:.1f}%")


if __name__ == "__main__":
    main()
