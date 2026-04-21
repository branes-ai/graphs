#!/usr/bin/env python
"""
Micro-architectural Validation Report Generator

Produces the Branes-branded HTML + JSON + PowerPoint report for the
9-layer micro-architectural hierarchy (Layers 1-7 populated during
M1-M7; Layers 8-9 deferred).

M0 ships the skeleton: empty layer panels, cross-SKU comparison shell,
landing page, and JSON data contract. Each subsequent milestone fills
in one layer's content.

See docs/plans/microarch-model-delivery-plan.md.

Usage:
    ./cli/microarch_validation_report.py --hardware jetson_orin_agx_64gb
    ./cli/microarch_validation_report.py --hardware jetson_orin_agx_64gb \
        --output tmp/report --format html
    ./cli/microarch_validation_report.py --hardware jetson_orin_agx_64gb \
        --format pptx --output tmp/report
    ./cli/microarch_validation_report.py --hardware jetson_orin_agx_64gb \
        --format json --output tmp/report

Multiple hardware may be specified; the CLI emits one per-SKU page plus
an index and a cross-SKU comparison page.
"""
from __future__ import annotations

import argparse
import sys
from datetime import datetime
from pathlib import Path
from typing import List

# Add repo root so `graphs.*` imports work when invoked directly.
_repo_root = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(_repo_root))

from graphs.reporting import (  # noqa: E402
    MicroarchReport,
    empty_report,
    render_sku_page,
    render_comparison_page,
    render_index_page,
)
from graphs.reporting.compare_archetypes import (  # noqa: E402
    build_default_comparison,
    render_archetype_page,
)
from graphs.reporting.native_op_energy import (  # noqa: E402
    build_default_comparison as build_native_op_comparison,
    render_native_op_page,
)
from graphs.reporting.generalized_architecture import (  # noqa: E402
    build_default_report as build_generalized_report,
    render_generalized_page,
)
from graphs.reporting.competitive_trajectory import (  # noqa: E402
    build_default_report as build_trajectory_report,
    render_trajectory_page,
)
from graphs.reporting.microarch_accounting import (  # noqa: E402
    build_default_report as build_accounting_report,
    render_accounting_page,
)
from graphs.hardware.resource_model import Precision  # noqa: E402


DEFAULT_SKU_LIST = [
    "jetson_orin_agx_64gb",
    "intel_core_i7_12700k",
    "ryzen_9_8945hs",
    "kpu_t64",
    "kpu_t128",
    "kpu_t256",
    "coral_edge_tpu",
    "hailo8",
    "hailo10h",
]


# Map SKU to archetype for the comparison view. Populated manually for
# M0 since M0.5 defines the archetype story more fully.
SKU_ARCHETYPE = {
    "jetson_orin_agx_64gb": "simt",
    "intel_core_i7_12700k": "cpu",
    "ryzen_9_8945hs": "cpu",
    "kpu_t64": "domainflow",
    "kpu_t128": "domainflow",
    "kpu_t256": "domainflow",
    "coral_edge_tpu": "systolic",
    "hailo8": "dsp",
    "hailo10h": "dsp",
}


def build_empty_report_for(sku: str) -> MicroarchReport:
    report = empty_report(sku=sku, display_name=sku)
    report.archetype = SKU_ARCHETYPE.get(sku, "")
    return report


def write_json_bundle(reports: List[MicroarchReport], out_dir: Path) -> List[Path]:
    """Write per-SKU JSON files under <out_dir>/data/."""
    data_dir = out_dir / "data"
    data_dir.mkdir(parents=True, exist_ok=True)
    written = []
    for r in reports:
        p = data_dir / f"{r.sku}.json"
        r.save(p)
        written.append(p)
    return written


def write_html_bundle(reports: List[MicroarchReport], out_dir: Path) -> List[Path]:
    """Write index.html, compare.html, compare_archetypes.html, and one hardware/<sku>.html per SKU."""
    hw_dir = out_dir / "hardware"
    hw_dir.mkdir(parents=True, exist_ok=True)
    written = []
    # Index
    index_path = out_dir / "index.html"
    index_path.write_text(render_index_page(reports, _repo_root))
    written.append(index_path)
    # Per-SKU
    for r in reports:
        p = hw_dir / f"{r.sku}.html"
        p.write_text(render_sku_page(r, _repo_root))
        written.append(p)
    # Comparison shell (M0 stub, M8 fills in)
    compare_path = out_dir / "compare.html"
    compare_path.write_text(render_comparison_page(reports, _repo_root))
    written.append(compare_path)
    # Compute-archetype comparison (M0.5)
    try:
        archetype_report = build_default_comparison(
            precision=Precision.INT8,
            kpu_sku="Stillwater-KPU-T128",
            kpu_display_name="KPU T128",
        )
        arch_path = out_dir / "compare_archetypes.html"
        arch_path.write_text(render_archetype_page(archetype_report, _repo_root))
        written.append(arch_path)
    except RuntimeError as exc:
        import sys as _sys
        print(f"warning: compare_archetypes.html skipped ({exc})", file=_sys.stderr)
    # Native-op energy breakdown (M0.5) - specific shipping products
    try:
        native_op_report = build_native_op_comparison(
            precision=Precision.INT8, kpu_sku="Stillwater-KPU-T128",
        )
        native_path = out_dir / "native_op_energy.html"
        native_path.write_text(render_native_op_page(native_op_report, _repo_root))
        written.append(native_path)
    except RuntimeError as exc:
        import sys as _sys
        print(f"warning: native_op_energy.html skipped ({exc})", file=_sys.stderr)
    # Generalized architecture comparison (M0.5) - process-normalized
    try:
        gen_report = build_generalized_report(
            reference_process_nm=16, power_budget_w=12.0,
        )
        gen_path = out_dir / "generalized_architecture.html"
        gen_path.write_text(render_generalized_page(gen_report, _repo_root))
        written.append(gen_path)
    except RuntimeError as exc:
        import sys as _sys
        print(f"warning: generalized_architecture.html skipped ({exc})", file=_sys.stderr)
    # Competitive trajectory (M0.5) - Jetson history + KPU target extrapolation
    try:
        traj_report = build_trajectory_report()
        traj_path = out_dir / "competitive_trajectory.html"
        traj_path.write_text(render_trajectory_page(traj_report, _repo_root))
        written.append(traj_path)
    except RuntimeError as exc:
        import sys as _sys
        print(f"warning: competitive_trajectory.html skipped ({exc})", file=_sys.stderr)
    # Per-structure micro-arch accounting (M0.5) - SM vs KPU tile validation
    try:
        acct_report = build_accounting_report()
        acct_path = out_dir / "microarch_accounting.html"
        acct_path.write_text(render_accounting_page(acct_report, _repo_root))
        written.append(acct_path)
    except RuntimeError as exc:
        import sys as _sys
        print(f"warning: microarch_accounting.html skipped ({exc})", file=_sys.stderr)
    return written


def write_pptx_bundle(reports: List[MicroarchReport], out_dir: Path) -> Path:
    """
    Emit a minimal branded PowerPoint deck. M8 replaces this with the
    full layer-by-layer walkthrough.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise RuntimeError(
            "python-pptx is not installed. Install it with "
            "`pip install python-pptx` or skip --format pptx."
        )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    logo = _repo_root / "docs" / "img" / "Branes_Logo.jpg"
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(0.4), Inches(0.3), height=Inches(0.6))
    title_tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    title_tf = title_tb.text_frame
    title_tf.text = "Micro-architectural model delivery"
    title_tf.paragraphs[0].runs[0].font.size = Pt(44)
    subtitle_tb = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(1))
    subtitle_tb.text_frame.text = (
        f"Layers 1-7 across {len(reports)} SKU"
        f"{'s' if len(reports) != 1 else ''}  -  "
        f"Generated {datetime.now().strftime('%Y-%m-%d')}"
    )
    subtitle_tb.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    # One slide per SKU (placeholder during M0)
    for r in reports:
        slide = prs.slides.add_slide(blank)
        if logo.exists():
            slide.shapes.add_picture(str(logo), Inches(0.4), Inches(0.3), height=Inches(0.5))
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(1))
        tb.text_frame.text = r.display_name or r.sku
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
        body = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12), Inches(5))
        body_tf = body.text_frame
        body_tf.text = (
            f"SKU: {r.sku}\n"
            f"Archetype: {r.archetype or 'unspecified'}\n"
            f"Overall confidence: {r.overall_confidence}"
        )
        body_tf.paragraphs[0].runs[0].font.size = Pt(16)
        body_tf.add_paragraph().text = ""
        body_tf.add_paragraph().text = (
            "Layer panels (1-7): NOT YET POPULATED at M0. "
            "Populated by milestones M1-M7."
        )

    deck_path = out_dir / "microarch_deck.pptx"
    prs.save(str(deck_path))
    return deck_path


def parse_args(argv: List[str]) -> argparse.Namespace:
    p = argparse.ArgumentParser(
        prog="microarch_validation_report",
        description=("Generate Branes-branded micro-architectural model "
                     "delivery report (Layers 1-7)."),
    )
    p.add_argument(
        "--hardware",
        nargs="+",
        default=DEFAULT_SKU_LIST,
        help="SKU IDs to include. Default: the locked M5 delivery SKU set.",
    )
    p.add_argument(
        "--layer",
        nargs="+",
        default=None,
        help=("Filter layers to include in the panel display. Not yet "
              "used at M0; reserved for M1-M7."),
    )
    p.add_argument(
        "--precision",
        nargs="+",
        default=None,
        help="Filter precisions. Reserved for M1-M7.",
    )
    p.add_argument(
        "--output", "-o",
        default="reports/microarch_model/latest",
        help="Output directory. Default: reports/microarch_model/latest",
    )
    p.add_argument(
        "--format",
        choices=["html", "pptx", "json", "all"],
        default="all",
        help=("Output format. 'all' emits html + json + pptx. "
              "Default: all."),
    )
    p.add_argument(
        "--serve",
        action="store_true",
        help="Reserved for future use (serve HTML over http.server).",
    )
    p.add_argument(
        "--verbose", "-v",
        action="store_true",
        help="Verbose logging.",
    )
    return p.parse_args(argv)


def main(argv: List[str]) -> int:
    args = parse_args(argv)
    out_dir = Path(args.output).resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    reports = [build_empty_report_for(sku) for sku in args.hardware]

    produced: List[Path] = []
    fmt = args.format
    if fmt in ("json", "all"):
        produced.extend(write_json_bundle(reports, out_dir))
    if fmt in ("html", "all"):
        produced.extend(write_html_bundle(reports, out_dir))
    if fmt in ("pptx", "all"):
        try:
            produced.append(write_pptx_bundle(reports, out_dir))
        except RuntimeError as exc:
            print(f"warning: PPT export skipped ({exc})", file=sys.stderr)
            if fmt == "pptx":
                return 1

    print(f"wrote {len(produced)} file(s) to {out_dir}")
    if args.verbose:
        for p in produced:
            print(f"  {p}")
    if args.serve:
        print("--serve is reserved for future use", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
