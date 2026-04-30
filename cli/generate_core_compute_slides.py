#!/usr/bin/env python3
"""
Generate PowerPoint slides comparing core compute engine energy across architectures.

This deck isolates the CORE COMPUTE ENGINE only -- no external memory, no NoC,
no memory controllers. All architectures evaluated at 7nm for fair comparison.

The three modeling domains:
  1. Core Compute Engine  <-- THIS DECK (workload-independent)
  2. On-chip memory resources (workload-dependent)
  3. Memory controller + I/O + external memory (system-dependent)

Slide 1: Title -- The Three Modeling Domains
Slide 2: CPU Core -- Scalar ALU + SIMD ALU + Register Files
Slide 3: GPU SM/CU -- SIMT Engine + TensorCores + 64K Register File
Slide 4: TPU Core -- Input Buffer + 128x128 Systolic Array
Slide 5: KPU Tile -- Input Buffer + 16x16 PE Array + SURE Network
Slide 6: Side-by-Side -- Core Compute Energy Equations at 7nm
"""

import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'src'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# -- Color palette -----------------------------------------------------------
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xD8)
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
MED_GRAY    = RGBColor(0x99, 0x99, 0x99)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)
RED_ACCENT  = RGBColor(0xE0, 0x40, 0x40)
GREEN       = RGBColor(0x2E, 0xCC, 0x71)
YELLOW      = RGBColor(0xF3, 0xC6, 0x23)

CPU_COLOR     = RGBColor(0x3A, 0x86, 0xC8)
GPU_COLOR     = RGBColor(0x76, 0xB9, 0x00)
TPU_COLOR     = RGBColor(0xDD, 0x88, 0x33)
KPU_COLOR     = RGBColor(0x00, 0x96, 0xD6)
ALU_COLOR     = RGBColor(0x2E, 0xCC, 0x71)
RF_COLOR      = RGBColor(0x6C, 0x5C, 0xE7)
SIMT_COLOR    = RGBColor(0x8B, 0x5C, 0xF6)
COLLECTOR_CLR = RGBColor(0xF3, 0x96, 0x23)
CROSSBAR_CLR  = RGBColor(0xCC, 0x55, 0x77)
BANK_COLOR    = RGBColor(0x48, 0x7D, 0xB3)
CACHE_COLOR   = RGBColor(0x48, 0x7D, 0xB3)
SURE_COLOR    = RGBColor(0x00, 0xCC, 0x88)
PE_COLOR      = RGBColor(0x3A, 0xB0, 0xD0)
BUF_COLOR     = RGBColor(0x55, 0x88, 0xBB)

ROW_EVEN = RGBColor(0x2A, 0x2A, 0x4E)
ROW_ODD  = RGBColor(0x22, 0x22, 0x3E)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=16, color=WHITE, alignment=PP_ALIGN.LEFT,
                          font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld, clr = item, False, color
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            clr = item[2] if len(item) > 2 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.bold = bld
        p.font.color.rgb = clr
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1.0) * 2)
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color,
                     text="", font_size=14, font_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
    return shape

def add_rect(slide, left, top, width, height, fill_color,
             text="", font_size=14, font_color=WHITE, bold=False,
             line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
    return shape

def add_arrow_right(slide, left, top, width=0.3, height=0.2, color=MED_GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_arrow_down(slide, left, top, width=0.2, height=0.3, color=MED_GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_line(slide, x1, y1, x2, y2, color=MED_GRAY, width=1.5):
    c = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


# ============================================================================
# SLIDE 1: Title -- The Three Modeling Domains
# ============================================================================

def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.5, 0.4, 12.0, 0.8,
                "Core Compute Engine Energy Comparison",
                font_size=34, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 1.1, 12.0, 0.5,
                "Isolating Architecture from System: All at 7nm, No External Memory",
                font_size=20, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    add_line(slide, 2.5, 1.8, 10.5, 1.8, color=ACCENT_BLUE, width=2)

    # Three domains
    dom_y = 2.1
    add_textbox(slide, 0.5, dom_y, 12.0, 0.35,
                "Three Distinct Modeling Domains:", font_size=18, bold=True, color=WHITE)

    domains = [
        ("1. Core Compute Engine", GREEN,
         "Scalar ALU, SIMD ALU, SM/CU,\nSystolic Array, PE Array\n+ their local register files/buffers\n\n"
         "NOT workload-dependent.\nDetermined by architecture.",
         "THIS DECK"),
        ("2. On-Chip Memory", ORANGE,
         "L1/L2/L3 caches (CPU)\nShared Memory + L2 (GPU)\nUnified Buffer (TPU)\nEDDO Scratchpads (KPU)\n\n"
         "Workload-dependent.\nData reuse pattern matters.",
         "Future deck"),
        ("3. Memory Controller + I/O", RED_ACCENT,
         "NoC / Crossbar to mem ctrl\nDDR/LPDDR/GDDR/HBM PHY\nExternal DRAM access\n\n"
         "System-dependent.\nPackaging + memory technology.",
         "Memory Tech deck"),
    ]

    for i, (title, color, desc, status) in enumerate(domains):
        x = 0.5 + i * 4.2
        add_rounded_rect(slide, x, dom_y + 0.4, 3.8, 0.45, color,
                         title, font_size=14, bold=True,
                         font_color=WHITE if color != ORANGE else RGBColor(0x1A, 0x1A, 0x2E))
        add_textbox(slide, x + 0.1, dom_y + 0.95, 3.6, 1.8,
                    desc, font_size=11, color=LIGHT_GRAY)
        add_rounded_rect(slide, x + 0.5, dom_y + 2.75, 2.8, 0.3,
                         RGBColor(0x2A, 0x2A, 0x4E), status,
                         font_size=10, bold=True, font_color=color)

    # Architectures being compared
    arch_y = 5.3
    add_line(slide, 0.5, arch_y, 12.5, arch_y, color=ACCENT_BLUE, width=1)
    add_textbox(slide, 0.5, arch_y + 0.1, 12.0, 0.3,
                "Core Compute Engines Compared (all at 7nm TSMC N7):",
                font_size=16, bold=True, color=WHITE)

    engines = [
        ("CPU", "Scalar ALU + RF\nSIMD ALU (AVX-512) + VRF\nGranularity: 1 or 16 ops/instr",
         CPU_COLOR),
        ("GPU SM", "32 CUDA Cores + 1 TensorCore\n+ 16K Register Bank\n+ Warp Sched + Scoreboard\n+ Operand Collector + Crossbar\nGranularity: 32 or 64 ops/instr\nx4 partitions = 1 SM",
         GPU_COLOR),
        ("TPU Tile", "Input Buffer (SRAM)\n+ 128x128 Systolic Array\n+ Accumulators\nGranularity: 16,384 MACs/cycle\nWeight-stationary only",
         TPU_COLOR),
        ("KPU Tile", "Input Buffer (EDDO L1)\n+ 16x16 PE Array\n+ SURE Network\nGranularity: 256 MACs/cycle\nA/B/C-stationary (any schedule)",
         KPU_COLOR),
    ]

    for i, (name, desc, color) in enumerate(engines):
        x = 0.5 + i * 3.2
        add_rounded_rect(slide, x, arch_y + 0.45, 2.9, 0.35, color,
                         name, font_size=13, bold=True)
        add_textbox(slide, x + 0.1, arch_y + 0.85, 2.7, 1.5,
                    desc, font_size=9, color=LIGHT_GRAY)

    return slide


# ============================================================================
# SLIDE 2: CPU Core
# ============================================================================

def build_slide_2_cpu(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.1, 12.5, 0.45,
                "CPU Core Compute Engine: Scalar + SIMD ALU",
                font_size=26, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # -- Scalar Path (left) --
    sx = 0.4
    sy = 0.7
    add_textbox(slide, sx, sy, 5.5, 0.3,
                "Scalar Path (1 op per instruction)", font_size=14, bold=True, color=CPU_COLOR)

    scalar_steps = [
        ("I-Cache Read", CACHE_COLOR, "0.540 pJ"),
        ("Instruction Decode", RGBColor(0x5E, 0x60, 0xCE), "0.288 pJ"),
        ("RF Read (src1)", RF_COLOR, "1.013 pJ"),
        ("RF Read (src2)", RF_COLOR, "1.013 pJ"),
        ("Scalar ALU (FP32 FMA)", ALU_COLOR, "1.800 pJ"),
        ("RF Write (result)", RGBColor(0xE0, 0x40, 0x40), "1.215 pJ"),
    ]

    step_h = 0.4
    for i, (name, color, energy) in enumerate(scalar_steps):
        y = sy + 0.35 + i * (step_h + 0.05)
        add_rounded_rect(slide, sx, y, 2.8, step_h, color,
                         name, font_size=10, bold=True)
        e_clr = GREEN if "ALU" in name else ORANGE
        add_textbox(slide, sx + 2.9, y + 0.05, 0.9, step_h,
                    energy, font_size=10, bold=True, color=e_clr)
        if i < len(scalar_steps) - 1:
            add_arrow_down(slide, sx + 1.3, y + step_h, 0.15, 0.05, MED_GRAY)

    # Scalar total
    sc_total_y = sy + 0.35 + len(scalar_steps) * (step_h + 0.05) + 0.1
    add_textbox(slide, sx, sc_total_y, 4.5, 0.3,
                "Total: 2 * (0.540 + 0.288) + 2 * 1.013 + 1.800 + 1.215",
                font_size=9, color=MED_GRAY, font_name="Consolas")
    add_textbox(slide, sx, sc_total_y + 0.25, 4.5, 0.25,
                "= 9.936 pJ/op  |  Compute eff: 18.1%",
                font_size=12, bold=True, color=CPU_COLOR, font_name="Consolas")

    # -- SIMD Path (right) --
    vx = 6.5
    vy = 0.7
    add_textbox(slide, vx, vy, 6.0, 0.3,
                "SIMD Path (AVX-512: 16 ops per instruction)", font_size=14, bold=True, color=CPU_COLOR)

    simd_steps = [
        ("I-Cache Read (1 instr)", CACHE_COLOR, "0.540 / 16 = 0.034 pJ/op"),
        ("Decode (1 instr)", RGBColor(0x5E, 0x60, 0xCE), "0.288 / 16 = 0.018 pJ/op"),
        ("VRF Read src1 (1 vec reg)", RF_COLOR, "1.013 / 16 = 0.063 pJ/op"),
        ("VRF Read src2 (1 vec reg)", RF_COLOR, "1.013 / 16 = 0.063 pJ/op"),
        ("SIMD ALU x16 (FP32 FMA)", ALU_COLOR, "16 x 1.800 / 16 = 1.800 pJ/op"),
        ("VRF Write (1 vec reg)", RGBColor(0xE0, 0x40, 0x40), "1.215 / 16 = 0.076 pJ/op"),
    ]

    for i, (name, color, energy) in enumerate(simd_steps):
        y = vy + 0.35 + i * (step_h + 0.05)
        add_rounded_rect(slide, vx, y, 2.8, step_h, color,
                         name, font_size=10, bold=True)
        e_clr = GREEN if "ALU" in name else GREEN
        add_textbox(slide, vx + 2.9, y + 0.02, 3.5, step_h,
                    energy, font_size=9, bold=True, color=e_clr if "ALU" in name else LIGHT_GRAY,
                    font_name="Consolas")
        if i < len(simd_steps) - 1:
            add_arrow_down(slide, vx + 1.3, y + step_h, 0.15, 0.05, MED_GRAY)

    si_total_y = sc_total_y
    add_textbox(slide, vx, si_total_y, 6.0, 0.3,
                "(0.540 + 0.288 + 2*1.013 + 1.215) / 16  +  1.800",
                font_size=9, color=MED_GRAY, font_name="Consolas")
    add_textbox(slide, vx, si_total_y + 0.25, 6.0, 0.25,
                "= 2.309 pJ/op  |  Compute eff: 78.0%  |  4.3x better than scalar",
                font_size=12, bold=True, color=CPU_COLOR, font_name="Consolas")

    # Key insight
    ins_y = si_total_y + 0.65
    add_line(slide, 0.5, ins_y, 12.5, ins_y, color=CPU_COLOR, width=1.5)
    add_rounded_rect(slide, 0.5, ins_y + 0.1, 12.0, 0.6, RGBColor(0x2A, 0x2A, 0x4E))
    add_multiline_textbox(slide, 0.8, ins_y + 0.15, 11.4, 0.5, [
        ("CPU Energy Equation:  E_cpu = (E_ifetch + E_decode + 2*E_rf_read + E_rf_write) / SIMD_width  +  E_alu",
         True, WHITE),
        ("SIMD amortizes the instruction overhead over 16 lanes. The ALU energy (1.800 pJ) is identical.",
         False, LIGHT_GRAY),
        ("Wider SIMD = higher efficiency, but the instruction pipeline ALWAYS exists. Every op needs an instruction.",
         True, ORANGE),
    ], font_size=11, font_name="Calibri")

    # Granularity box
    add_textbox(slide, 0.5, ins_y + 0.8, 12.0, 0.3,
                "Granularity: 1 core = 1 scalar ALU + 1 SIMD ALU  |  "
                "Multi-core: N cores share L3  |  Each core has private L1i/L1d + L2",
                font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    return slide


# ============================================================================
# SLIDE 3: GPU SM/CU
# ============================================================================

def build_slide_3_gpu(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.1, 12.5, 0.45,
                "GPU SM: SIMT Engine + TensorCores + 64K Register File",
                font_size=26, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.3, 0.55, 12.5, 0.3,
                "One SM partition (x4 per SM): the indivisible GPU compute block",
                font_size=13, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    # -- SM Partition diagram (left) --
    px = 0.4
    py = 0.95
    pw = 5.5

    add_rect(slide, px, py, pw, 5.0, RGBColor(0x22, 0x22, 0x3E), "",
             line_color=GPU_COLOR)
    add_textbox(slide, px, py + 0.03, pw, 0.25,
                "SM Partition (1 of 4)", font_size=11, bold=True,
                color=GPU_COLOR, alignment=PP_ALIGN.CENTER)

    # Steps with energy
    parts = [
        ("Warp Scheduler", SIMT_COLOR, "0.875 pJ/instr"),
        ("Scoreboard (RAW/WAW/WAR)", RED_ACCENT, "0.525 pJ/instr"),
        ("Register Addr Gen (x3)", BANK_COLOR, "1.050 pJ/instr"),
        ("Operand Collector", COLLECTOR_CLR, "1.400 pJ (CUDA) / 0.875 pJ (TC)"),
        ("Bank Arbitration (32 banks)", BANK_COLOR, "0.525 pJ (CUDA) / 0.350 pJ (TC)"),
        ("Register File Bank (16K regs)", RF_COLOR, "0.281 pJ x 2 reads (CUDA)"),
        ("Crossbar (32-wide)", CROSSBAR_CLR, "0.700 pJ (CUDA) / 0.525 pJ (TC)"),
        ("32 CUDA Cores / 1 TensorCore", ALU_COLOR, "1.530 pJ/op (both paths)"),
        ("Result Route + RF Write", RED_ACCENT, "0.525 + 0.338 pJ (CUDA)"),
    ]

    sh = 0.38
    sg = 0.04
    for i, (name, color, energy) in enumerate(parts):
        y = py + 0.35 + i * (sh + sg)
        add_rounded_rect(slide, px + 0.1, y, 2.8, sh, color,
                         name, font_size=8, bold=True)
        add_textbox(slide, px + 3.0, y + 0.04, 2.3, sh,
                    energy, font_size=7, bold=True,
                    color=GREEN if "CUDA Cores" in name else LIGHT_GRAY,
                    font_name="Consolas")

    # -- Right side: Energy equations --
    ex = 6.3
    ey = 0.95

    add_textbox(slide, ex, ey, 6.5, 0.3,
                "Energy Equations (per MAC, amortized over warp)", font_size=14,
                bold=True, color=GPU_COLOR)

    # CUDA Core path
    add_textbox(slide, ex, ey + 0.35, 6.5, 0.25,
                "CUDA Core Path (amortized over 32 threads):", font_size=12,
                bold=True, color=LIGHT_GRAY)
    add_textbox(slide, ex, ey + 0.6, 6.5, 0.3,
                "E_cuda = (E_sched + E_score + E_addr + E_coll + E_arb + E_rf + E_xbar + E_wb) / 32  +  E_alu",
                font_size=9, bold=True, color=WHITE, font_name="Consolas")
    add_textbox(slide, ex, ey + 0.9, 6.5, 0.25,
                "= (0.875+0.525+1.05+1.40+0.525+0.562+0.70+0.863) / 32  +  1.530",
                font_size=9, color=MED_GRAY, font_name="Consolas")
    add_textbox(slide, ex, ey + 1.15, 6.5, 0.25,
                "= 0.203 (overhead)  +  1.530 (ALU)  =  1.733 pJ/op  |  Eff: 88.3%",
                font_size=11, bold=True, color=GREEN, font_name="Consolas")

    # TensorCore path
    add_textbox(slide, ex, ey + 1.55, 6.5, 0.25,
                "TensorCore Path (amortized over 64 MACs per MMA):", font_size=12,
                bold=True, color=LIGHT_GRAY)
    add_textbox(slide, ex, ey + 1.8, 6.5, 0.3,
                "E_tc = (E_sched + E_score + E_addr + E_coll + E_arb + E_rf + E_xbar + E_wb) / 64  +  E_alu",
                font_size=9, bold=True, color=WHITE, font_name="Consolas")
    add_textbox(slide, ex, ey + 2.1, 6.5, 0.25,
                "= (0.875+0.525+0.788+0.875+0.35+0.90+0.525+0.978) / 64  +  1.530",
                font_size=9, color=MED_GRAY, font_name="Consolas")
    add_textbox(slide, ex, ey + 2.35, 6.5, 0.25,
                "= 0.075 (overhead)  +  1.530 (ALU)  =  1.605 pJ/op  |  Eff: 95.3%",
                font_size=11, bold=True, color=GREEN, font_name="Consolas")

    # What's critical
    add_line(slide, ex, ey + 2.75, ex + 6.3, ey + 2.75, color=ORANGE, width=1)
    add_textbox(slide, ex, ey + 2.85, 6.5, 0.25,
                "The Critical Point: Per-MAC vs Per-SM", font_size=14,
                bold=True, color=ORANGE)
    add_multiline_textbox(slide, ex, ey + 3.15, 6.5, 2.0, [
        ("Per-MAC overhead looks small (0.07-0.20 pJ) because it's amortized", False, LIGHT_GRAY),
        ("over 32-64 MACs. But the SM has 4 partitions, each with:", True, WHITE),
        ("", False, WHITE),
        ("  * 16K registers (64 KB SRAM) -- always powered", False, LIGHT_GRAY),
        ("  * Warp scheduler + scoreboard -- tracking 16 warps", False, LIGHT_GRAY),
        ("  * Operand collector -- continuous buffering", False, LIGHT_GRAY),
        ("  * Crossbar -- always ready for next warp switch", False, LIGHT_GRAY),
        ("", False, WHITE),
        ("This infrastructure exists for ALL 2048 resident threads.", False, LIGHT_GRAY),
        ("Static power of 256 KB RF is NOT in the per-instruction model.", True, YELLOW),
    ], font_size=10)

    # Granularity
    add_textbox(slide, 0.3, 6.7, 12.5, 0.3,
                "Granularity: 1 SM = 4 partitions = 128 CUDA + 4 TC  |  "
                "Full GPU: 132 SMs (H100) share L2  |  SM has private L1/Shared Mem (192 KB)",
                font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    return slide


# ============================================================================
# SLIDE 4: TPU Core
# ============================================================================

def build_slide_4_tpu(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.1, 12.5, 0.45,
                "TPU Core: Input Buffer + 128x128 Systolic Array",
                font_size=26, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.3, 0.55, 12.5, 0.3,
                "Weight-stationary dataflow: weights loaded once, inputs stream through, outputs drain",
                font_size=13, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    # -- Diagram (left) --
    dx = 0.5
    dy = 1.0

    # Weight Buffer
    add_rect(slide, dx, dy, 2.0, 0.6, BUF_COLOR,
             "Weight Buffer\n(SRAM)", font_size=10, bold=True, line_color=TPU_COLOR)
    add_textbox(slide, dx + 2.1, dy + 0.1, 1.5, 0.4,
                "Load once\nper tile", font_size=8, color=MED_GRAY)

    add_arrow_down(slide, dx + 0.9, dy + 0.6, 0.2, 0.2, TPU_COLOR)

    # Input Buffer
    add_rect(slide, dx + 2.5, dy + 1.0, 2.0, 0.6, BUF_COLOR,
             "Input Buffer\n(L1 SRAM)", font_size=10, bold=True, line_color=TPU_COLOR)
    add_arrow_right(slide, dx + 2.2, dy + 1.2, 0.25, 0.2, TPU_COLOR)

    # Systolic Array
    add_rect(slide, dx, dy + 1.0, 2.2, 2.5, RGBColor(0x2A, 0x3A, 0x4E),
             "", line_color=TPU_COLOR)
    add_textbox(slide, dx, dy + 1.05, 2.2, 0.25,
                "128 x 128 Systolic Array", font_size=10, bold=True,
                color=TPU_COLOR, alignment=PP_ALIGN.CENTER)

    # PE grid (small representation)
    for r in range(6):
        for c in range(6):
            px = dx + 0.15 + c * 0.33
            py_pe = dy + 1.4 + r * 0.3
            add_rect(slide, px, py_pe, 0.28, 0.24,
                     RGBColor(0xAA, 0x66 + r*5, 0x22),
                     "MAC", font_size=5, font_color=WHITE)

    add_textbox(slide, dx, dy + 3.2, 2.2, 0.25,
                "16,384 MACs/cycle", font_size=10, bold=True,
                color=TPU_COLOR, alignment=PP_ALIGN.CENTER)

    # Accumulators + Output
    add_arrow_down(slide, dx + 0.9, dy + 3.5, 0.2, 0.2, TPU_COLOR)
    add_rect(slide, dx, dy + 3.75, 2.2, 0.5, BUF_COLOR,
             "Accumulators +\nOutput Buffer", font_size=9, bold=True,
             line_color=TPU_COLOR)

    # -- Energy equation (right) --
    ex = 5.0
    ey = 1.0

    add_textbox(slide, ex, ey, 7.5, 0.3,
                "TPU Core Energy Equation (per MAC):", font_size=14,
                bold=True, color=TPU_COLOR)

    add_textbox(slide, ex, ey + 0.35, 7.5, 0.3,
                "E_tpu = E_weight_load/reuse + E_input_feed + E_systolic_mac + E_output_drain/array_size",
                font_size=11, bold=True, color=WHITE, font_name="Consolas")

    add_textbox(slide, ex, ey + 0.7, 7.5, 0.3,
                "E_tpu = 0.500/64 + 0.300 + 1.440 + 0.300/128",
                font_size=11, color=MED_GRAY, font_name="Consolas")

    add_textbox(slide, ex, ey + 1.0, 7.5, 0.25,
                "= 0.008 + 0.300 + 1.440 + 0.002 = 1.750 pJ/op  (input-feed dominated overhead)",
                font_size=11, bold=True, color=GREEN, font_name="Consolas")

    # But model says 2.085
    add_textbox(slide, ex, ey + 1.4, 7.5, 0.25,
                "Model output (from build_tpu_cycle_energy at 7nm, L1-resident):",
                font_size=12, bold=True, color=LIGHT_GRAY)

    tpu_events = [
        ("Config (one-time)", "0.003 pJ/op", MED_GRAY),
        ("Weight load (amortized 64x)", "0.016 pJ/op", BUF_COLOR),
        ("Input feed (one edge, per cycle)", "0.610 pJ/op", ORANGE),
        ("Systolic MAC (128x128)", "1.440 pJ/op", ALU_COLOR),
        ("Partial sum accumulation", "0.006 pJ/op", MED_GRAY),
        ("Output drain (one edge)", "0.010 pJ/op", BUF_COLOR),
        ("TOTAL (core only)", "2.085 pJ/op", TPU_COLOR),
    ]

    for i, (name, val, color) in enumerate(tpu_events):
        y = ey + 1.7 + i * 0.28
        is_total = "TOTAL" in name
        add_textbox(slide, ex + 0.2, y, 4.0, 0.25,
                    name, font_size=10, bold=is_total, color=color)
        add_textbox(slide, ex + 4.5, y, 2.0, 0.25,
                    val, font_size=10, bold=is_total, color=YELLOW if is_total else LIGHT_GRAY,
                    font_name="Consolas")

    # Key characteristics
    char_y = ey + 3.8
    add_line(slide, ex, char_y, ex + 7.5, char_y, color=TPU_COLOR, width=1)
    add_multiline_textbox(slide, ex, char_y + 0.1, 7.5, 1.5, [
        ("Key Characteristics:", True, TPU_COLOR),
        ("  * No instruction fetch per operation (fixed-function)", False, GREEN),
        ("  * No register file arbitration (data flows spatially)", False, GREEN),
        ("  * Weight-stationary ONLY -- cannot do input/output stationary", False, ORANGE),
        ("  * Input feed (0.610 pJ/op) dominates overhead -- each input must enter at array edge", False, LIGHT_GRAY),
        ("  * Granularity is LARGE: 128x128 = 16,384 MACs per cycle. Small workloads waste MACs.", False, ORANGE),
    ], font_size=10)

    add_textbox(slide, 0.3, 6.7, 12.5, 0.3,
                "Granularity: 1 MXU = 128x128 systolic array  |  "
                "Full TPU: 2 MXUs + Unified Buffer (24 MB)  |  Weight-stationary only",
                font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    return slide


# ============================================================================
# SLIDE 5: KPU Tile
# ============================================================================

def build_slide_5_kpu(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.1, 12.5, 0.45,
                "KPU Tile: Input Buffer + 16x16 PE Array + SURE Network",
                font_size=26, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.3, 0.55, 12.5, 0.3,
                "Programmable domain flow: A/B/C-stationary, any SURE schedule, compile-time routed",
                font_size=13, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    # -- Diagram (left) --
    dx = 0.5
    dy = 1.0

    # EDDO L1 Buffer
    add_rect(slide, dx, dy, 2.5, 0.6, BUF_COLOR,
             "EDDO L1 Stream Buffer\n(software-managed SRAM)", font_size=9, bold=True,
             line_color=KPU_COLOR)

    add_arrow_down(slide, dx + 1.15, dy + 0.6, 0.2, 0.15, KPU_COLOR)

    # Streamer
    add_rect(slide, dx, dy + 0.8, 2.5, 0.4, RGBColor(0x44, 0x77, 0x55),
             "Streamer (vector -> token stream)", font_size=9, bold=True,
             line_color=SURE_COLOR)

    add_arrow_down(slide, dx + 1.15, dy + 1.2, 0.2, 0.15, SURE_COLOR)

    # PE Array
    add_rect(slide, dx, dy + 1.45, 2.5, 2.2, RGBColor(0x1A, 0x3A, 0x4E),
             "", line_color=KPU_COLOR)
    add_textbox(slide, dx, dy + 1.5, 2.5, 0.25,
                "16 x 16 PE Array", font_size=10, bold=True,
                color=KPU_COLOR, alignment=PP_ALIGN.CENTER)

    # PE grid
    for r in range(5):
        for c in range(5):
            px = dx + 0.15 + c * 0.44
            py_pe = dy + 1.85 + r * 0.3
            add_rect(slide, px, py_pe, 0.38, 0.24, PE_COLOR,
                     "PE", font_size=6, font_color=WHITE, bold=True)
            # Horizontal wire
            if c < 4:
                add_line(slide, px + 0.38, py_pe + 0.12, px + 0.44, py_pe + 0.12,
                         color=SURE_COLOR, width=0.75)

    add_textbox(slide, dx, dy + 3.4, 2.5, 0.2,
                "256 MACs/cycle per tile", font_size=10, bold=True,
                color=KPU_COLOR, alignment=PP_ALIGN.CENTER)

    # SURE network annotation
    add_textbox(slide, dx, dy + 3.6, 2.5, 0.2,
                "SURE network: PE-to-PE wires", font_size=9, bold=True,
                color=SURE_COLOR, alignment=PP_ALIGN.CENTER)

    # -- Energy equation (right) --
    ex = 5.0
    ey = 1.0

    add_textbox(slide, ex, ey, 7.5, 0.3,
                "KPU Tile Energy Equation (per MAC):", font_size=14,
                bold=True, color=KPU_COLOR)

    add_textbox(slide, ex, ey + 0.35, 7.5, 0.3,
                "E_kpu = E_config/N_ops + 2 * E_pe_wire + E_domain_mac + E_local_reg",
                font_size=11, bold=True, color=WHITE, font_name="Consolas")

    add_textbox(slide, ex, ey + 0.7, 7.5, 0.3,
                "E_kpu = ~0 + 2 * 0.022 + E_mac(precision) + 0.009",
                font_size=11, color=MED_GRAY, font_name="Consolas")

    add_textbox(slide, ex, ey + 1.0, 7.5, 0.25,
                "BF16: 0.053 + 1.080 = 1.133 pJ/op  |  INT8: 0.053 + 0.675 = 0.728 pJ/op",
                font_size=11, bold=True, color=GREEN, font_name="Consolas")

    # Model output
    add_textbox(slide, ex, ey + 1.4, 7.5, 0.25,
                "Model output (from build_kpu_cycle_energy at 7nm, L1-resident, mixed precision):",
                font_size=12, bold=True, color=LIGHT_GRAY)

    kpu_events = [
        ("Domain program load (1-time)", "0.000 pJ/op (amortized)", MED_GRAY),
        ("Domain tracker (64 tiles)", "0.001 pJ/op (amortized)", MED_GRAY),
        ("PE-to-PE transfer (x2 operands)", "0.044 pJ/op", SURE_COLOR),
        ("Local register write", "0.009 pJ/op", PE_COLOR),
        ("INT8 MACs (70%)", "0.473 pJ/op (weighted)", ALU_COLOR),
        ("BF16 MACs (20%)", "0.216 pJ/op (weighted)", ALU_COLOR),
        ("FP32 MACs (10%)", "0.216 pJ/op (weighted)", ALU_COLOR),
        ("TOTAL (core only, mixed)", "0.958 pJ/op", KPU_COLOR),
    ]

    for i, (name, val, color) in enumerate(kpu_events):
        y = ey + 1.7 + i * 0.27
        is_total = "TOTAL" in name
        add_textbox(slide, ex + 0.2, y, 4.0, 0.25,
                    name, font_size=10, bold=is_total, color=color)
        add_textbox(slide, ex + 4.5, y, 2.5, 0.25,
                    val, font_size=10, bold=is_total, color=YELLOW if is_total else LIGHT_GRAY,
                    font_name="Consolas")

    # Key characteristics
    char_y = ey + 3.9
    add_line(slide, ex, char_y, ex + 7.5, char_y, color=KPU_COLOR, width=1)
    add_multiline_textbox(slide, ex, char_y + 0.1, 7.5, 1.5, [
        ("Key Characteristics:", True, KPU_COLOR),
        ("  * No instruction fetch per operation (domain program loaded once)", False, GREEN),
        ("  * No register file -- PE-to-PE wire transfer at 0.022 pJ (vs GPU's ~3 pJ/warp infrastructure)", False, GREEN),
        ("  * Programmable: A/B/C-stationary, any SURE schedule (vs TPU weight-stationary only)", False, GREEN),
        ("  * Small tile (256 MACs) = better utilization at batch=1 than TPU's 16,384 MACs", False, LIGHT_GRAY),
        ("  * Overhead is 5.5% of total -- 94.4% of energy is useful compute", True, YELLOW),
    ], font_size=10)

    add_textbox(slide, 0.3, 6.7, 12.5, 0.3,
                "Granularity: 1 tile = 16x16 = 256 MACs  |  "
                "Full KPU: 64 tiles (KPU-T64) to 768 tiles (KPU-T768)  |  "
                "Each tile has private L1 stream buffer",
                font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    return slide


# ============================================================================
# SLIDE 6: Side-by-Side Summary
# ============================================================================

def build_slide_6_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.3, 0.1, 12.5, 0.45,
                "Core Compute Energy: Side-by-Side at 7nm (No External Memory)",
                font_size=26, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Summary table
    headers = ["", "CPU Scalar", "CPU AVX-512", "GPU CUDA", "GPU TensorCore",
               "TPU 128x128", "KPU 16x16"]
    data = [
        ["MAC energy (pJ)",      "1.800", "1.800", "1.530", "1.530", "1.440", "0.905*"],
        ["Overhead (pJ/op)",     "8.136", "0.509", "0.418", "0.284", "0.639", "0.053"],
        ["Total (pJ/op)",        "9.936", "2.309", "1.948", "1.814", "2.085", "0.958"],
        ["Compute Efficiency",   "18.1%", "78.0%", "78.5%", "84.4%", "69.4%", "94.4%"],
        ["Amortization",         "1",     "16",    "32",    "64",    "128-16K","256"],
        ["MACs per cycle",       "1",     "16",    "128",   "256",   "16,384", "256"],
        ["Has instruction fetch","Yes",   "Yes",   "Yes",   "Yes",   "No",     "No"],
        ["Has register file",    "Yes",   "Yes",   "256KB", "256KB", "No (PE)","No (PE)"],
        ["Programmable schedule","N/A",   "N/A",   "N/A",   "N/A",   "Weight only","A/B/C"],
    ]

    tbl = slide.shapes.add_table(
        len(data) + 1, len(headers),
        Inches(0.3), Inches(0.6), Inches(12.6), Inches(3.5)).table

    arch_colors = [RGBColor(0x33, 0x33, 0x55), CPU_COLOR, CPU_COLOR,
                   GPU_COLOR, GPU_COLOR, TPU_COLOR, KPU_COLOR]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = arch_colors[j]

    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.bold = (j == 0 or i <= 3)
                p.alignment = PP_ALIGN.CENTER
                if i == 2:  # Total row
                    p.font.color.rgb = YELLOW
                elif i == 3:  # Efficiency row
                    p.font.color.rgb = GREEN
                elif j == 0:
                    p.font.color.rgb = ACCENT_TEAL
                else:
                    p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_EVEN if i % 2 == 0 else ROW_ODD

    add_textbox(slide, 0.3, 4.15, 12.5, 0.2,
                "* KPU mixed precision: 70% INT8 (0.675) + 20% BF16 (1.080) + 10% FP32 (2.160) = 0.905 pJ avg",
                font_size=9, color=MED_GRAY)

    # Stacked energy bars
    bar_y = 4.45
    add_textbox(slide, 0.5, bar_y, 12.0, 0.25,
                "Energy per Op (pJ) -- Green = Compute, Gray = Overhead:",
                font_size=13, bold=True, color=WHITE)

    bar_data = [
        ("CPU Scalar",   1.800, 8.136, 9.936, CPU_COLOR),
        ("CPU AVX-512",  1.800, 0.509, 2.309, CPU_COLOR),
        ("GPU CUDA",     1.530, 0.418, 1.948, GPU_COLOR),
        ("GPU TC",       1.530, 0.284, 1.814, GPU_COLOR),
        ("TPU 128x128",  1.440, 0.639, 2.085, TPU_COLOR),
        ("KPU 16x16",    0.905, 0.053, 0.958, KPU_COLOR),
    ]

    max_e = 10.0  # scale
    bar_w = 8.0
    bar_start = 3.0
    bar_h = 0.32

    for i, (name, compute, overhead, total, color) in enumerate(bar_data):
        by = bar_y + 0.3 + i * (bar_h + 0.08)

        # Label
        add_textbox(slide, 0.5, by, 2.4, bar_h,
                    name, font_size=10, bold=True, color=color,
                    alignment=PP_ALIGN.RIGHT)

        # Compute bar
        cw = bar_w * (compute / max_e)
        add_rect(slide, bar_start, by, cw, bar_h, ALU_COLOR,
                 f"{compute:.2f}", font_size=8, font_color=WHITE if cw > 0.5 else WHITE, bold=True)

        # Overhead bar
        ow = bar_w * (overhead / max_e)
        if ow > 0.15:
            add_rect(slide, bar_start + cw, by, ow, bar_h,
                     RGBColor(0x55, 0x55, 0x77),
                     f"{overhead:.2f}" if ow > 0.3 else "",
                     font_size=7, font_color=MED_GRAY, bold=True)

        # Total label
        add_textbox(slide, bar_start + cw + ow + 0.05, by, 1.5, bar_h,
                    f"{total:.2f} pJ", font_size=9, bold=True, color=LIGHT_GRAY)

    # Key takeaway
    tk_y = bar_y + 0.3 + 6 * (bar_h + 0.08) + 0.15
    add_line(slide, 0.5, tk_y, 12.5, tk_y, color=ACCENT_BLUE, width=2)
    add_rounded_rect(slide, 0.5, tk_y + 0.1, 12.0, 0.8, RGBColor(0x2A, 0x2A, 0x4E))
    add_multiline_textbox(slide, 0.8, tk_y + 0.15, 11.4, 0.7, [
        ("Core Compute Takeaway (all at 7nm, no memory):", True, ACCENT_TEAL),
        ("The MAC energy itself is comparable: 0.9-1.8 pJ across all architectures (2x range).", False, LIGHT_GRAY),
        ("The overhead varies 150x (0.053 to 8.136 pJ). CPU scalar is worst; KPU is best.", True, WHITE),
        ("KPU achieves 94% compute efficiency vs GPU TensorCore's 84% -- a 1.9x total energy advantage at same node.",
         True, KPU_COLOR),
        ("The 20x argument requires adding memory hierarchy (domain 2) and system I/O (domain 3) advantages.",
         False, ORANGE),
    ], font_size=11)

    return slide


# ============================================================================
# MAIN
# ============================================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1(prs)
    build_slide_2_cpu(prs)
    build_slide_3_gpu(prs)
    build_slide_4_tpu(prs)
    build_slide_5_kpu(prs)
    build_slide_6_comparison(prs)

    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "..", "core_compute_energy.pptx")
    output_path = os.path.normpath(output_path)
    prs.save(output_path)
    print(f"Saved presentation to: {output_path}")
    print(f"  6 slides generated:")
    print(f"  1. Title: The Three Modeling Domains")
    print(f"  2. CPU Core: Scalar ALU + SIMD ALU + Register Files")
    print(f"  3. GPU SM: SIMT Engine + TensorCores + 64K Register File")
    print(f"  4. TPU Core: Input Buffer + 128x128 Systolic Array")
    print(f"  5. KPU Tile: Input Buffer + 16x16 PE Array + SURE Network")
    print(f"  6. Side-by-Side: Core Compute Energy at 7nm")


if __name__ == "__main__":
    main()
